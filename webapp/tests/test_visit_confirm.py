"""자동분류 전 **차수 확인** — 사람이 보고 필요하면 고친다.

화면은 인식된 차수(슬라이드 번호 포함)를 보여주고, 이번 차수와 새 슬라이드
위치를 미리 채워 둔다. 맞으면 그대로 진행하고, 고치면 그 값으로 세션이 열린다.
프런트가 먼저 거르지만 **서버가 한 번 더 검증한다** — 의료 기록물이라 이중 잠금이다.
"""
import io
import os
import shutil
import sys

import numpy as np
import pytest
from PIL import Image
from starlette.testclient import TestClient

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "backend"))

import main as M                                                  # noqa: E402
import naming as N                                                # noqa: E402

NAME, HOSP, ORTHO = "차수확인", "111222333", "54321"
IDS = N.Identifiers(NAME, HOSP, ORTHO)
FOLDER = N.folder_name(IDS, M.cfg.naming.folder_pattern)
PPT = N.ppt_filename(IDS, M.cfg.naming.ppt_pattern)


def _png() -> io.BytesIO:
    buf = io.BytesIO()
    Image.fromarray(np.full((300, 400, 3), 128, np.uint8)).save(buf, "PNG")
    buf.seek(0)
    return buf


def _deck(path, labels):
    """라벨 + 십자뷰(폭 8.4cm 사진 5장)짜리 슬라이드로 덱을 만든다."""
    from pptx import Presentation
    from pptx.util import Emu
    prs = Presentation()
    w = int(round(8.4 * M.EMU_PER_CM))
    for label in labels:
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        tb = sl.shapes.add_textbox(Emu(0), Emu(0), Emu(3000000), Emu(500000))
        tb.text_frame.text = label
        for i in range(5):
            sl.shapes.add_picture(_png(), Emu(i * w), Emu(1000000), width=Emu(w))
    prs.save(str(path))


@pytest.fixture
def patient():
    with TestClient(M.app) as c:
        d = M.ROOT / FOLDER
        shutil.rmtree(d, ignore_errors=True)
        d.mkdir(parents=True)
        _deck(d / PPT, ["24.06.05 (초진 A)", "24.09.04 (재진 B)"])
        yield c, d
        shutil.rmtree(d, ignore_errors=True)


def test_확인줄_재료가_목록에_실린다(patient):
    """인식된 차수·슬라이드 번호·제안 위치 — 화면이 이걸로 확인 줄을 그린다."""
    c, _ = patient
    me = [p for p in c.get("/api/patients").json()["patients"]
          if p["folder"] == FOLDER]
    assert me, "환자가 목록에 없다"
    p = me[0]
    assert [v["visit"] for v in p["visit_slides"]] == ["A", "B"]
    assert [v["slide_no"] for v in p["visit_slides"]] == [1, 2]
    assert p["next_visit"] == "C"
    assert p["suggest_after"] == 2           # 마지막 차수 슬라이드(2번) **뒤**
    assert p["ppt_slides"] == 2
    assert p["label_excluded"] == [] and p["label_fallback"] is False


def test_기본값은_그대로_진행된다(patient):
    c, _ = patient
    r = c.post("/api/session", json={"folder": FOLDER})
    assert r.status_code == 200, r.text
    assert r.json()["visit"] == "C"


def test_고친_차수로_세션이_열린다(patient):
    """구내촬영 없이 지나간 회차가 있어 사람이 차수를 올려 잡는 경우."""
    c, _ = patient
    r = c.post("/api/session", json={"folder": FOLDER, "visit": "E", "insert_after": 1})
    assert r.status_code == 200, r.text
    assert r.json()["visit"] == "E"
    s = M.get_session(r.json()["session_id"])
    assert s.insert_after == 1               # 1번 장 뒤 = 새 장이 2번이 된다


def test_이미_있는_차수는_서버가_막는다(patient):
    c, _ = patient
    r = c.post("/api/session", json={"folder": FOLDER, "visit": "B"})
    assert r.status_code == 400
    assert "이미" in r.json()["detail"]


def test_슬라이드_위치가_범위를_벗어나면_막는다(patient):
    """0 = 맨 앞, 최대 = 덱의 마지막 장 뒤. 그 밖은 서버가 거른다."""
    c, _ = patient
    assert c.post("/api/session",
                  json={"folder": FOLDER, "insert_after": 99}).status_code == 400
    assert c.post("/api/session",
                  json={"folder": FOLDER, "insert_after": -1}).status_code == 400
    assert c.post("/api/session",
                  json={"folder": FOLDER, "insert_after": 0}).status_code == 200


def test_차수_형식이_아니면_막는다(patient):
    c, _ = patient
    for bad in ("3", "재진", "ABC"):
        r = c.post("/api/session", json={"folder": FOLDER, "visit": bad})
        assert r.status_code == 400, (bad, r.text)
