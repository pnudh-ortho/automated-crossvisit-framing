"""수제 덱의 사진 크기·자리를 **어느 장에서** 물려받는가.

정합 기준영상 · 표기 지문 · 계측선 · 슬롯 레이아웃은 모두 같은 한 장 —
**직전 차수 슬라이드**(차수 글자가 가장 큰 장) — 에서 와야 한다. 규칙을 각자 다시
쓰면 어긋난다. 실제로 레이아웃만 "장 순서상 마지막"을 골라서, 정합은 D 를 보는데
사진은 C 의 크기로 들어가는 일이 있었다.

실행: cd webapp && python -m pytest tests/test_layout_source.py -q
"""
import os
import pathlib
import sys
import tempfile

import pytest
from pptx import Presentation
from pptx.util import Cm

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "backend"))

import main as M  # noqa: E402
import ppt_reader as Rd  # noqa: E402
from coords import EMU_PER_CM  # noqa: E402

_PNG = (b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
        b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc```\x00\x00"
        b"\x00\x04\x00\x01\xf6\x178U\x00\x00\x00\x00IEND\xaeB`\x82")


@pytest.fixture
def png(tmp_path):
    p = tmp_path / "d.png"
    p.write_bytes(_PNG)
    return str(p)


def _deck(png, sw=33.87, sh=19.05):
    p = Presentation()
    p.slide_width, p.slide_height = Cm(sw), Cm(sh)
    return p


def _cross(p, png, w_cm, h_cm=None, label="24.06.05 (재진 B)"):
    """십자로 사진 5장을 놓은 수제 슬라이드 한 장."""
    h_cm = h_cm if h_cm is not None else w_cm * 0.75
    s = p.slides.add_slide(p.slide_layouts[6])
    sw, sh = M.emu_to_cm(p.slide_width), M.emu_to_cm(p.slide_height)
    cx, cy = sw / 2 - w_cm / 2, sh / 2 - h_cm / 2
    for dx, dy in [(0, 0), (0, -h_cm - .2), (0, h_cm + .2),
                   (-w_cm - .2, 0), (w_cm + .2, 0)]:
        s.shapes.add_picture(png, Cm(cx + dx), Cm(cy + dy), Cm(w_cm), Cm(h_cm))
    if label:
        tb = s.shapes.add_textbox(Cm(0.5), Cm(0.2), Cm(10), Cm(0.9))
        tb.text_frame.text = label
    return s


def _front(p):
    """세션이 하는 그대로 — 장부를 세고 그 장의 창을 받는다."""
    scan = Rd.scan_ppt_visits(p, M.cfg)
    return M._layout_from_ppt(p, M._prev_visit_slide_no(scan))["SLOT_FRONT"]


def test_직전_차수의_레이아웃을_따른다(png):
    """장 순서가 아니라 차수 글자다. 확인 줄에서 삽입 자리를 직접 고르거나 사람이
    장을 옮기면 둘이 어긋난다."""
    p = _deck(png)
    _cross(p, png, 9.0, label="24.07.18 (재진 D)")   # 직전 차수
    _cross(p, png, 8.6, label="24.06.14 (재진 C)")   # 장 순서로는 마지막
    assert abs(_front(p).w - 9.0) < 1e-3


def test_라벨_없는_장은_기준이_되지_않는다(png):
    """사진 다섯 장이 십자로 놓였다는 것만으로는 차수 슬라이드가 아니다 —
    얼굴 장·비교용 장이 그럴 수 있다."""
    p = _deck(png)
    _cross(p, png, 9.0, label="24.07.18 (재진 D)")
    _cross(p, png, 8.6, label=None)
    assert abs(_front(p).w - 9.0) < 1e-3


def test_문턱은_차수_장부와_같은_값이다(png):
    """레이아웃 쪽만 8.0 을 따로 들고 있어서, 7.9cm 덱이 차수로는 세어지는데
    레이아웃은 안 물려받는 어중간한 상태가 됐다 — 사진이 덱보다 크게 들어갔다."""
    p = _deck(png)
    _cross(p, png, Rd.CROSS_MIN_W_CM + 0.1)
    got = _front(p)
    assert abs(got.w - (Rd.CROSS_MIN_W_CM + 0.1)) < 1e-3
    assert abs(got.w - M.SLOT_WINDOWS["SLOT_FRONT"].w) > 1e-3   # 템플릿 기본값이 아니다


def test_4대3_이_아니어도_그대로_들어간다(png):
    """8.38 x 6.28 은 4:3(6.285)이 아니다. 비율을 맞추려 들면 6.29 가 되어
    직전 차수 사진과 크기가 달라 보인다."""
    p = _deck(png)
    _cross(p, png, 8.38, 6.28)
    win = _front(p)
    pl = M._exact_placement(win)
    assert abs(pl.ext_cx / EMU_PER_CM - 8.38) < 1e-3
    assert abs(pl.ext_cy / EMU_PER_CM - 6.28) < 1e-3


def test_슬라이드_크기가_달라도_물려받는다(png):
    """양식과 다른 크기로 만든 덱 — 창은 슬라이드 안의 절대 좌표라 그대로 읽힌다."""
    p = _deck(png, sw=24.866, sh=18.516)
    _cross(p, png, 8.38, 6.28)
    scan = Rd.scan_ppt_visits(p, M.cfg)
    assert [v["visit"] for v in scan["visits"]] == ["B"]
    win = M._layout_from_ppt(p, M._prev_visit_slide_no(scan))["SLOT_FRONT"]
    assert abs(win.w - 8.38) < 1e-3 and abs(win.h - 6.28) < 1e-3


def test_차수를_못_읽으면_종전처럼_훑는다(png):
    """번호를 못 주는 덱 — 그래도 사진이 든 장을 찾아 물려받아야 한다."""
    p = _deck(png)
    _cross(p, png, 9.0, label=None)
    assert abs(M._layout_from_ppt(p, None)["SLOT_FRONT"].w - 9.0) < 1e-3


# ── 검수 화면이 그리는 슬라이드 크기 ───────────────────────────────────────
def test_화면은_그_덱의_슬라이드_크기로_그린다(png, tmp_path):
    """판 위에 겹쳐 그리는 계측선·노트 상자는 슬라이드 cm 를 좌표계로 쓴다.
    양식 크기로 굳혀 두면 크기가 다른 덱에서 그 겹침이 통째로 어긋난 자리에
    그려진다 — 결과물은 맞는데 화면이 거짓말을 한다."""
    import shutil
    from starlette.testclient import TestClient
    import naming as N

    ids = N.Identifiers("크기검사", "111222777", "54327")
    folder = N.folder_name(ids, M.cfg.naming.folder_pattern)

    p = _deck(png, sw=24.866, sh=18.516)
    _cross(p, png, 8.38, 6.28, label="24.06.05 (초진 A)")
    with TestClient(M.app) as c:
        d = M.ROOT / folder
        shutil.rmtree(d, ignore_errors=True)
        d.mkdir(parents=True)
        try:
            p.save(str(d / f"{folder}.pptx"))
            sid = c.post("/api/session", json={"folder": folder}).json()["session_id"]
            slide = c.get(f"/api/notes/{sid}").json()["slide"]
            assert abs(slide["w"] - 24.866) < 1e-2, slide
            assert abs(slide["h"] - 18.516) < 1e-2, slide
        finally:
            shutil.rmtree(d, ignore_errors=True)


def test_새_덱은_양식_크기_그대로다(png):
    """초진은 양식으로 만든다 — 기본값이 바뀌면 안 된다."""
    s = M.Session("first", None, "A")
    assert s.slide_cm == M.CASE_SLIDE_CM
