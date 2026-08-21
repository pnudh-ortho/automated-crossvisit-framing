"""
덱 저장 속도 — 이미 압축된 사진을 다시 압축하지 않는다 (2026-08-21)

131MB 짜리 실측 덱에서 저장이 미리보기 시간의 61% 를 먹고 있었다. 그런데 그
131MB 는 압축해 봐야 줄지가 않는다:

    ppt/media (사진)   130.8 MB → 130.8 MB    0.0% 줄어듦
    xml 등               2.5 MB →   0.1 MB   94.5% 줄어듦

이미 JPEG 이기 때문이다. 줄지도 않는 것을 훑느라 쓰는 시간이 노트북에서는
그대로 사람이 기다리는 시간이 됐다(같은 덱에서 2.11초 → 0.15초).

이 조정은 python-pptx 의 **사설 API**(_ZipPkgWriter.write)를 갈아 끼워 이뤄진다.
그쪽 판이 바뀌면 조용히 종전 동작으로 물러나도록 해 두었는데, 조용히 물러나면
아무도 모르게 다시 느려진다 — 그래서 여기서 못박는다.

실행: cd webapp && python -m pytest tests/test_deck_save.py -q
"""
import io
import os
import sys
import zipfile

import pytest
from PIL import Image

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "backend"))

import template as T          # noqa: E402


def _deck_with_a_photo(tmp_path):
    """사진 한 장이 든 최소 덱. 실제 저장 경로를 그대로 탄다."""
    from pptx import Presentation
    from pptx.util import Cm

    img = tmp_path / "p.jpg"
    # 잡음이라 JPEG 으로 이미 압축돼 있고 deflate 로는 더 줄지 않는다 — 실제 사진과 같은 처지
    Image.frombytes("RGB", (600, 450), os.urandom(600 * 450 * 3)).save(img, quality=90)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(img), 0, 0, Cm(10), Cm(7.5))
    out = tmp_path / "deck.pptx"
    prs.save(str(out))
    return out


def test_the_fast_save_patch_is_in_place():
    """python-pptx 가 바뀌어 갈아 끼우기가 실패하면 여기서 잡는다."""
    assert T.FAST_SAVE, "저장 최적화가 적용되지 않았다 — python-pptx 쪽 판이 바뀌었나"


def test_photos_are_stored_but_xml_is_still_compressed(tmp_path):
    """사진은 그대로 담고(STORED), XML 은 종전대로 압축한다(DEFLATED)."""
    z = zipfile.ZipFile(_deck_with_a_photo(tmp_path))
    media = [i for i in z.infolist() if i.filename.startswith("ppt/media/")]
    xml = [i for i in z.infolist() if i.filename.endswith(".xml")]
    assert media and xml

    for i in media:
        assert i.compress_type == zipfile.ZIP_STORED, i.filename
    assert any(i.compress_type == zipfile.ZIP_DEFLATED for i in xml), \
        "XML 까지 압축을 끄면 파일이 쓸데없이 커진다"


def test_the_deck_still_opens_and_keeps_its_shapes(tmp_path):
    """빨라진 것과 별개로 결과물이 온전해야 한다 — zip 도, 덱도."""
    out = _deck_with_a_photo(tmp_path)
    assert zipfile.ZipFile(out).testzip() is None, "zip 이 깨졌다"

    prs = T.load_presentation(out)
    assert len(prs.slides) == 1
    pics = [sh for sh in prs.slides[0].shapes if getattr(sh, "shape_type", None) == 13]
    assert len(pics) == 1, "사진이 사라졌다"


def test_uncompressed_media_does_not_bloat_the_file(tmp_path):
    """압축을 꺼도 파일이 커지지 않는다 — 애초에 줄던 것이 아니기 때문이다."""
    out = _deck_with_a_photo(tmp_path)
    z = zipfile.ZipFile(out)
    media = [i for i in z.infolist() if i.filename.startswith("ppt/media/")]
    grew = sum(i.compress_size for i in media) - sum(i.file_size for i in media)
    assert grew == 0, f"사진이 {grew} 바이트 늘었다"
