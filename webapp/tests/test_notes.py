"""
차수 노트 검증.

슬라이드의 텍스트 박스를 직접 타이핑하지 않고, 칸을 채우면 서버가 서식에 끼워
박스 텍스트를 만든다. 여기서 확인하는 것:

  - 비워 둔 칸이 만든 줄은 통째로 빠지는가 ("U: " 같은 껍데기를 안 남기는가)
  - 모르는 칸은 거절하는가
  - 확정하면 십자뷰 슬라이드의 박스에 실제로 적히는가
  - 아무것도 안 채우면 양식의 원래 문구를 건드리지 않는가

실행: cd webapp && python -m pytest tests/test_notes.py -q
"""
import io
import os
import shutil
import sys

import cv2
import numpy as np
import pytest
from starlette.testclient import TestClient

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "backend"))

import case_deck as CD  # noqa: E402
import main as M  # noqa: E402
import naming as N  # noqa: E402
import template as T  # noqa: E402

client = TestClient(M.app)

NAME, HOSP, ORTHO = "노트검사", "111222666", "54324"
IDS = N.Identifiers(NAME, HOSP, ORTHO)
FOLDER = N.folder_name(IDS, M.cfg.naming.folder_pattern)


@pytest.fixture
def patient():
    d = M.ROOT / FOLDER
    shutil.rmtree(d, ignore_errors=True)
    d.mkdir(parents=True)
    yield d
    shutil.rmtree(d, ignore_errors=True)


@pytest.fixture
def sid(patient):
    s = client.post("/api/session", json={"folder": FOLDER}).json()["session_id"]
    img = np.full((1200, 900, 3), 40, np.uint8)
    cv2.putText(img, "x", (100, 600), cv2.FONT_HERSHEY_SIMPLEX, 4, (255, 255, 255), 8)
    blob = cv2.imencode(".jpg", img)[1].tobytes()
    client.post(f"/api/photos/{s}", files=[("files", ("a.jpg", io.BytesIO(blob), "image/jpeg"))])
    client.post(f"/api/classify/{s}")
    return s


def test_fields_are_declared(sid):
    js = client.get(f"/api/notes/{sid}").json()
    keys = [f["key"] for f in js["fields"]]
    assert "tx_period" in keys and "rx_period" in keys
    # 앱이 채우는 것은 기간 세 줄뿐이다 — 자유 기입 칸도 와이어도 두지 않는다
    assert not ({"subj", "plan", "next", "ll", "wire_u", "wire_l"} & set(keys)), keys
    # preview 는 서식 박스 + 날짜 칸이다. 날짜는 서식이 아니라 차수 라벨에서
    # 오지만, 검수 화면 오버레이가 네 칸을 한자리에서 보여줘야 해서 같이 나간다.
    assert set(js["preview"]) == set(M.cfg.notes.boxes) | {"NOTE_DATE"}


def test_empty_fields_drop_their_lines():
    """빈 칸이 만든 줄은 사라져야 한다 — 슬라이드에 'Rx. Period:' 만 남으면 안 된다."""
    out = M._render_notes({"tx_period": "3 month"})
    status = out["NOTE_STATUS"]
    assert "Tx. Period: 3 month" in status
    assert "Rx. Period" not in status         # 안 채웠다
    assert "App. Period" not in status

    # 아무것도 안 채우면 내용이 없다 — 커밋이 박스를 건드리지 않는 신호.
    # 서식의 빈 줄은 그대로 남으므로(양식의 글 시작 높이) 공백을 걷어내고 본다.
    assert all(not v.strip() for v in M._render_notes({}).values())


def test_whitespace_only_counts_as_empty():
    out = M._render_notes({"tx_period": "   ", "rx_period": "3 month"})
    assert "Tx. Period" not in out["NOTE_STATUS"]
    assert "Rx. Period: 3 month" in out["NOTE_STATUS"]


def test_unknown_field_rejected(sid):
    r = client.post("/api/notes", json={"session_id": sid, "values": {"헛소리": "x"}})
    assert r.status_code == 400


def test_values_round_trip(sid):
    r = client.post("/api/notes", json={"session_id": sid,
                                        "values": {"tx_period": "3 month"}})
    js = r.json()
    assert js["values"]["tx_period"] == "3 month"
    assert "Tx. Period: 3 month" in js["preview"]["NOTE_STATUS"]
    # 다시 읽어도 남아 있다
    assert client.get(f"/api/notes/{sid}").json()["values"]["tx_period"] == "3 month"


def test_free_note_fields_are_gone(sid):
    """②④⑤ 와 와이어는 앱에서 적지 않는다 — 값을 보내면 거절한다."""
    for k in ("subj", "plan", "next", "ll", "wire_u", "wire_l"):
        r = client.post("/api/notes", json={"session_id": sid, "values": {k: "x"}})
        assert r.status_code == 400, (k, r.text)


@pytest.mark.skipif(not M.CASE_ANCHORS, reason="케이스 양식이 없습니다")
def test_commit_writes_notes_onto_cross_slide(sid, patient):
    client.post("/api/notes", json={"session_id": sid,
                                    "values": {"tx_period": "0 month"}})
    assert client.post(f"/api/commit/{sid}?allow_missing=true").status_code == 200

    prs = T.load_presentation(patient / N.ppt_filename(IDS, M.cfg.naming.ppt_pattern))
    cross = prs.slides[CD.cross_slide_index(prs)]
    # 자유 기입 상자는 만들어지되 앱이 글을 넣지 않는다
    assert CD.get_note_text(cross, CD.NOTE_SOAP) == ""
    status = CD.get_note_text(cross, CD.NOTE_STATUS)
    assert "Tx. Period: 0 month" in status and "U:" not in status
    # 차수 표시는 노트와 별개로 자기 박스에 적힌다
    assert "초진" in CD.get_note_text(cross, CD.NOTE_DATE)


@pytest.mark.skipif(not M.CASE_ANCHORS, reason="케이스 양식이 없습니다")
def test_untouched_notes_start_empty(sid, patient):
    """한 칸도 안 채우면 빈 칸으로 남는다.

    예전에는 양식 슬라이드의 텍스트박스를 글자째 복제해서 안내문("n) dx. consult")이
    남았다. 그러려면 양식이 **남의 진료기록 슬라이드**를 들고 다녀야 했고(s)/p) 칸이
    거기에만 있었다), 그 파일이 배포되면 기록이 함께 나갔다. 이제 칸은 좌표로 만들고
    글자는 임상의가 채운다.
    """
    assert client.post(f"/api/commit/{sid}?allow_missing=true").status_code == 200
    prs = T.load_presentation(patient / N.ppt_filename(IDS, M.cfg.naming.ppt_pattern))
    cross = prs.slides[CD.cross_slide_index(prs)]
    assert (CD.get_note_text(cross, CD.NOTE_NEXT) or "").strip() == ""
    # 다섯 칸이 모두 만들어져 있어야 한다 — 적을 곳이 없으면 노트가 사라진다
    names = {sh.name for sh in cross.shapes}
    assert set(CD.NOTE_ORDER) <= names


# ── 박스 번호(①~⑤)와 좌하단 칸 ──────────────────────────────────────────────
def test_note_key_splits_left_column_by_height():
    """왼쪽 열은 높이로 갈린다 — ①(날짜) / ②(s,p) / ④(좌하단)."""
    W = 25.4
    assert CD.note_key(0.37, 0.63, W) == CD.NOTE_DATE      # ① 맨 위
    assert CD.note_key(0.15, 1.53, W) == CD.NOTE_SOAP      # ② 그 아래
    assert CD.note_key(0.15, 12.90, W) == CD.NOTE_LL       # ④ 십자뷰 아래칸
    assert CD.note_key(17.44, 0.63, W) == CD.NOTE_STATUS   # ③
    assert CD.note_key(16.88, 12.58, W) == CD.NOTE_NEXT    # ⑤
    print("PASS 위치 → 박스 판정 (①②③④⑤)")


def test_note_order_is_the_screen_numbering():
    assert CD.NOTE_ORDER == [CD.NOTE_DATE, CD.NOTE_SOAP, CD.NOTE_STATUS,
                             CD.NOTE_LL, CD.NOTE_NEXT]
    # 화면이 붙이는 번호와 같은 순서로 위치도 내려간다
    assert list(M.NOTE_BOXES) == CD.NOTE_ORDER, list(M.NOTE_BOXES)
    print("PASS 번호 순서 =", " ".join(f"{i}.{k}" for i, k in enumerate(CD.NOTE_ORDER, 1)))


def test_lower_left_box_has_a_place_even_though_template_lacks_it():
    """④는 양식에 도형이 없다 — 그래도 화면에 자리를 알려 줘야 고칠 수 있다."""
    r = M.NOTE_BOXES[CD.NOTE_LL]
    assert r["w"] > 0 and r["h"] > 0
    # 십자뷰 좌하단 칸(MASK_R3_L: 0,12.72 8.50x6.30) 안에 들어가야 한다
    assert r["x"] >= 0 and r["y"] >= 12.72 - 0.01
    assert r["x"] + r["w"] <= 8.50 + 0.01
    print(f"PASS ④ 좌하단 자리 = {r['x']},{r['y']} {r['w']}x{r['h']}cm")


def test_add_note_box_creates_a_real_textbox():
    """기증할 도형이 없을 때 실제로 만들어지는지."""
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    assert CD.add_note_box(slide, CD.NOTE_LL, CD.NOTE_LL_WINDOW)
    made = [sh for sh in slide.shapes if sh.name == CD.NOTE_LL]
    assert len(made) == 1, made
    assert made[0].has_text_frame and made[0].text_frame.text == ""
    print("PASS ④ 텍스트박스 생성")


def test_blank_lines_survive_rendering():
    """양식의 빈 줄은 글의 시작 높이를 잡는다 — 다듬어 없애면 안 된다."""
    out = M._render_notes({"tx_period": "3 month"},
                          {"B": "U: {rx_period}\n\n\nTx. Period: {tx_period}"})
    # 칸이 빈 줄(U:)만 사라지고 빈 줄 두 개는 남는다
    assert out["B"] == "\n\nTx. Period: 3 month", repr(out["B"])
    print("PASS 빈 줄 보존, 빈 칸 줄만 제거")


def test_trailing_date_is_written_smaller():
    """양식은 본문 15pt / 줄 끝 날짜 9pt 다 — 갈아끼울 때도 그 구분을 지킨다."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Emu, Pt
    prs = Presentation()
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    tb = sl.shapes.add_textbox(Emu(0), Emu(0), Emu(3000000), Emu(2000000))
    tb.name = CD.NOTE_STATUS
    seed = tb.text_frame.paragraphs[0].add_run()
    seed.text = "seed"
    seed.font.size = Pt(15)
    seed.font.color.rgb = RGBColor.from_string("FFFF00")

    CD.set_note_text(sl, CD.NOTE_STATUS,
                     "Tx. Period: 23 month\nRx. Period: 23 month (24.08.12)", small_pt=9)
    paras = tb.text_frame.paragraphs
    # 날짜가 없는 줄은 한 덩어리, 본문 크기 그대로
    assert [r.text for r in paras[0].runs] == ["Tx. Period: 23 month"]
    assert paras[0].runs[0].font.size.pt == 15
    # 날짜가 붙은 줄은 둘로 갈리고 뒤쪽만 작다
    sizes = [(r.text, r.font.size.pt) for r in paras[1].runs]
    assert sizes == [("Rx. Period: 23 month", 15.0), (" (24.08.12)", 9.0)], sizes
    # 색은 원래 것을 그대로 물려받는다
    assert str(paras[1].runs[1].font.color.rgb) == "FFFF00"
    # 읽을 때는 한 줄로 합쳐진다(기존 계약 유지)
    assert CD.get_note_text(sl, CD.NOTE_STATUS).endswith("(24.08.12)")
    print("PASS 줄 끝 날짜만 9pt, 본문·색은 유지")


def test_paren_only_line_is_not_split():
    """본문 없이 괄호만 있는 줄은 통째로 본문 크기다 — 잘못 쪼개면 안 된다."""
    from pptx import Presentation
    from pptx.util import Emu, Pt
    prs = Presentation()
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    tb = sl.shapes.add_textbox(Emu(0), Emu(0), Emu(3000000), Emu(2000000))
    tb.name = CD.NOTE_SOAP
    seed = tb.text_frame.paragraphs[0].add_run()
    seed.text = "seed"
    seed.font.size = Pt(12)

    CD.set_note_text(sl, CD.NOTE_SOAP, "(메모만 있는 줄)", small_pt=9)
    runs = tb.text_frame.paragraphs[0].runs
    assert [r.text for r in runs] == ["(메모만 있는 줄)"], [r.text for r in runs]
    assert runs[0].font.size.pt == 12
    print("PASS 괄호만 있는 줄은 안 쪼갠다")


def test_Rx는_초진일을_자동으로_넣지_않는다():
    """Rx 기준일도 Tx·App 과 같은 규칙 — 고르기 전에는 괄호가 없다.

    예전에는 Rx 만 초진일을 기본 기준일로 박아 넣어서, 기준일 체크를 풀어도
    날짜가 남았다 (2026-08-14 결정으로 없앰).
    """
    import main as _M
    import naming as _N
    ids = _N.Identifiers("홍길동", "", "12345")
    s = _M.Session("first", ids, "A")
    s.first_date = "26.08.14"
    status = _M._note_text(s)["NOTE_STATUS"]
    assert "Rx. Period: 0 month\n" in status + "\n", status
    assert "(26.08.14)" not in status, status
    # 기준일을 고르면 그때 괄호가 붙는다 — Tx·App 과 같은 동작
    s.period_start["rx"] = "26.08.14"
    assert "Rx. Period: 0 month (26.08.14)" in _M._note_text(s)["NOTE_STATUS"]


def test_TxRxApp_상자는_줄_앞_간격을_갖고_태어난다():
    """초진 덱은 물려받을 원본이 없다 — 세 줄이 붙어 나오지 않게 기본값을 새긴다.

    간격은 문단이 아니라 상자의 목록서식(lstStyle)에 둔다. 글을 갈아끼워도(문단을
    새로 만들어도) 남고, 글자 크기를 바꿔도 비율이 유지된다.
    """
    import re
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from lxml import etree

    prs = Presentation()
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    CD.add_note_box(sl, CD.NOTE_STATUS, {"x": 1, "y": 1, "w": 8, "h": 5})
    CD.add_note_box(sl, CD.NOTE_SOAP, {"x": 1, "y": 7, "w": 8, "h": 3})
    CD.set_note_text(sl, CD.NOTE_STATUS, "Tx. Period: 0 month\nRx. Period: 3 month")

    def spc(name):
        sh = next(x for x in sl.shapes if x.name == name)
        lst = sh.text_frame._txBody.find(qn("a:lstStyle"))
        if lst is None:
            return None
        m = re.search(r'<a:spcBef>\s*<a:spcPct val="(\d+)"',
                      etree.tostring(lst, encoding="unicode"))
        return m.group(1) if m else None

    assert spc(CD.NOTE_STATUS) == CD.STATUS_SPACE_PCT   # 글을 쓴 뒤에도 남는다
    assert spc(CD.NOTE_SOAP) is None                    # 다른 칸은 종전 그대로


def test_수제_덱의_기간_상자도_읽는다():
    """이름 규약이 없는 상자를 **글 내용**으로 찾는다 — 라벨과 같은 방식.

    이 폴백이 없으면 수제 PPT 에서 기준일 이력이 통째로 비어, Tx/Rx/App 이 차수마다
    0 month 로 다시 시작한다(이어붙인 슬라이드에서 기간이 어긋난 원인).
    """
    from pptx import Presentation
    from pptx.util import Emu
    import ppt_reader as _Rd

    prs = Presentation()
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    box = sl.shapes.add_textbox(Emu(0), Emu(0), Emu(3000000), Emu(2000000))
    box.name = "TextBox 7"                       # 수제 — 규약명이 아니다
    box.text_frame.text = ("Tx. Period: 3 month (24.06.05)\n"
                           "Rx. Period: 3 month (24.05.01)")
    assert "Rx. Period" in _Rd._status_text(sl)

    vs = _Rd.VisitSlide(slide_index=0, visit="B", date="24.09.04", kind="revisit",
                        status_text=_Rd._status_text(sl))
    hist = M._period_history([vs], "Rx. Period")
    assert hist["dates"] == ["24.05.01"], hist
    assert hist["last"] == "Rx. Period: 3 month (24.05.01)"


# ── 슬라이드 크기가 양식과 다른 덱 ─────────────────────────────────────────
# 자리표는 양식(25.4x19.05) 좌표다. 직전 차수에 없던 칸은 여기서 새로 만들어지는데,
# 좌표를 그대로 쓰면 더 작은 덱에서 오른쪽·아래 칸이 슬라이드 밖으로 밀려난다.
# 이 칸들은 원래 모서리에 붙어 있는 것이라, 옮길 것은 좌표가 아니라 여백이다.
def test_슬라이드가_같으면_자리표가_그대로다():
    import case_deck as CD

    ref = M.CASE_SLIDE_CM
    for win in (M.NOTE_BOXES or {}).values():
        got = CD.anchored_window(win, ref, ref)
        assert (got["x"], got["y"]) == (win["x"], win["y"])


def test_작은_덱에서는_모서리_여백을_지킨다():
    import case_deck as CD

    ref = M.CASE_SLIDE_CM
    now = (ref[0] - 0.534, ref[1] - 0.534)
    for name, win in (M.NOTE_BOXES or {}).items():
        got = CD.anchored_window(win, ref, now)
        assert got["x"] + got["w"] <= now[0] + 1e-6, name    # 오른쪽으로 안 넘친다
        assert got["y"] + got["h"] <= now[1] + 1e-6, name    # 아래로도
        assert (got["w"], got["h"]) == (win["w"], win["h"]), name   # 크기는 그대로
        # 붙어 있던 모서리에서의 여백이 유지된다
        if win["x"] + win["w"] / 2 > ref[0] / 2:
            assert abs((now[0] - (got["x"] + got["w"]))
                       - (ref[0] - (win["x"] + win["w"]))) < 2e-3, name
        else:
            assert got["x"] == win["x"], name


def test_크기를_모르면_손대지_않는다():
    """양식이 없는 설치본 등 — 알 수 없으면 종전 좌표 그대로 둔다."""
    import case_deck as CD

    win = {"x": 17.05, "y": 12.87, "w": 8.2, "h": 6.0}
    assert CD.anchored_window(win, None, (24.0, 18.0)) == win
    assert CD.anchored_window(win, (25.4, 19.05), None) == win
