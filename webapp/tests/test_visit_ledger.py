"""차수 장부 — **십자뷰 슬라이드의 라벨** 하나로 통일한다.

라벨 확인 → 사진 기하 확인 순서로 판정한다. 라벨은 있어도 십자뷰가 아닌 장
(얼굴·엑스레이·사진 4장짜리)은 차수가 아니다 — 그런 장의 라벨이 차수를 가져가면
겹쳐보기·정합이 엉뚱한 회차를 기준으로 삼는다. 다만 십자뷰를 하나도 못 알아본
덱은 이력이 통째로 사라지므로 그 덱만 라벨 장 전체로 물러난다(fallback).
"""
import io
import os
import sys

import numpy as np
import pytest
from PIL import Image

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "backend"))

import main                                                       # noqa: E402
import ppt_reader as Rd                                           # noqa: E402


def _png(w=400, h=300) -> io.BytesIO:
    buf = io.BytesIO()
    Image.fromarray(np.full((h, w, 3), 128, np.uint8)).save(buf, "PNG")
    buf.seek(0)
    return buf


def _slide(prs, label: str, pics: int, width_cm: float = 8.4):
    """라벨 한 줄 + 사진 n장짜리 슬라이드."""
    from pptx.util import Emu
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    tb = sl.shapes.add_textbox(Emu(0), Emu(0), Emu(3000000), Emu(500000))
    tb.text_frame.text = label
    w = int(round(width_cm * main.EMU_PER_CM))
    for i in range(pics):
        sl.shapes.add_picture(_png(), Emu(i * w), Emu(1000000), width=Emu(w))
    return sl


@pytest.fixture
def prs():
    from pptx import Presentation
    return Presentation()


def test_십자뷰만_차수로_센다(prs):
    """사진 4장짜리 라벨 장은 제외되고, 왜 제외됐는지 알려준다."""
    _slide(prs, "24.06.05 (초진 A)", 5)      # 십자뷰
    _slide(prs, "24.08.14 (재진)", 4)        # 얼굴 장 등 — 차수 아님
    _slide(prs, "24.09.04 (재진 B)", 5)      # 십자뷰
    out = Rd.scan_ppt_visits(prs, main.cfg)
    assert [v["visit"] for v in out["visits"]] == ["A", "B"]
    assert [v["slide_no"] for v in out["visits"]] == [1, 3]
    assert out["excluded"] == [{"slide_no": 2, "visit": None, "date": "24.08.14"}]
    assert out["fallback"] is False


def test_글자없는_십자뷰는_날짜순으로_이어받는다(prs):
    """"(재진)" 라벨 — 글자를 아는 마지막 차수 뒤로 이어 붙인다."""
    _slide(prs, "24.06.05 (초진 A)", 5)
    _slide(prs, "24.09.04 (재진)", 5)
    _slide(prs, "24.12.01 (재진)", 5)
    out = Rd.scan_ppt_visits(prs, main.cfg)
    assert [v["visit"] for v in out["visits"]] == ["A", "B", "C"]


def test_글자차수_사이에_끼는_날짜는_매기지_않는다(prs):
    """어느 글자인지 확정할 수 없다 — 잘못 매기면 정합 기준이 엉뚱해진다."""
    _slide(prs, "24.06.05 (초진 A)", 5)
    _slide(prs, "24.12.01 (재진 C)", 5)
    _slide(prs, "24.09.04 (재진)", 5)        # A 와 C 사이 — 그대로 둔다
    out = Rd.scan_ppt_visits(prs, main.cfg)
    assert [v["visit"] for v in out["visits"]] == ["A", "C"]


def test_십자뷰를_못_알아본_덱은_라벨로_물러난다(prs):
    """사진이 문턱(8cm)에 못 미치는 수제 덱 — 이력이 통째로 사라지면 안 된다."""
    _slide(prs, "24.06.05 (초진 A)", 5, width_cm=6.0)
    _slide(prs, "24.09.04 (재진 B)", 5, width_cm=6.0)
    out = Rd.scan_ppt_visits(prs, main.cfg)
    assert [v["visit"] for v in out["visits"]] == ["A", "B"]
    assert out["fallback"] is True
    assert out["excluded"] == []


def test_같은_차수가_여러_장이어도_안전하다(prs):
    """글이 많아 슬라이드를 복제해 쓰는 관행 — 차수는 하나로 세고, 새 장은 뒤로."""
    _slide(prs, "24.06.05 (초진 A)", 5)
    _slide(prs, "24.08.12 (재진 C)", 5)
    _slide(prs, "24.08.12 (재진 C)", 5)      # 같은 차수 복제본
    out = Rd.scan_ppt_visits(prs, main.cfg)
    letters = sorted({v["visit"] for v in out["visits"]})
    assert letters == ["A", "C"]              # 중복은 한 차수로
    assert main._revisit_insert_index(prs) == 3   # 마지막 복제본 **뒤**
    seen = Rd.read_all_visits(prs, main.cfg, main.PPC)
    refs = Rd.references_for_registration(seen)
    assert sorted(next(iter(refs.values()))) == ["A", "C"]   # 겹쳐보기도 중복 없음


def test_라벨_없는_장은_차수가_아니다(prs):
    _slide(prs, "환자정보", 5)
    out = Rd.scan_ppt_visits(prs, main.cfg)
    assert out["visits"] == [] and out["excluded"] == []


def test_삽입_자리는_날짜가_아니라_차수_글자로_고른다(prs):
    """라벨 날짜는 손으로 적다 보니 오타가 난다 — J 가 K 보다 뒤 날짜인 덱이 있었다.

    날짜로 고르면 새 장이 K 앞으로 들어간다. 차수 글자는 순서 그 자체라 안 흔들린다.
    """
    _slide(prs, "25.03.17 (재진 I)", 5)
    _slide(prs, "25.05.19 (재진 J)", 5)     # 오타 — K 보다 뒤 날짜
    _slide(prs, "25.04.30 (재진 K)", 5)
    assert main._revisit_insert_index(prs) == 3      # K(3번 장) 뒤 = 0-기반 3


def test_조금_작게_놓인_사진도_십자뷰로_본다(prs):
    """수제 덱마다 사진을 조금씩 작게 잡는다. 8.0cm 문턱에서는 7.9cm 로 놓인 덱이
    통째로 "십자뷰 아님" 으로 밀려나 차수 이력이 비었다."""
    _slide(prs, "24.06.05 (초진 A)", 5, width_cm=7.9)
    out = Rd.scan_ppt_visits(prs, main.cfg)
    assert [v["visit"] for v in out["visits"]] == ["A"]
    assert out["excluded"] == [] and out["fallback"] is False


def test_로고나_썸네일은_여전히_안_센다(prs):
    """문턱을 낮춰도 작은 그림과는 확연히 갈려야 한다 — 안 그러면 아무 장이나
    십자뷰가 되고, 차수 장부가 통째로 어긋난다."""
    _slide(prs, "24.06.05 (초진 A)", 5, width_cm=Rd.CROSS_MIN_W_CM - 0.1)
    out = Rd.scan_ppt_visits(prs, main.cfg)
    # 십자뷰로 인정되는 장이 하나도 없으니 라벨 장 전체로 물러난다(폴백)
    assert out["fallback"] is True


def test_표기_지문은_직전_차수에서_온다(prs):
    """장 순서가 아니라 **차수 글자**가 기준이다.

    확인 줄에서 새 장을 넣을 자리를 직접 고를 수 있고 사람이 파워포인트에서 장을
    옮기기도 해서, 순서와 차수가 어긋난 덱이 실제로 나온다. 그때 정합은 D 를
    기준 삼는데 표기는 C 를 따라 쓰면 한 화면이 서로 다른 두 장을 가리킨다.
    """
    _slide(prs, "24.05.17 (재진 B)", 5)
    _slide(prs, "24.07.18 (F/U D)", 5)      # 차수 글자가 가장 큼 — 정합 기준
    _slide(prs, "24.06.14 (재진 C)", 5)      # 장 순서로는 이쪽이 마지막
    out = Rd.scan_ppt_visits(prs, main.cfg)
    assert out["label_fp"]["word"] == "F/U"


def test_글자가_없는_덱은_순서로_물러난다(prs):
    """글자가 하나도 없으면 고를 기준이 없다 — 마지막 장을 쓴다."""
    _slide(prs, "24.05.17 (재진)", 5)
    _slide(prs, "24.06.14 (F/U)", 5)
    out = Rd.scan_ppt_visits(prs, main.cfg)
    assert out["label_fp"]["word"] == "F/U"
