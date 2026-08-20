"""겹쳐 놓인 십자뷰에서 **정면이 맨 위**여야 한다.

수제 십자뷰는 사진끼리 조금씩 포개진다 — 실측한 덱 대부분이 상악과 정면이
1.6mm 가량 겹쳐 있었다. 예전에는 `slot_names` 순으로 넣어 **상악이 맨 위**가
됐고(먼저 넣은 것이 맨 위로 간다), 그만큼 정면 사진 윗변이 가려졌다. 십자뷰의
중심은 정면이므로 그쪽이 위여야 한다.

실행: cd webapp && python -m pytest tests/test_z_order.py -q
"""
import io
import os
import sys

from PIL import Image
from pptx import Presentation
from pptx.util import Cm
from starlette.testclient import TestClient

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "backend"))

import main as M  # noqa: E402
import template as T  # noqa: E402
import ppt_writer as W  # noqa: E402
from coords import emu_to_cm  # noqa: E402

client = TestClient(M.app)

NAME, HOSP, ORTHO = "겹침검사", "123456789", "12345"
FOLDER = f"{NAME}_{HOSP}_{ORTHO}"
DECK = f"{NAME}({ORTHO}).pptx"

W_CM, H_CM = 8.4, 6.3
OVERLAP = 0.16          # 실측 덱에서 상악과 정면이 포개져 있던 만큼


def _png(tmp_path):
    p = tmp_path / "dot.png"
    Image.new("RGB", (8, 6), (90, 90, 90)).save(p)
    return str(p)


def _jpg(seed: int) -> bytes:
    im = Image.new("RGB", (1200, 900), (40 + seed * 30, 90, 200 - seed * 30))
    buf = io.BytesIO()
    im.save(buf, "JPEG", quality=90)
    return buf.getvalue()


def _cross_deck(path, png, order):
    """사진이 서로 포개진 십자뷰 한 장짜리 덱. `order` 순서로 놓는다(뒤가 위)."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Cm(33.87), Cm(19.05)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cx = emu_to_cm(prs.slide_width) / 2 - W_CM / 2
    cy = emu_to_cm(prs.slide_height) / 2 - H_CM / 2
    at = {"SLOT_FRONT": (cx, cy),
          "SLOT_UPPER": (cx, cy - H_CM + OVERLAP),
          "SLOT_LOWER": (cx, cy + H_CM - OVERLAP),
          "SLOT_LEFT":  (cx - W_CM + OVERLAP, cy),
          "SLOT_RIGHT": (cx + W_CM - OVERLAP, cy)}
    for slot in order:
        x, y = at[slot]
        slide.shapes.add_picture(png, Cm(x), Cm(y), Cm(W_CM), Cm(H_CM))
    tb = slide.shapes.add_textbox(Cm(0.5), Cm(0.2), Cm(10), Cm(0.9))
    tb.text_frame.text = "24.06.05 (초진 A)"
    prs.save(str(path))


def _photo_z(slide) -> list[str]:
    """앱이 넣은 사진의 z-순서 (아래 → 위). 도형 이름이 곧 자리다."""
    out = []
    for sh in slide.shapes:
        nm = str(getattr(sh, "name", ""))
        if nm.startswith(W.PHOTO_NAME_PREFIX):
            out.append(nm[len(W.PHOTO_NAME_PREFIX):])
    return out


def _commit_one_visit(root, tmp_path, source_order):
    """`source_order` 로 놓인 덱에 한 차수를 이어붙이고, 새 장의 z-순서를 준다."""
    d = root / FOLDER
    if not d.exists():
        d.mkdir(parents=True)
    _cross_deck(d / DECK, _png(tmp_path), source_order)

    r = client.post("/api/session", json={"folder": FOLDER})
    assert r.status_code == 200, r.text
    sid = r.json()["session_id"]
    s = M.SESSIONS[sid]

    files = [("files", (f"p{i}.jpg", io.BytesIO(_jpg(i)), "image/jpeg"))
             for i in range(5)]
    assert client.post(f"/api/photos/{sid}", files=files).status_code == 200
    # 분류기에 기대지 않는다 — 여기서 보려는 것은 배정이 아니라 쌓는 순서다
    for slot, photo in zip(M.cfg.ppt.slot_names, s.photos):
        M._put(s, photo, slot, at=0)

    r = client.post(f"/api/commit/{sid}?allow_missing=true")
    assert r.status_code == 200, r.text
    prs = T.load_presentation(d / r.json()["ppt"])
    fresh = [_photo_z(sl) for sl in prs.slides if _photo_z(sl)]
    assert fresh, "사진이 들어간 장이 없다"
    return fresh


def test_정면이_맨_위에_들어간다(_isolate_paths, tmp_path):
    """사람이 정면을 맨 위에 둔 덱 — 새 장도 같아야 한다."""
    order = ["SLOT_LEFT", "SLOT_RIGHT", "SLOT_UPPER", "SLOT_LOWER", "SLOT_FRONT"]
    for z in _commit_one_visit(_isolate_paths, tmp_path, order):
        assert z[-1] == "SLOT_FRONT", z
        assert len(z) == 5, z


def test_직전_장이_상악을_위에_뒀어도_정면이_위다(_isolate_paths, tmp_path):
    """앱이 예전에 반대로 쌓아 둔 덱에서도 새 장은 정면이 위다.

    직전 차수만 따라가면 그 잘못이 영원히 대물림된다 — 실측 덱 10개 중 4개가
    수제 장은 전부 정면이 위인데 앱이 쓴 장만 상악이 위였다.
    """
    order = ["SLOT_LOWER", "SLOT_RIGHT", "SLOT_LEFT", "SLOT_FRONT", "SLOT_UPPER"]
    for z in _commit_one_visit(_isolate_paths, tmp_path, order):
        assert z[-1] == "SLOT_FRONT", z


def test_배경판은_제_사진_바로_아래에_깔린다(_isolate_paths, tmp_path):
    """배경판(레터박스)이 제 사진보다 위로 오면 사진이 덮인다."""
    order = ["SLOT_LEFT", "SLOT_RIGHT", "SLOT_UPPER", "SLOT_LOWER", "SLOT_FRONT"]
    d = _isolate_paths / FOLDER
    d.mkdir(parents=True)
    _cross_deck(d / DECK, _png(tmp_path), order)
    sid = client.post("/api/session", json={"folder": FOLDER}).json()["session_id"]
    s = M.SESSIONS[sid]
    files = [("files", (f"p{i}.jpg", io.BytesIO(_jpg(i)), "image/jpeg"))
             for i in range(5)]
    client.post(f"/api/photos/{sid}", files=files)
    for slot, photo in zip(M.cfg.ppt.slot_names, s.photos):
        M._put(s, photo, slot, at=0)
    r = client.post(f"/api/commit/{sid}?allow_missing=true")
    assert r.status_code == 200, r.text
    prs = T.load_presentation(d / r.json()["ppt"])
    for sl in prs.slides:
        names = [str(sh.name) for sh in sl.shapes]
        if not any(n.startswith(W.PHOTO_NAME_PREFIX) for n in names):
            continue
        for slot in M.cfg.ppt.slot_names:
            pi = names.index(W.PHOTO_NAME_PREFIX + slot)
            bi = names.index(W.backdrop_shape_name(slot))
            assert bi < pi, f"{slot}: 배경판이 사진 위에 있다"
