"""직전 차수 슬라이드 복사 — 선만 / 전부 / 복사 안 함.

정중선 같은 기준선은 매 차수 같은 자리를 가리키므로 따라오는 편이 맞다(기본값).
"전부"는 글상자와 그 안의 글까지 가져온다 — 지난 차수 내용을 이어 고쳐 쓰는
방식이다. "복사 안 함"은 상속으로 딸려 온 자유 기입 글도 지운다.
"""
import os
import sys

import pytest

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "backend"))

import case_deck as CD                                            # noqa: E402
import main as M                                                  # noqa: E402
import template as T                                              # noqa: E402


def _slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _textbox(slide, name, text=""):
    from pptx.util import Emu
    tb = slide.shapes.add_textbox(Emu(500000), Emu(500000), Emu(2000000), Emu(800000))
    tb.name = name
    tb.text_frame.text = text
    return tb


def _line(slide, name="정중선"):
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_CONNECTOR
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Emu(1000000), Emu(0), Emu(1000000), Emu(4000000))
    ln.name = name
    return ln


@pytest.fixture
def deck():
    """직전 장(선 + 자유 메모 + 규약 상자) / 새 장(빈 규약 상자)."""
    from pptx import Presentation
    prs = Presentation()
    src, dst = _slide(prs), _slide(prs)
    _line(src)
    _textbox(src, "메모", "발치 부위 확인")
    _textbox(src, CD.NOTE_LL, "지난 차수 좌하단")
    _textbox(src, CD.NOTE_NEXT, "지난 차수 우하단")
    _textbox(src, M.cfg.ppt.info_box_name, "24.08.12 (재진 C)")
    for nm in (CD.NOTE_LL, CD.NOTE_NEXT):
        _textbox(dst, nm, "")
    return src, dst


def _names(slide):
    return [str(sh.name) for sh in slide.shapes]


def test_선만_복사(deck):
    src, dst = deck
    M._inherit_shapes(src, dst, "lines")
    assert "정중선" in _names(dst)
    assert "메모" not in _names(dst)                  # 글상자는 안 온다
    assert T.find_shape(dst, CD.NOTE_LL).text_frame.text.strip() == ""


def test_전부_복사는_글까지(deck):
    src, dst = deck
    M._inherit_shapes(src, dst, "all")
    assert "정중선" in _names(dst) and "메모" in _names(dst)
    # 자유 기입 상자는 **글 내용**이 물려 내려온다
    assert T.find_shape(dst, CD.NOTE_LL).text_frame.text == "지난 차수 좌하단"
    assert T.find_shape(dst, CD.NOTE_NEXT).text_frame.text == "지난 차수 우하단"
    # 날짜/차수 상자는 이 경로로 오지 않는다 — 매 차수 새로 쓴다
    assert _names(dst).count(M.cfg.ppt.info_box_name) == 0


def test_복사_안_함(deck):
    src, dst = deck
    _textbox(dst, CD.NOTE_SOAP, "상속으로 딸려 온 지난 글")
    M._inherit_shapes(src, dst, "none")
    assert "정중선" not in _names(dst) and "메모" not in _names(dst)
    # 상속으로 딸려 온 글도 지운다 — 안 지우면 설정이 거짓말이 된다
    assert T.find_shape(dst, CD.NOTE_SOAP).text_frame.text.strip() == ""


def test_기본값은_선만_복사(tmp_path, monkeypatch):
    monkeypatch.setattr(M, "SETTINGS_FILE", tmp_path / "settings.json", raising=False)
    assert M._copy_shapes() == "lines"


def test_이름_없는_규약_상자는_두_번_넣지_않는다():
    """수제 덱의 라벨·기간·노트 상자는 이미 통째로 물려받은 것들이다.

    이름으로만 걸러내면 규약명이 없는 그 상자들이 '전부 복사' 로 한 벌 더 얹혀
    글자가 겹친다. 이름이 아니라 **역할**(내용·자리)로 가려낸다.
    """
    from pptx import Presentation
    from pptx.util import Emu
    import main as _M
    import ppt_reader as _Rd

    EMU, SW = _M.EMU_PER_CM, 25.4
    prs = Presentation()
    src, dst = (prs.slides.add_slide(prs.slide_layouts[6]) for _ in range(2))

    def box(sl, name, text, x, y):
        b = sl.shapes.add_textbox(Emu(int(x * EMU)), Emu(int(y * EMU)),
                                  Emu(int(8 * EMU)), Emu(int(2 * EMU)))
        b.name, b.text_frame.text = name, text
        return b

    box(src, "TextBox 1", "24.09.04 (재진 C)", 0.4, 0.4)              # 라벨
    box(src, "TextBox 2", "Tx. Period: 3 month (24.06.05)", 17.1, 0.6)  # 기간
    box(src, "TextBox 3", "s) n/s", 0.1, 1.5)                          # 좌상단 노트
    _line(src)
    # 새 슬라이드에는 규약 상자가 이미 있다(통째 복사 상속의 결과)
    box(dst, _M.cfg.ppt.info_box_name, "26.08.14 (재진 D)", 0.4, 0.4)
    box(dst, CD.NOTE_STATUS, "Tx. Period: 5 month (24.06.05)", 17.1, 0.6)

    assert _Rd.note_role(src.shapes[0], SW) == "label"
    assert _Rd.note_role(src.shapes[1], SW) == "status"

    # 실제 커밋은 "이번에 물려받은 역할" 을 함께 넘긴다
    _M._inherit_shapes(src, dst, "all", SW, {"label", "status", CD.NOTE_SOAP})
    names = [str(sh.name) for sh in dst.shapes]
    assert "정중선" in names                       # 선은 넘어온다
    assert not [n for n in names if n.startswith("TextBox")], names
    assert names.count(_M.cfg.ppt.info_box_name) == 1
    assert names.count(CD.NOTE_STATUS) == 1


def test_물려받은_것이_없으면_글상자를_버리지_않는다():
    """앱이 만든 직전 슬라이드에는 물려받을 상자가 없다.

    그런데도 '노트 자리에 있다' 는 이유로 손으로 그려 둔 글상자를 버렸다 —
    전부 복사인데 아무것도 안 따라오는 것처럼 보였다.
    """
    from pptx import Presentation
    from pptx.util import Emu
    import main as _M

    EMU, SW = _M.EMU_PER_CM, 25.4
    prs = Presentation()
    src, dst = (prs.slides.add_slide(prs.slide_layouts[6]) for _ in range(2))
    for name, x, y in (("INFO_BOX", 0.4, 0.4), (CD.NOTE_STATUS, 17.1, 0.6)):
        b = src.shapes.add_textbox(Emu(int(x * EMU)), Emu(int(y * EMU)),
                                   Emu(int(8 * EMU)), Emu(int(2 * EMU)))
        b.name, b.text_frame.text = name, "24.09.04 (재진 C)"
    memo = src.shapes.add_textbox(Emu(int(0.2 * EMU)), Emu(int(13 * EMU)),
                                  Emu(int(6 * EMU)), Emu(int(2 * EMU)))
    memo.name, memo.text_frame.text = "TextBox 9", "손으로 적은 메모"

    _M._inherit_shapes(src, dst, "all", SW, set())      # 물려받은 역할 없음
    names = [str(sh.name) for sh in dst.shapes]
    assert "TextBox 9" in names, names                  # 그대로 따라온다
    assert "INFO_BOX" not in names and CD.NOTE_STATUS not in names
