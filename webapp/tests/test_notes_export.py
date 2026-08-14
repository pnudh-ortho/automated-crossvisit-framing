"""
노트가 PPT 에 실제로 **보이게** 남는가 + 양식 첫 장(환자정보) 자동 입력.

여기서 확인하는 것:

  - 초진 십자뷰: 노트 칸이 검은 마스크보다 **앞에** 있는가
    (뒤에 있으면 글자는 파일에 남아도 화면·인쇄물에서는 안 보인다)
  - 재진 슬라이드: 십자뷰 양식에는 없는 노트 칸이 만들어지고 글이 적히는가
  - 양식 첫 장의 "Hospital No./Case No./Pt. name/C/C" 가 환자 값으로 채워지는가
  - 빈 칸이 만든 껍데기("(/)" 같은 것)를 남기지 않는가

실행: cd webapp && python -m pytest tests/test_notes_export.py -q
"""
import io
import os
import shutil
import sys

import cv2
import numpy as np
import pytest
from starlette.testclient import TestClient

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "backend"))

import case_deck as CD  # noqa: E402
import main as M  # noqa: E402
import naming as N  # noqa: E402
import template as T  # noqa: E402

client = TestClient(M.app)

NAME, HOSP, ORTHO = "노트내보내기", "111222777", "54325"
IDS = N.Identifiers(NAME, HOSP, ORTHO)
FOLDER = N.folder_name(IDS, M.cfg.naming.folder_pattern)
PPT = N.ppt_filename(IDS, M.cfg.naming.ppt_pattern)

needs_case = pytest.mark.skipif(not M.CASE_ANCHORS, reason="케이스 양식이 없습니다")


@pytest.fixture
def patient():
    d = M.ROOT / FOLDER
    shutil.rmtree(d, ignore_errors=True)
    d.mkdir(parents=True)
    yield d
    shutil.rmtree(d, ignore_errors=True)


def _photo_blob() -> bytes:
    img = np.full((1200, 900, 3), 40, np.uint8)
    cv2.putText(img, "x", (100, 600), cv2.FONT_HERSHEY_SIMPLEX, 4, (255, 255, 255), 8)
    return cv2.imencode(".jpg", img)[1].tobytes()


def _stage_photo(sid: str, name: str = "a.jpg", slot: str = "SLOT_FRONT") -> None:
    """사진 한 장을 담고 슬롯에 붙인다.

    분류 결과에 기대지 않는다 — 실제 모델이 있는 설치본에서는 합성 이미지가
    어느 슬롯으로 갈지 알 수 없고, 슬롯이 비면 차수 파일이 안 생겨 재진 판정이
    흔들린다. 여기서 보려는 것은 분류가 아니라 노트다.
    """
    r = client.post(f"/api/photos/{sid}",
                    files=[("files", (name, io.BytesIO(_photo_blob()), "image/jpeg"))])
    client.post(f"/api/classify/{sid}")
    pid = r.json()["photos"][0]["id"]
    client.post("/api/assign", json={"session_id": sid, "photo_id": pid, "slot": slot, "at": 0})


def _session_with_photo(payload: dict) -> str:
    """세션을 열고 사진 한 장을 슬롯에 담는다. payload 는 /api/session 본문."""
    sid = client.post("/api/session", json=payload).json()["session_id"]
    _stage_photo(sid)
    return sid


def _commit(sid) -> None:
    r = client.post(f"/api/commit/{sid}?allow_missing=true")
    assert r.status_code == 200, r.text


def _shape_names(slide) -> list[str]:
    return [sh.name for sh in slide.shapes]


# ── 초진: 노트가 마스크에 가리지 않는가 ───────────────────────────────────────
@needs_case
def test_note_boxes_sit_in_front_of_the_masks(patient):
    """마스크(MASK_*)는 확정 슬라이드에서 **제거된다** (2026-08-12 결정).

    사진을 창 크기로 구워 넣어 초과가 없고 배경이 검정이라 마스크는 잉여였다.
    수제 레이아웃을 상속하면 템플릿 좌표의 마스크가 사진을 가리기까지 했다.
    노트 칸은 그대로 있어야 한다 — 마스크가 사라져도 글이 보이는 구조다.
    """
    sid = _session_with_photo({"folder": FOLDER})
    client.post("/api/notes", json={"session_id": sid, "values": {"subj": "n/s"}})
    _commit(sid)

    prs = T.load_presentation(patient / PPT)
    cross = prs.slides[CD.cross_slide_index(prs)]
    names = _shape_names(cross)
    masks = [n for n in names if n.startswith("MASK_")]
    assert not masks, f"마스크가 남아 있다: {masks}"
    for box in CD.NOTE_ORDER:
        assert box in names, f"{box} 가 없다: {names}"


# ── 재진: 노트 칸이 아예 없던 슬라이드 ────────────────────────────────────────
@needs_case
def test_revisit_slide_gets_note_boxes_and_keeps_the_text(patient):
    """십자뷰 양식에는 노트 칸이 없다 — 만들어 주지 않으면 노트가 조용히 사라진다."""
    _commit(_session_with_photo({"folder": FOLDER}))          # A 차수

    r = client.post("/api/session/revisit", json={"ppt_path": f"{FOLDER}/{PPT}"})
    assert r.status_code == 200, r.text
    sid2 = r.json()["session_id"]
    assert r.json()["visit"] == "B", r.json()

    _stage_photo(sid2, "b.jpg")
    client.post("/api/notes", json={"session_id": sid2,
                                    "values": {"tx_period": "1 month"}})
    _commit(sid2)

    prs = T.load_presentation(patient / PPT)
    # 이번 차수 슬라이드 = 방금 늘어난 마지막 십자뷰
    added = [sl for sl in prs.slides if CD.NOTE_SOAP in _shape_names(sl)]
    assert len(added) >= 2, "재진 슬라이드에 노트 칸이 생기지 않았다"
    new = added[-1]
    # 자유 기입 상자는 자리만 만들어지고 글은 앱이 넣지 않는다
    assert CD.get_note_text(new, CD.NOTE_SOAP) == ""
    assert "Tx. Period: 1 month" in CD.get_note_text(new, CD.NOTE_STATUS)
    # 날짜는 십자뷰 양식이 이미 들고 있는 INFO_BOX 가 맡는다 — 겹쳐 적지 않는다
    assert CD.NOTE_DATE not in _shape_names(new)
    assert "재진" in T.find_shape(new, M.cfg.ppt.info_box_name).text_frame.text


@needs_case
def test_box_overrides_reach_the_slide(patient):
    """판 위 오버레이에서 통째로 고쳐 쓴 글이 그대로 슬라이드에 남는가.

    화면이 이 경로(boxes)로 보내는데, 오래도록 같은 이름의 함수가 둘이라
    payload 가 버려지고 values 만 갔다 — 고쳐 쓴 노트가 PPT 에 없던 이유다.
    """
    sid = _session_with_photo({"folder": FOLDER})
    r = client.post("/api/notes", json={"session_id": sid,
                                        "boxes": {CD.NOTE_NEXT: "n) 손으로 고쳐 쓴 줄"}})
    assert r.status_code == 200, r.text
    assert r.json()["overrides"][CD.NOTE_NEXT] == "n) 손으로 고쳐 쓴 줄"
    _commit(sid)

    prs = T.load_presentation(patient / PPT)
    cross = prs.slides[CD.cross_slide_index(prs)]
    assert CD.get_note_text(cross, CD.NOTE_NEXT) == "n) 손으로 고쳐 쓴 줄"


# ── 양식 첫 장(환자정보) ──────────────────────────────────────────────────────
@needs_case
def test_patient_info_slide_is_filled_from_the_patient(patient):
    sid = _session_with_photo({"folder": FOLDER})
    client.post("/api/notes", json={"session_id": sid,
                                    "values": {"sex": "M", "age": "24y 3m",
                                               "cc": "앞니가 튀어나왔어요"}})
    _commit(sid)

    prs = T.load_presentation(patient / PPT)
    text = "\n".join(sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame)
    assert f"Hospital No. {HOSP}" in text, text
    assert f"Case No. {ORTHO}" in text, text
    assert f"Pt. name : {NAME} (M/24y 3m)" in text, text
    assert "C/C : 앞니가 튀어나왔어요" in text, text
    # 양식의 안내 문구는 갈아끼워진다
    assert "(M/y m)" not in text, text


@needs_case
def test_untouched_patient_info_keeps_the_template_hints(patient):
    """성별·나이·C/C 를 안 채우면 그 줄은 양식 그대로 둔다 — 안내가 사라지면
    나중에 무엇을 적을 자리인지 알 수 없다. 이름·번호는 이미 아는 값이라 채운다."""
    _commit(_session_with_photo({"folder": FOLDER}))

    prs = T.load_presentation(patient / PPT)
    text = "\n".join(sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame)
    assert f"Hospital No. {HOSP}" in text, text
    assert f"Pt. name : {NAME}" in text, text
    assert "(/)" not in text and "( )" not in text, f"빈 껍데기가 남았다: {text}"
    assert "C/C :" in text, text


# ── 글자 크기·글머리 점 ───────────────────────────────────────────────────────
@needs_case
def test_visit_label_is_not_shrunk(patient):
    """날짜 칸의 "(초진 A)" 는 날짜가 아니라 차수 표시다 — 본문과 같은 크기여야 한다.

    줄 끝 괄호를 무조건 9pt 로 줄이면 차수 글자만 깨알같이 작아진다.
    """
    _commit(_session_with_photo({"folder": FOLDER}))
    prs = T.load_presentation(patient / PPT)
    cross = prs.slides[CD.cross_slide_index(prs)]
    date_box = [sh for sh in cross.shapes if sh.name == CD.NOTE_DATE][0]
    runs = [r for p in date_box.text_frame.paragraphs for r in p.runs if r.text.strip()]
    assert "초진" in "".join(r.text for r in runs)
    sizes = {r.font.size.pt for r in runs if r.font.size}
    assert sizes and M.cfg.notes.date_pt not in sizes, \
        f"차수 표시가 작게 쓰였다: {[(r.text, r.font.size.pt) for r in runs]}"
    # 한편 상태 칸의 줄 끝 날짜는 여전히 작아야 한다(양식이 그렇다)
    assert M._small_pt(CD.NOTE_DATE) is None
    assert M._small_pt(CD.NOTE_STATUS) == M.cfg.notes.date_pt


@needs_case
def test_no_bullet_appears_in_front_of_the_note_lines(patient):
    """Tx./Rx./App. 앞에 글머리 점(•)이 붙으면 안 된다.

    양식의 문단은 저마다 buNone 을 들고 있는데, 글을 갈아끼울 때 문단서식까지
    새로 만들면 개체 틀의 목록서식(buChar •)을 물려받아 없던 점이 생긴다.
    """
    from pptx.oxml.ns import qn
    sid = _session_with_photo({"folder": FOLDER})
    client.post("/api/notes", json={"session_id": sid, "values": {
        "tx_period": "3 month", "rx_period": "3 month (26.05.02)", "subj": "n/s"}})
    _commit(sid)

    prs = T.load_presentation(patient / PPT)
    cross = prs.slides[CD.cross_slide_index(prs)]
    for name in CD.NOTE_ORDER:
        for sh in [s for s in cross.shapes if s.name == name]:
            for para in sh.text_frame.paragraphs:
                if not para.runs:
                    continue
                pPr = para._p.find(qn("a:pPr"))
                assert pPr is not None, f"{name}: 문단서식이 없어 목록서식을 물려받는다"
                assert pPr.find(qn("a:buChar")) is None, f"{name}: 글머리 점이 붙었다"
                assert pPr.find(qn("a:buAutoNum")) is None, f"{name}: 번호 글머리가 붙었다"
                assert pPr.find(qn("a:buNone")) is not None, \
                    f"{name}: 글머리 없음이 명시돼 있지 않다 — 물려받을 위험"


# ── 첫 장 화면 미리보기 ───────────────────────────────────────────────────────
@needs_case
def test_info_preview_matches_what_gets_written(patient):
    """화면이 보여주는 첫 장 글 = 확정 뒤 실제로 적히는 글."""
    sid = _session_with_photo({"folder": FOLDER})
    client.post("/api/notes", json={"session_id": sid,
                                    "values": {"sex": "F", "age": "19y 2m", "cc": "덧니"}})
    pi = client.get(f"/api/notes/{sid}").json()["patient_info"]
    assert pi["enabled"] and pi["slide"] == 1
    assert pi["fields"] == ["sex", "hospital_id", "age", "cc"]
    # 고칠 수 있는 박스가 정확히 하나 (우리가 채우는 줄을 든 박스)
    hot = [n for n, b in pi["boxes"].items() if b["editable"]]
    assert len(hot) == 1, pi["boxes"]
    shown = pi["preview"][hot[0]]
    _commit(sid)

    prs = T.load_presentation(patient / PPT)
    written = [sh for sh in prs.slides[0].shapes if sh.name == hot[0]][0].text_frame.text
    assert shown == written, f"화면과 결과가 다르다\n화면: {shown!r}\n결과: {written!r}"


# ── 줄 만들기 규칙 (파일을 안 건드리는 단위 검사) ─────────────────────────────
def test_empty_parenthetical_group_disappears():
    f = {"name": "김하늘", "sex": "", "age": ""}
    assert M._render_info_line(" {name} ({sex}/{age})", f) == " 김하늘"


def test_half_filled_parenthetical_keeps_only_what_is_there():
    assert M._render_info_line(" {name} ({sex}/{age})",
                               {"name": "김하늘", "sex": "M", "age": ""}) == " 김하늘 (M)"
    assert M._render_info_line(" {name} ({sex}/{age})",
                               {"name": "김하늘", "sex": "", "age": "24y"}) == " 김하늘 (24y)"


def test_notes_json_carries_the_small_date_size(patient):
    """줄 끝 날짜를 작게 쓰는 크기(9pt)를 화면도 알아야 양식과 같게 보인다.

    화면 오버레이는 이 값으로 날짜만 줄여 그린다 — 서버가 안 알려주면 화면은
    본문 크기로 그려서 슬라이드에서 보일 모습과 어긋난다.
    """
    sid = _session_with_photo({"folder": FOLDER})
    js = client.get(f"/api/notes/{sid}").json()
    assert js["date_pt"] == M.cfg.notes.date_pt == 9
    # 본문 크기는 양식에서 읽은 값이 그대로 온다 (우상단 Tx./Rx. = 15pt)
    assert js["layout"][CD.NOTE_STATUS]["font"]["size_pt"] == 15.0
