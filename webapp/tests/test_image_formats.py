"""
투입 이미지 포맷 검증.

카메라가 내놓는 .jpg 중에는 실제로 **MPO**(다중 프레임 JPEG)인 것이 있다.
확장자만 믿고 원본 바이트를 그대로 넘기면 투입·분류까지는 멀쩡히 지나가고
**확정 저장에서** python-pptx 가 거부해 통째로 롤백된다. 그 경로를 여기서 막는다.

실행: cd webapp && python -m pytest tests/test_image_formats.py -q
"""
import io
import os
import shutil
import sys

import numpy as np
import pytest
from PIL import Image
from starlette.testclient import TestClient

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "backend"))

import main as M  # noqa: E402
import naming as N  # noqa: E402

client = TestClient(M.app)

NAME, HOSP, ORTHO = "포맷검사", "111222777", "54325"
IDS = N.Identifiers(NAME, HOSP, ORTHO)
FOLDER = N.folder_name(IDS, M.cfg.naming.folder_pattern)


@pytest.fixture
def patient():
    d = M.ROOT / FOLDER
    shutil.rmtree(d, ignore_errors=True)
    d.mkdir(parents=True)
    yield d
    shutil.rmtree(d, ignore_errors=True)


def _mpo_bytes(w=900, h=1200):
    """카메라가 내놓는 것과 같은 다중 프레임 JPEG(MPO)."""
    a = Image.fromarray(np.uint8(np.random.rand(h, w, 3) * 255))
    b = Image.fromarray(np.uint8(np.random.rand(h, w, 3) * 255))
    buf = io.BytesIO()
    a.save(buf, format="MPO", append_images=[b])
    return buf.getvalue()


def test_mpo_is_rejected_by_pptx():
    """전제 확인 — 이 포맷이 정말 PPT 에 안 들어가는가."""
    from pptx import Presentation
    from pptx.util import Emu
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    with pytest.raises(Exception) as e:
        slide.shapes.add_picture(io.BytesIO(_mpo_bytes(60, 40)),
                                 Emu(0), Emu(0), Emu(100000), Emu(100000))
    assert "MPO" in str(e.value), e.value
    print("PASS 전제: python-pptx 는 MPO 를 거부한다")


def test_mpo_named_jpg_is_reencoded_on_ingest(patient):
    """.jpg 로 위장한 MPO 는 투입 시점에 JPEG 으로 다시 인코딩돼야 한다."""
    sid = client.post("/api/session", json={"folder": FOLDER}).json()["session_id"]
    r = client.post(f"/api/photos/{sid}",
                    files=[("files", ("IMG_0001.jpg", io.BytesIO(_mpo_bytes()), "image/jpeg"))])
    assert r.status_code == 200, r.text
    photos = r.json()["photos"]
    assert len(photos) == 1, photos

    s = M.SESSIONS[sid]
    stored = s.photos[0].path
    with Image.open(stored) as im:
        assert im.format in M.PPTX_IMAGE_FORMATS, f"{im.format} 는 PPT 에 못 넣는다"
        assert im.format == "JPEG", im.format
    print(f"PASS MPO(.jpg) → 저장 포맷 {Image.open(stored).format}")


def test_plain_jpeg_still_passes_through_untouched(patient):
    """멀쩡한 JPEG 은 종전대로 원본 바이트 그대로 둔다(화질·EXIF 보존)."""
    img = Image.fromarray(np.uint8(np.random.rand(1200, 900, 3) * 255))
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=95)
    raw = buf.getvalue()

    sid = client.post("/api/session", json={"folder": FOLDER}).json()["session_id"]
    client.post(f"/api/photos/{sid}",
                files=[("files", ("IMG_0002.jpg", io.BytesIO(raw), "image/jpeg"))])
    stored = M.SESSIONS[sid].photos[0].path
    assert stored.read_bytes() == raw, "다시 인코딩되어 원본이 아니다"
    print("PASS 정상 JPEG 은 원본 그대로")


if __name__ == "__main__":
    test_mpo_is_rejected_by_pptx()
    print("나머지는 pytest 로 실행하세요 (픽스처 필요)")
