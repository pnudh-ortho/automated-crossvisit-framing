"""
PPT 사진 삽입 (Stage 1)

사양 §7.4 삽입 로직:
 1. 슬롯 앵커 위치·크기를 읽음
 2. cover-fit(또는 §5 정합 변환)으로 사진 배치, 슬롯 중심 정렬
 3. z-order 맨 뒤 → 마스크가 초과영역을 가림
 4. 슬롯 안내 라벨/점선 윤곽 제거
 5. 재진은 템플릿 슬라이드를 임포트해 올바른 위치에 추가

python-pptx + lxml deepcopy 사용.
"""

from __future__ import annotations

import shutil
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from coords import Placement, WindowCm, cm_to_emu, cover_fit_placement
from template import find_shape, load_presentation, shape_window_cm

# 삽입된 사진 도형 이름: 슬롯을 역추적할 수 있도록 슬롯명을 접미로.
PHOTO_NAME_PREFIX = "PHOTO_"
# 사진 뒤에 까는 레터박스 배경 도형. ppt_reader는 PHOTO_ 접두만 보므로 간섭 없음.
BACKDROP_NAME_PREFIX = "BACKDROP_"
DEFAULT_LETTERBOX_COLOR = "000000"


def photo_shape_name(slot_name: str) -> str:
    return PHOTO_NAME_PREFIX + slot_name


def backdrop_shape_name(slot_name: str) -> str:
    return BACKDROP_NAME_PREFIX + slot_name


# ── z-order ──────────────────────────────────────────────────────────────────
def send_to_back(slide, shape) -> None:
    """spTree에서 사진 요소를 맨 앞 index(=가장 뒤에 그려짐)로 이동. 마스크 뒤."""
    spTree = slide.shapes._spTree
    el = shape._element
    spTree.remove(el)
    # 0: nvGrpSpPr, 1: grpSpPr, 그 다음이 첫 도형 → index 2에 삽입
    spTree.insert(2, el)


# ── 사진 삽입 ─────────────────────────────────────────────────────────────────
def set_flip_v(pic, on: bool = True) -> None:
    """도형에 상하반전을 건다 (`a:xfrm/@flipV`).

    이미지 픽셀은 그대로 두고 표시만 뒤집는다 — 원본 파일이 PPT 안에서도
    원본으로 남으므로 다음 차수 정합·재학습이 원본 기준을 그대로 쓴다.
    DrawingML 은 flip 을 도형 bbox 안에서 먼저, rot 을 그 다음에 적용하므로
    off/ext/rot 은 손대지 않아도 된다(coords.flip_editor_v 주석 참고).
    """
    xfrm = pic._element.spPr.xfrm
    if on:
        xfrm.set("flipV", "1")
    elif xfrm.get("flipV") is not None:
        del xfrm.attrib["flipV"]


def insert_photo(slide, window: WindowCm, image_path: str | Path,
                 placement: Placement, slot_name: str, flip_v: bool = False):
    """한 슬롯에 사진 1장을 placement대로 삽입하고 맨 뒤로 보낸다."""
    pic = slide.shapes.add_picture(
        str(image_path),
        left=placement.off_x, top=placement.off_y,
        width=placement.ext_cx, height=placement.ext_cy,
    )
    if placement.rot:
        pic.rotation = placement.rot / 60000.0  # 도 단위
    if flip_v:
        set_flip_v(pic)
    pic.name = photo_shape_name(slot_name)
    send_to_back(slide, pic)
    return pic


def add_slot_backdrop(slide, window: WindowCm, slot_name: str,
                      color: str = DEFAULT_LETTERBOX_COLOR):
    """
    슬롯 크기의 단색 사각형을 깔아, 사진이 창을 다 덮지 못할 때(회전·축소·이동)
    슬라이드 배경(흰색) 대신 이 색이 보이게 한다.

    반드시 사진을 넣은 **뒤에** 호출할 것 — send_to_back이 맨 뒤로 보내므로
    나중에 보낸 쪽이 더 아래에 깔린다(배경 < 사진 < 마스크).
    """
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        cm_to_emu(window.x), cm_to_emu(window.y),
        cm_to_emu(window.w), cm_to_emu(window.h),
    )
    rect.name = backdrop_shape_name(slot_name)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor.from_string(color)
    rect.line.fill.background()       # 테두리 없음
    rect.shadow.inherit = False       # 테마 그림자 상속 방지
    if rect.has_text_frame:           # 자동도형 기본 텍스트 여백 제거
        rect.text_frame.text = ""
    send_to_back(slide, rect)
    return rect


def remove_slot_anchor(slide, slot_name: str) -> None:
    """슬롯 안내 앵커(라벨 텍스트 + 점선 윤곽)를 제거."""
    sh = find_shape(slide, slot_name)
    if sh is not None:
        sh._element.getparent().remove(sh._element)


def read_slot_windows(slide, slot_names) -> dict[str, WindowCm]:
    """
    이미 만들어진 슬라이드에서 슬롯 창을 읽는다.

    템플릿의 이상적인 좌표 대신 그 PPT가 실제로 쓰던 레이아웃을 따르기 위한 것이다.
    환자마다 상단 여백·사진 간격이 조금씩 다를 수 있는데, 새 차수만 템플릿 좌표로
    넣으면 그 슬라이드만 어긋난다.

    우선순위:
      1) BACKDROP_<슬롯>  — 사진을 넣을 때 창 크기 그대로 깔리고 지워지지 않는다
      2) SLOT_<슬롯>      — 그 차수에 비어 있던 슬롯은 앵커가 남아 있다
    둘 다 없으면 그 슬롯은 결과에서 빠진다(호출측이 템플릿 기본값으로 메운다).
    """
    out: dict[str, WindowCm] = {}
    for slot in slot_names:
        sh = find_shape(slide, backdrop_shape_name(slot)) or find_shape(slide, slot)
        if sh is not None:
            out[slot] = shape_window_cm(sh)
    return out


def place_photo_in_window(slide, slot_name: str, window: WindowCm,
                          image_path: str | Path, photo_wh: tuple[int, int],
                          placement: Placement | None = None,
                          letterbox_color: str | None = DEFAULT_LETTERBOX_COLOR,
                          flip_v: bool = False):
    """
    창을 직접 받아 사진을 배치한다. 앵커가 없어도 된다.
    앵커가 남아 있으면 함께 제거한다.
    반환: (pic, 사용된 placement)
    """
    if placement is None:
        placement = cover_fit_placement(photo_wh[0], photo_wh[1], window)
    pic = insert_photo(slide, window, image_path, placement, slot_name, flip_v=flip_v)
    if letterbox_color is not None:
        # 사진 뒤에 깐다. 부수효과로, 먼저 배치된 다른 슬롯 사진이 이 슬롯으로
        # 흘러넘친 부분도 이 배경이 가려준다.
        add_slot_backdrop(slide, window, slot_name, letterbox_color)
    remove_slot_anchor(slide, slot_name)
    return pic, placement


def place_photo_in_slot(slide, slot_name: str, image_path: str | Path,
                        photo_wh: tuple[int, int], placement: Placement | None = None,
                        letterbox_color: str | None = DEFAULT_LETTERBOX_COLOR,
                        flip_v: bool = False):
    """
    슬롯 앵커를 읽어 사진을 배치(placement 없으면 cover-fit)하고 앵커 제거.
    letterbox_color가 None이 아니면 사진 뒤에 그 색의 배경 사각형을 깐다.
    반환: (pic, 사용된 placement)
    """
    anchor = find_shape(slide, slot_name)
    if anchor is None:
        raise KeyError(f"슬라이드에 슬롯 앵커 '{slot_name}'가 없습니다")
    return place_photo_in_window(slide, slot_name, shape_window_cm(anchor),
                                 image_path, photo_wh, placement, letterbox_color,
                                 flip_v=flip_v)


# ── INFO_BOX ─────────────────────────────────────────────────────────────────
def set_info_box(slide, info_box_name: str, text: str) -> None:
    sh = find_shape(slide, info_box_name)
    if sh is None:
        raise KeyError(f"슬라이드에 '{info_box_name}'가 없습니다")
    sh.text_frame.text = text


# ── 초진: 템플릿 복사로 새 PPT ────────────────────────────────────────────────
def new_ppt_from_template(template_path: str | Path, dest_path: str | Path) -> Presentation:
    """템플릿 파일을 그대로 복사해 새 환자 PPT 생성 후 연다(원본 무손상)."""
    dest_path = Path(dest_path)
    dest_path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copyfile(str(template_path), str(dest_path))
    return load_presentation(dest_path)


# ── 재진: 템플릿 슬라이드 임포트 ──────────────────────────────────────────────
def _find_layout_by_name(prs: Presentation, name: str):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    return prs.slide_layouts[0]


def _slide_id_list(prs: Presentation):
    return prs.slides._sldIdLst


def import_template_slide(dest_prs: Presentation, template_prs: Presentation,
                          insert_index: int):
    """
    템플릿의 (깨끗한) 첫 슬라이드를 dest_prs로 복제해 insert_index 위치에 삽입.
    템플릿 슬라이드는 이미지가 없으므로(도형만) 미디어 rel 복사 불필요.
    반환: 새 slide.
    """
    src_slide = template_prs.slides[0]
    dest_layout = _find_layout_by_name(dest_prs, src_slide.slide_layout.name)
    new_slide = dest_prs.slides.add_slide(dest_layout)

    # add_slide가 만든 레이아웃 placeholder들을 제거
    for ph in list(new_slide.placeholders):
        ph._element.getparent().remove(ph._element)

    # 템플릿 도형들을 deepcopy로 복제
    dst_spTree = new_slide.shapes._spTree
    for sp in src_slide.shapes:
        dst_spTree.append(deepcopy(sp._element))

    # sldIdLst에서 맨 끝(방금 추가됨) → insert_index로 이동
    sldIdLst = _slide_id_list(dest_prs)
    ids = list(sldIdLst)
    moved = ids[-1]
    sldIdLst.remove(moved)
    sldIdLst.insert(insert_index, moved)
    return new_slide


def count_slides(prs: Presentation) -> int:
    return len(prs.slides)
