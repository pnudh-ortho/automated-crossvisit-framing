"""
케이스 덱 조립 (초진)

초진 PPT는 template.pptx(케이스 프레젠테이션 양식)를 바탕으로 만든다.
좌표는 하드코딩하지 않는다 — template.py와 같은 원칙으로, 템플릿에 이미 놓여
있는 '자리 잡는 도형'(사진 플레이스홀더 / 내용 개체 틀)의 기하를 읽어 앵커로 쓴다.
템플릿에서 사진 크기를 옮기면 코드 수정 없이 따라간다.

템플릿 구조 (25.40 x 19.05 cm, 30장):
    1        환자정보 (빈 양식)
    2, 3     Cc.Dx. / Tx.  — 완전히 빈 슬라이드. 사용자가 스크린샷을 붙여넣는다.
    4, 5, 6  사진 2장 좌/우           앵커: PICTURE 2개
    7, 8, 9  사진 1장 중앙            앵커: PICTURE 1개
    10, 11   사진 1장 (위아래로 큰)   앵커: 큰 PICTURE. tracing 선·격자표는 남긴다.
    12~16    구내 개별 1장씩          앵커: 내용 개체 틀 (전면)
    17~29    타 환자의 실제 진료기록  — 버린다.
    30       빈 노트 슬라이드         십자뷰의 원본으로 재사용한다.

조립 결과 (초진):
    1~16 은 템플릿 그대로, 17 = 십자뷰(구내 5슬롯) + 노트 텍스트박스.
    십자뷰 도형(SLOT_*, MASK_*)은 기존 십자 템플릿에서 복제해 30번 슬라이드에 얹는다.
    두 템플릿은 슬라이드 크기가 같고, 십자뷰의 MASK_R1_L / MASK_R1_R 이 덮던
    좌·우 상단 코너가 정확히 노트 텍스트박스 자리다.
"""

from __future__ import annotations

from copy import deepcopy
from dataclasses import dataclass
from pathlib import Path
import re
import shutil

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from coords import EMU_PER_CM, WindowCm, emu_to_cm
from template import load_presentation


# 노트 텍스트박스 — 템플릿의 30번 슬라이드에 있는 도형 이름 그대로.
# 위치로 구분한다: 이름이 '내용 개체 틀 2'로 겹치기 때문.
NOTE_DATE = "NOTE_DATE"      # ① 좌상단 (0.40, 0.40)  "24.08.12 (초진 A)"
NOTE_SOAP = "NOTE_SOAP"      # ② ① 아래 (0.15, 1.53)  진단 목록 또는 s)/p)
NOTE_STATUS = "NOTE_STATUS"  # ③ 우상단 (17.44, 0.63) 와이어 + Tx./Rx./App. Period
NOTE_LL = "NOTE_LL"          # ④ 좌하단 (양식에 없다 — 우리가 만든다)
NOTE_NEXT = "NOTE_NEXT"      # ⑤ 우하단 (16.88, 12.5) n) 다음 계획

# 화면에 붙는 번호 순서. 사람이 "3번 박스"라고 부를 수 있게 한 곳에 못박는다.
NOTE_ORDER = [NOTE_DATE, NOTE_SOAP, NOTE_STATUS, NOTE_LL, NOTE_NEXT]
# 노트 칸 다섯의 자리(cm)와 글꼴. **양식에서 실측한 값**이다.
#
# 예전에는 양식의 빈 노트 슬라이드에서 도형을 옮겨오고, 거기 없는 칸(s)/p))은
# **진료기록이 있는 슬라이드에서 복제**했다. 그래서 양식이 남의 진료기록 13장을
# 좌표 하나 때문에 들고 다녀야 했다 (24.08~26.04 실제 경과). 그 파일이 배포되면
# 링크를 가진 사람이 그 기록을 받는다.
#
# 필요한 것은 좌표와 글꼴뿐이므로 여기 적는다. 양식은 앞 16장(진료 서식)만 남는다.
# 값을 바꾸면 결과물과 화면 오버레이가 함께 움직인다 — 둘 다 여기를 읽는다.
NOTE_WINDOWS = {
    # 날짜/차수는 (0.40, 0.40) 고정 — 재진은 사람이 만든 PPT의 박스를 그대로
    # 쓰는데 사람마다 자리가 조금씩 달랐다. 초진(여기)과 재진(_write_visit_label
    # 이 기존 박스를 이 좌표로 옮긴다)이 같은 자리를 쓰도록 통일한다 (2026-08-12).
    NOTE_DATE:   {"x":  0.40, "y":  0.40, "w":  7.22, "h": 0.90,
                  "size_pt": 15.0, "color": "FFFFFF", "bold": False},
    NOTE_SOAP:   {"x":  0.15, "y":  1.53, "w":  8.12, "h": 3.11,
                  "size_pt": 12.0, "color": "FFFFFF", "bold": False},
    # 양식 원본은 17.44,0.63 에 8.41 폭이라 우측이 0.45cm 삐져나갔다.
    # 우상단 영역(MASK_R1_R 16.90,0.03 8.50x6.30) 안으로 넣는다.
    NOTE_STATUS: {"x": 17.05, "y":  0.63, "w":  8.20, "h": 5.28,
                  "size_pt": 15.0, "color": "FFFF00", "bold": False},
    # ── 아래 두 칸은 십자뷰의 빈 영역을 **꽉 채운다** ──────────────────
    # 십자뷰는 3×3 인데 아래 줄의 좌·우가 비어 있다 (가운데만 SLOT_LOWER):
    #     MASK_R3_L   0.00, 12.72  8.50×6.30      MASK_R3_R  16.90, 12.72  8.50×6.30
    # 양식의 NOTE_NEXT 는 16.88,12.58 에 12.93×1.11 이라 **가로가 슬라이드 밖으로
    # 4.4cm 나가고** 세로는 1.11cm 뿐이었다. 좌하단은 애초에 양식에 도형이 없었다.
    # 둘을 대칭으로 맞춰 영역을 채운다 — 마스크 경계에 0.15cm 여백만 둔다.
    NOTE_LL:     {"x":  0.15, "y": 12.87, "w":  8.20, "h": 6.00,
                  "size_pt": 14.0, "color": "FFFFFF", "bold": False},
    NOTE_NEXT:   {"x": 17.05, "y": 12.87, "w":  8.20, "h": 6.00,
                  "size_pt": 14.0, "color": "FFFFFF", "bold": False},
}
NOTE_LL_WINDOW = NOTE_WINDOWS[NOTE_LL]      # 옛 이름 유지


def apply_note_sizes(sizes: dict) -> None:
    """설정의 노트 글자 크기를 상수에 반영 — 화면 오버레이와 생성이 함께 따라온다."""
    for k, v in (sizes or {}).items():
        if k in NOTE_WINDOWS and v:
            NOTE_WINDOWS[k]["size_pt"] = float(v)

_PICTURE = 13
_PLACEHOLDER = 14


@dataclass(frozen=True)
class Anchor:
    """사진이 들어갈 자리. window는 cm, shape_name은 템플릿의 원래 도형 이름."""
    slide: int          # 1-based 슬라이드 번호
    pos: str            # "L" | "R" | "C" | "BIG" | "FULL"
    window: WindowCm
    shape_name: str


def _win(shape) -> WindowCm:
    return WindowCm(x=emu_to_cm(shape.left), y=emu_to_cm(shape.top),
                    w=emu_to_cm(shape.width), h=emu_to_cm(shape.height))


def read_photo_anchors(prs: Presentation, slide_no: int) -> list[Anchor]:
    """
    한 슬라이드의 사진 앵커를 읽는다.

    사진 플레이스홀더(PICTURE)가 있으면 그것을, 없으면 내용 개체 틀을 앵커로 쓴다.
    좌/우 구분은 x좌표 순서로 정한다 — 도형 이름에 의미가 없기 때문이다.

    사진이 여러 장이어도 좌우로 갈라져 있을 때만 두 자리로 본다. 슬라이드 10·11처럼
    같은 자리에 겹쳐 있는 경우는 한 자리다 — 겹친 장수는 원본 케이스의 흔적일 뿐이라
    가장 큰 것 하나만 남긴다.
    """
    slide = prs.slides[slide_no - 1]
    pics = [sh for sh in slide.shapes if sh.shape_type == _PICTURE]

    if pics:
        pics.sort(key=lambda s: (s.left, s.top))
        first, last = pics[0], pics[-1]
        split = len(pics) >= 2 and (first.left + first.width) <= last.left + Emu(1)
        if split:
            return [Anchor(slide_no, "L", _win(first), first.name),
                    Anchor(slide_no, "R", _win(last), last.name)]
        p = max(pics, key=lambda s: (s.width or 0) * (s.height or 0))
        # 슬라이드 높이를 거의 다 채우면 '위아래로 큰' 자리로 본다.
        big = emu_to_cm(p.height) >= emu_to_cm(prs.slide_height) * 0.95
        return [Anchor(slide_no, "BIG" if big else "C", _win(p), p.name)]

    phs = [sh for sh in slide.shapes if sh.shape_type == _PLACEHOLDER]
    if phs:
        ph = max(phs, key=lambda s: (s.width or 0) * (s.height or 0))
        return [Anchor(slide_no, "FULL", _win(ph), ph.name)]
    return []


def read_deck_anchors(prs: Presentation, face_slides: list[int],
                      big_slides: list[int], intraoral_slides: list[int]) -> dict:
    """
    덱 전체의 앵커 표. {(slide, pos): Anchor}

    구내 슬라이드 중에는 자리 잡는 도형이 아예 없는 것도 있다(격자표만 있는 슬라이드).
    그때는 바로 앞 구내 슬라이드의 창을 물려받는다 — 같은 크기로 놓이는 자리들이라
    슬라이드마다 빈 도형을 넣어 두라고 템플릿에 요구하지 않는다.
    """
    out: dict[tuple[int, str], Anchor] = {}
    for n in list(face_slides) + list(big_slides):
        for a in read_photo_anchors(prs, n):
            out[(a.slide, a.pos)] = a

    inherited: WindowCm | None = None
    for n in intraoral_slides:
        found = read_photo_anchors(prs, n)
        if found:
            inherited = found[0].window
            out[(found[0].slide, found[0].pos)] = found[0]
        elif inherited is not None:
            out[(n, "FULL")] = Anchor(n, "FULL", inherited, "(상속)")
    return out


# ── 슬라이드 삭제 ─────────────────────────────────────────────────────────────
def drop_slides(prs: Presentation, keep: set[int]) -> None:
    """
    keep에 없는 슬라이드를 모두 지운다 (1-based).
    python-pptx에 삭제 API가 없어 sldIdLst와 관계(rel)를 직접 정리한다.
    """
    sldIdLst = prs.slides._sldIdLst
    for i, sldId in reversed(list(enumerate(list(sldIdLst), start=1))):
        if i in keep:
            continue
        prs.part.drop_rel(sldId.rId)
        sldIdLst.remove(sldId)


def move_slide(prs: Presentation, frm: int, to: int) -> None:
    """슬라이드를 frm(1-based)에서 to(1-based)로 옮긴다."""
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    node = ids[frm - 1]
    sldIdLst.remove(node)
    sldIdLst.insert(to - 1, node)


# ── 노트 텍스트박스 ───────────────────────────────────────────────────────────
def note_key(x_cm: float, y_cm: float, slide_w_cm: float) -> str:
    """텍스트박스의 위치로 어느 노트 칸인지 정한다 (십자뷰의 네 모서리).

    템플릿에는 '내용 개체 틀 2'가 두 개라 이름으로는 구분할 수 없다. 위치가
    유일한 단서이므로 판정을 여기 한 곳에만 둔다 — 이름 붙이기·보충·화면
    오버레이가 전부 같은 규칙을 써야 한 칸이 두 뜻을 갖는 일이 없다.
    """
    half = slide_w_cm / 2
    if x_cm < half and y_cm < 1.2:
        return NOTE_DATE            # ① 좌상단 — 날짜/차수
    if x_cm >= half and y_cm < 6.0:
        return NOTE_STATUS          # ③ 우상단 — 와이어 · Tx./Rx./App. Period
    if x_cm < half:
        # 십자뷰 아래칸(12.72cm~)까지 내려온 것만 좌하단이다. 그 위는 ① 아래 칸.
        return NOTE_LL if y_cm >= 12.0 else NOTE_SOAP
    return NOTE_NEXT                # ⑤ 우하단 — n) 다음


def name_note_boxes(slide, slide_w_cm: float) -> dict[str, str]:
    """
    노트 슬라이드의 텍스트박스에 의미 있는 이름을 붙인다.
    반환: {의미이름: 원래이름}
    """
    renamed: dict[str, str] = {}
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        key = note_key(emu_to_cm(sh.left), emu_to_cm(sh.top), slide_w_cm)
        if key in renamed:      # 같은 자리에 둘이면 첫 번째만 쓴다
            continue
        renamed[key] = sh.name
        sh.name = key
    return renamed


def _run_font(shape) -> dict:
    """박스 첫 글자의 글꼴(크기 pt · 색 · 굵기). 화면 오버레이를 양식과 같게 보이려는 것.

    테마색은 실제 색을 파일에서 알 수 없다(테마 파트를 따라가야 한다). 노트 박스에
    쓰이는 것은 배경 대비색 하나뿐이라 흰색으로 둔다 — 틀리면 화면만 달라 보이고
    슬라이드 결과에는 영향이 없다.
    """
    for para in shape.text_frame.paragraphs:
        for r in para.runs:
            f = r.font
            col = None
            try:
                col = str(f.color.rgb)
            except Exception:
                if f.color is not None and f.color.type is not None:
                    col = "FFFFFF"          # 테마색 → 배경 대비색으로 간주
            return {"size_pt": f.size.pt if f.size else None,
                    "color": col, "bold": bool(f.bold)}
    return {"size_pt": None, "color": None, "bold": False}


def note_box_windows(prs: Presentation = None, note_slide_no: int = 0,
                     donor_slides=None) -> dict[str, dict]:
    """노트 칸 다섯의 위치·크기·글꼴. 화면 오버레이가 결과물과 같아 보이게 하는 값.

    `NOTE_WINDOWS` 를 그대로 돌려준다 — 결과물도 같은 상수로 만들어지므로 화면과
    슬라이드가 어긋날 수 없다. 예전에는 양식 슬라이드를 뒤져 좌표를 재고, 화면과
    결과물이 각각 그 일을 했다. 인자는 옛 호출부 호환용이다.
    """
    return {k: {**v, "font": {"size_pt": v["size_pt"], "color": v["color"],
                              "bold": v["bold"]}}
            for k, v in NOTE_WINDOWS.items()}


def ensure_note_boxes(prs: Presentation, cross_slide, donor_slides=None) -> list[str]:
    """노트 박스 5종(`NOTE_ORDER`)이 모두 있도록 **좌표로 만든다**.

    예전에는 양식 슬라이드에서 도형을 복제했다. 그러려면 양식이 그 칸을 가진
    진료기록 슬라이드를 들고 있어야 했고, 그 기록이 배포 파일에 실려 나갔다.
    좌표와 글꼴은 `NOTE_WINDOWS` 에 실측해 두었으므로 복제할 이유가 없다.

    `donor_slides` 는 옛 호출부 호환용이며 쓰이지 않는다.
    """
    have = {sh.name for sh in cross_slide.shapes}
    added: list[str] = []
    for want in NOTE_ORDER:
        if want in have:
            continue
        w = NOTE_WINDOWS[want]
        if add_note_box(cross_slide, want, w, size_pt=w["size_pt"],
                        color_rgb=w["color"], bold=w["bold"]):
            added.append(want)
    return added


def pin_shape(slide, name: str, x_cm: float, y_cm: float) -> bool:
    """이름이 맞는 도형을 지정 좌표(cm)로 옮긴다.

    재진의 날짜/차수 박스는 사람이 만든 PPT 에서 오므로 자리가 제각각이다 —
    쓸 때마다 초진과 같은 자리로 고정한다 (`NOTE_WINDOWS[NOTE_DATE]` 주석 참조).
    """
    for sh in slide.shapes:
        if sh.name == name:
            sh.left = Emu(int(round(x_cm * EMU_PER_CM)))
            sh.top = Emu(int(round(y_cm * EMU_PER_CM)))
            return True
    return False


def replace_with_copied_box(slide, name: str, sp_xml) -> bool:
    """원본 PPT 의 텍스트 상자를 **통째로 복사**해 오고 이름만 규약명으로 바꾼다.

    위치·크기·글꼴·문단 규칙(lstStyle)·배경까지 전부 따라오므로 속성 개별 상속이
    필요 없다. 내용은 이후 `set_note_text` 가 원래 글꼴을 되살려 갈아끼운다.
    같은 이름의 기존 박스(템플릿 것)는 제거한다.
    """
    if not sp_xml:
        return False
    spTree = slide.shapes._spTree
    for sh in list(slide.shapes):
        if sh.name == name:
            sh._element.getparent().remove(sh._element)
    new = etree.fromstring(sp_xml)
    nv = new.find(qn("p:nvSpPr"))
    cnv = nv.find(qn("p:cNvPr")) if nv is not None else None
    if cnv is None:
        return False
    ids = [int(e.get("id")) for e in spTree.iter(qn("p:cNvPr"))
           if (e.get("id") or "").isdigit()]
    cnv.set("id", str(max(ids, default=1) + 1))
    cnv.set("name", name)
    spTree.append(new)                  # 맨 뒤 = 맨 앞(z-order 위) — 글자가 보인다
    return True


def force_run_color(slide, name: str, rgb: str = "FFFFFF") -> None:
    """박스의 글자 색을 명시 RGB 로 고정한다.

    템플릿 라벨(INFO_BOX)은 색이 테마(bg1)로 지정돼 있어, 수제 PPT 에 슬라이드를
    끼워 넣으면 그 파일의 테마에 따라 검정으로 풀릴 수 있다. 명시 흰색이면
    어떤 테마에서도 같은 모습이다.
    """
    for sh in slide.shapes:
        if sh.name != name or not getattr(sh, "has_text_frame", False):
            continue
        for para in sh.text_frame.paragraphs:
            for r in para.runs:
                r.font.color.rgb = RGBColor.from_string(rgb)
        return


def fix_empty_para_sizes(slide, name: str) -> None:
    """빈 문단의 endParaRPr 크기를 본문 크기로 맞춘다.

    `set_note_text` 가 문단을 다시 만들면 빈 줄에는 run 이 없어, 복사해 온
    lstStyle 의 기본 크기(옛 마스터 32pt 등)로 떨어질 수 있다.
    """
    for sh in slide.shapes:
        if sh.name != name or not getattr(sh, "has_text_frame", False):
            continue
        base = None
        for para in sh.text_frame.paragraphs:
            for r in para.runs:
                if r.font.size:
                    base = r.font.size.pt
                    break
            if base:
                break
        if not base:
            return
        for para in sh.text_frame.paragraphs:
            if para.runs:
                continue
            p_el = para._p
            end = p_el.find(qn("a:endParaRPr"))
            if end is None:
                end = p_el.makeelement(qn("a:endParaRPr"), {})
                p_el.append(end)
            end.set("sz", str(int(base * 100)))
        return


def style_note_box(slide, name: str, style: dict) -> bool:
    """라벨 박스에 수제 PPT 의 자리·크기·폰트를 입힌다.

    `ppt_reader.last_label_style` 의 출력을 받는다 — 수제 PPT 에 이어붙일 때
    원본 라벨 모습을 따라가기 위한 것 (앱 PPT 는 (0.4, 0.4) 고정 규칙 유지).
    """
    for sh in slide.shapes:
        if sh.name != name:
            continue
        sh.left = Emu(int(round(style["x"] * EMU_PER_CM)))
        sh.top = Emu(int(round(style["y"] * EMU_PER_CM)))
        sh.width = Emu(int(round(style["w"] * EMU_PER_CM)))
        sh.height = Emu(int(round(style["h"] * EMU_PER_CM)))
        if sh.has_text_frame:
            tf = sh.text_frame
            try:
                # 원본 박스에는 자동 맞춤이 없다 — spAutoFit 이 남아 있으면
                # 상자 높이가 텍스트에 끌려가 상속한 크기가 무의미해진다.
                tf.auto_size = None
            except Exception:                                     # noqa: BLE001
                pass
            xml = style.get("lst_style_xml")
            if xml:
                # 원본 상자의 lstStyle 에서 **간격 규칙만** 이식한다.
                # 통째로 옮기면 옛 마스터 기본값(defRPr sz=3200=32pt, 불릿)까지
                # 따라와서, 명시 크기가 없는 빈 줄이 32pt 로 커진다 — 원본은
                # 모든 문단에 명시값이 있어 그 기본값이 드러나지 않았을 뿐이다.
                new = etree.fromstring(xml)
                keep = {qn("a:spcBef"), qn("a:spcAft"), qn("a:lnSpc")}
                for lvl in list(new):
                    for child in list(lvl):
                        if child.tag not in keep:
                            lvl.remove(child)
                tb = tf._txBody
                old = tb.find(qn("a:lstStyle"))
                if old is not None:
                    tb.remove(old)
                bp = tb.find(qn("a:bodyPr"))
                tb.insert(list(tb).index(bp) + 1, new)
            base = style.get("size_pt")
            for para in tf.paragraphs:
                # 줄간격은 **원본 값 그대로** — 미지정(None)이면 None 으로 둔다.
                # 1.0 을 강제로 박으면 개체 틀 상속 기본과 미묘하게 달라진다.
                para.line_spacing = style.get("line_spacing")
                if base and not para.runs:
                    # 빈 문단은 run 이 없어 endParaRPr 가 줄 높이를 정한다 —
                    # 없으면 기본 18pt 로 떨어져 원본(15pt)과 달라진다.
                    p_el = para._p
                    end = p_el.find(qn("a:endParaRPr"))
                    if end is None:
                        end = p_el.makeelement(qn("a:endParaRPr"), {})
                        p_el.append(end)
                    end.set("sz", str(int(base * 100)))
                for r in para.runs:
                    cur = r.font.size.pt if r.font.size else None
                    # 줄 끝 괄호(9pt) 같은 '작은 글씨' run 만 보존한다 — 기준의
                    # 80% 미만일 때다. (템플릿 14pt vs 원본 15pt 처럼 근소하게
                    # 작은 것은 '작은 글씨'가 아니라 그냥 다른 기본값이다.)
                    if base and (cur is None or cur >= base * 0.8):
                        r.font.size = Pt(base)
                    if style.get("font"):
                        r.font.name = style["font"]
                    if style.get("bold") is not None:
                        r.font.bold = style["bold"]
                    if style.get("color"):
                        r.font.color.rgb = RGBColor.from_string(style["color"])
        return True
    return False


def add_note_box(slide, name: str, win: dict, size_pt: float = 12.0,
                 color_rgb: str = "FFFFFF", bold: bool = False) -> bool:
    """빈 노트 텍스트박스를 새로 만든다 (양식에 없는 칸용).

    글꼴은 같은 쪽 칸(s)/p))에 맞춰 둔다 — 양식을 흉내 내는 것이지 새 서식을
    들이는 게 아니다. 채우기·테두리는 두지 않아 슬라이드 배경이 그대로 비친다.
    """
    try:
        box = slide.shapes.add_textbox(
            Emu(int(win["x"] * EMU_PER_CM)), Emu(int(win["y"] * EMU_PER_CM)),
            Emu(int(win["w"] * EMU_PER_CM)), Emu(int(win["h"] * EMU_PER_CM)))
    except Exception:
        return False
    box.name = name
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _apply_para(p, None)          # 글머리 점 없음 — 양식의 노트 칸과 같게
    r = p.add_run()
    r.text = ""
    r.font.size = Pt(size_pt)
    r.font.bold = bold
    r.font.color.rgb = RGBColor.from_string(color_rgb)
    # 빈 상자는 PowerPoint 가 빈 run 을 무시하고 커서 글꼴을 기본 18pt 로
    # 보여준다 — endParaRPr 에도 크기·색을 새겨 타이핑이 설정 크기를 따르게 한다.
    end = p._p.find(qn("a:endParaRPr"))
    if end is None:
        end = p._p.makeelement(qn("a:endParaRPr"), {})
        p._p.append(end)
    end.set("sz", str(int(size_pt * 100)))
    return True


def add_note_boxes_from_layout(slide, layout: dict[str, dict],
                               skip: set[str] = frozenset()) -> list[str]:
    """자리표(`note_box_windows` 가 양식에서 읽어 둔 것)대로 빠진 노트 칸을 만든다.

    재진 슬라이드는 십자뷰 양식에서 임포트되는데 거기엔 노트 칸이 아예 없다 —
    그대로 두면 이번 차수에 적은 노트가 갈 곳이 없어 조용히 사라진다.
    화면 오버레이가 쓰는 자리표를 그대로 쓰므로 화면과 결과물이 어긋나지 않는다.
    반환: 새로 만든 칸 이름 목록.
    """
    have = {sh.name for sh in slide.shapes}
    added: list[str] = []
    for name in NOTE_ORDER:
        win = layout.get(name)
        if win is None or name in have or name in skip:
            continue
        f = win.get("font") or {}
        if add_note_box(slide, name, win,
                        size_pt=f.get("size_pt") or 12.0,
                        color_rgb=f.get("color") or "FFFFFF",
                        bold=bool(f.get("bold"))):
            added.append(name)
    return added


def bring_to_front(slide, names) -> list[str]:
    """도형을 spTree 맨 끝으로 옮긴다 — 맨 앞에 그려진다.

    십자뷰를 얹으면(`graft_cross_view`) 검은 마스크가 노트 칸 **위에** 쌓인다.
    MASK_R1_L/R·MASK_R3_L/R 이 덮는 네 모서리가 정확히 노트 칸 자리라, 글자는
    파일에 남아 있어도 화면과 인쇄물에서는 보이지 않는다. 얹은 뒤 다시 올린다.
    """
    spTree = slide.shapes._spTree
    moved: list[str] = []
    for name in names:
        for sh in slide.shapes:
            if sh.name != name:
                continue
            el = sh._element
            spTree.remove(el)
            spTree.append(el)
            moved.append(name)
            break
    return moved


_TAIL_PAREN = re.compile(r"\s*\([^()]*\)\s*$")


def _copy_font(src: dict, dst_font, size_pt: float | None = None) -> None:
    """떠 둔 서식을 새 run 에 입힌다. size_pt 를 주면 크기만 그걸로 바꾼다."""
    want = size_pt if size_pt is not None else src["size_pt"]
    if want:
        dst_font.size = Pt(want)
    if src["bold"] is not None:
        dst_font.bold = src["bold"]
    if src["rgb"]:
        dst_font.color.rgb = RGBColor.from_string(src["rgb"])
    elif src["theme"] is not None:
        dst_font.color.theme_color = src["theme"]


def _base_para(tf):
    """박스 첫 문단의 문단서식(`a:pPr`). 글을 갈아끼운 뒤 되살리려고 떠 둔다.

    `tf.clear()` 는 문단을 통째로 새로 만들어 pPr 이 사라진다. 그러면 문단이
    개체 틀의 목록서식(lvl1pPr)을 물려받아 **없던 글머리 점(•)** 이 Tx./Rx./App.
    앞에 생긴다 — 양식의 문단들은 저마다 `buNone` 을 들고 있어서 점이 없었다.
    들여쓰기·정렬도 같이 들어 있으므로 통째로 떠서 그대로 다시 입힌다.
    """
    for para in tf.paragraphs:
        pPr = para._p.find(qn("a:pPr"))
        if pPr is not None:
            return deepcopy(pPr)
    return None


def _apply_para(para, pPr) -> None:
    """떠 둔 문단서식을 새 문단에 입힌다. 뜬 게 없으면 글머리 점만 끈다."""
    p = para._p
    old = p.find(qn("a:pPr"))
    if old is not None:
        p.remove(old)
    if pPr is not None:
        p.insert(0, deepcopy(pPr))
        return
    new = p.makeelement(qn("a:pPr"), {})
    new.append(new.makeelement(qn("a:buNone"), {}))
    p.insert(0, new)


def _base_font(tf) -> dict:
    """박스의 첫 글자 서식. 글을 갈아끼운 뒤 되살리려고 미리 떠 둔다."""
    for para in tf.paragraphs:
        for r in para.runs:
            f, rgb, theme = r.font, None, None
            # 색을 아예 안 지정한 run 은 rgb 도 theme_color 도 예외를 던진다
            # (_NoneColor). 그 경우 색은 건드리지 않고 문단 서식을 물려받게 둔다.
            try:
                rgb = str(f.color.rgb)
            except Exception:
                try:
                    theme = f.color.theme_color
                except Exception:
                    theme = None
            return {"size_pt": f.size.pt if f.size else None,
                    "bold": f.bold, "rgb": rgb, "theme": theme}
    return {"size_pt": None, "bold": None, "rgb": None, "theme": None}


def set_note_text(slide, box_name: str, text: str, small_pt: float | None = None) -> bool:
    """노트 박스의 텍스트를 통째로 갈아끼운다. 원래 글꼴을 되살려 쓴다.

    `small_pt` 를 주면 **줄 끝 괄호**(예: "Rx. Period: 23 month (24.08.12)" 의
    날짜)만 그 크기로 줄여 쓴다 — 양식이 그렇게 돼 있다(본문 15pt, 날짜 9pt).
    괄호만 있는 줄은 본문이 없으니 건드리지 않는다.
    """
    for sh in slide.shapes:
        if sh.name != box_name or not sh.has_text_frame:
            continue
        tf = sh.text_frame
        base = _base_font(tf)
        base_p = _base_para(tf)
        tf.clear()
        for i, ln in enumerate((text or "").split("\n")):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            _apply_para(para, base_p)
            head, tail = ln, ""
            if small_pt:
                m = _TAIL_PAREN.search(ln)
                if m and ln[:m.start()].strip():
                    head, tail = ln[:m.start()], ln[m.start():]
            r = para.add_run()
            r.text = head
            _copy_font(base, r.font)
            if tail:
                r2 = para.add_run()
                r2.text = tail
                _copy_font(base, r2.font, size_pt=small_pt)
        return True
    return False


# ── 환자정보 슬라이드 ─────────────────────────────────────────────────────────
def fill_prefixed_lines(slide, filled: dict[str, str]) -> dict[str, str]:
    """머리말로 문단을 찾아 **그 뒤만** 갈아끼운다.

    양식 첫 장의 "Hospital No. / Case No. / Pt. name : / C/C :" 처럼, 사람이
    적을 자리만 비워 둔 줄을 채우기 위한 것이다. 문단을 통째로 새로 쓰면 양식의
    글꼴·줄간격·정렬이 날아가므로, 첫 run 의 글을 갈아끼우고 나머지 run 만
    걷어낸다 — 서식은 첫 run 이 들고 있다.

    filled: {머리말: 머리말 뒤에 올 글}. 머리말은 양식에 적힌 글자 그대로.
    반환: 실제로 채운 {머리말: 최종 문단 글}.
    """
    todo = dict(filled)
    done: dict[str, str] = {}
    for sh in slide.shapes:
        if not sh.has_text_frame or not todo:
            continue
        for para in sh.text_frame.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            key = next((k for k in todo if text.startswith(k)), None)
            if key is None:
                continue
            runs = para.runs
            if not runs:            # 서식을 물려받을 run 이 없으면 건드리지 않는다
                continue
            runs[0].text = key + todo[key]
            for extra in runs[1:]:
                extra._r.getparent().remove(extra._r)
            done[key] = key + todo.pop(key)
            if not todo:
                break
    return done


def get_note_text(slide, box_name: str) -> str | None:
    for sh in slide.shapes:
        if sh.name == box_name and sh.has_text_frame:
            return sh.text_frame.text
    return None


# ── 십자뷰 슬라이드 조립 ──────────────────────────────────────────────────────
def graft_cross_view(note_slide, cross_prs: Presentation,
                     drop_shape_names: set[str] = frozenset()) -> None:
    """
    십자 템플릿(구내 5슬롯)의 도형들을 노트 슬라이드 위에 복제해 얹는다.

    십자 템플릿 슬라이드에는 이미지가 없고 도형뿐이라 미디어 관계 복사가 필요 없다
    (ppt_writer.import_template_slide와 같은 전제).
    노트 텍스트박스가 이미 날짜를 들고 있으므로 십자 템플릿의 INFO_BOX 등
    겹치는 도형은 drop_shape_names로 걸러낸다.
    """
    src = cross_prs.slides[0]
    spTree = note_slide.shapes._spTree
    for sp in src.shapes:
        if sp.name in drop_shape_names:
            continue
        spTree.append(deepcopy(sp._element))


def build_first_visit_deck(case_template: str | Path, cross_template: str | Path,
                           dest_path: str | Path, keep_slides: int,
                           note_slide_no: int = 0,
                           cross_drop_shapes: set[str] = frozenset()) -> Presentation:
    """초진 덱을 만든다 — **케이스 양식 + 십자뷰 양식**.

    1) 케이스 양식을 복사 (원본 무손상)
    2) `keep_slides` 뒤가 있으면 버린다 (옛 30장짜리 양식 호환)
    3) 십자뷰가 얹힐 **빈 슬라이드를 새로 추가**한다
    4) 십자뷰 도형을 얹고 노트 칸 다섯을 좌표로 만든다

    예전에는 양식 안의 '빈 노트 슬라이드'를 옮겨와 바탕으로 쓰고, s)/p) 칸은
    진료기록 슬라이드에서 복제했다. 배포 파일에 남의 진료기록이 실려 나가는
    원인이었다. 이제 양식은 앞 16장만 있으면 된다.

    `note_slide_no` 는 옛 호출부 호환용이며 쓰이지 않는다.

    반환: 열린 Presentation. 십자뷰는 마지막 슬라이드다.
    """
    dest_path = Path(dest_path)
    dest_path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copyfile(str(case_template), str(dest_path))
    prs = load_presentation(dest_path)

    total = len(prs.slides._sldIdLst)
    if not (1 <= keep_slides <= total):
        raise ValueError(f"keep_slides 가 범위를 벗어났습니다: {keep_slides}/{total}")
    if total > keep_slides:
        drop_slides(prs, keep=set(range(1, keep_slides + 1)))

    # 십자뷰 바탕. 양식 마지막 장과 같은 레이아웃을 써서 배경·머리말을 잇는다.
    cross = prs.slides.add_slide(prs.slides[keep_slides - 1].slide_layout)
    for ph in list(cross.placeholders):      # 레이아웃이 딸려 보낸 빈 자리표시자 제거
        ph._element.getparent().remove(ph._element)

    graft_cross_view(cross, load_presentation(cross_template),
                     drop_shape_names=cross_drop_shapes)
    ensure_note_boxes(prs, cross)
    # 십자뷰의 검은 마스크가 노트 칸을 덮는다 — 다시 맨 앞으로.
    bring_to_front(cross, NOTE_ORDER)
    return prs


def cross_slide_index(prs: Presentation) -> int:
    """십자뷰(마지막) 슬라이드의 0-based 인덱스."""
    return len(prs.slides._sldIdLst) - 1
