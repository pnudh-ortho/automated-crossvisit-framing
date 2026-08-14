"""
FastAPI 로컬 서버 (Stage 5)

완전 오프라인. 브라우저 단일 페이지 UI(localhost)와 통신.
세션 상태는 서버 메모리에 보관하고 업로드 사진은 세션 임시폴더에 저장한다.
확정(commit) 전에는 환자 폴더/PPT에 아무것도 쓰지 않는다.
"""

from __future__ import annotations

import base64
import contextlib
import copy
import io
from collections import Counter
import json
import os
import re
import shutil
import string
import subprocess
import sys
import threading
import time
import uuid
from datetime import datetime
from pathlib import Path

import cv2
import numpy as np
from fastapi import Body, FastAPI, File, HTTPException, UploadFile
from fastapi.responses import FileResponse, JSONResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from pptx.oxml.ns import qn
from pydantic import BaseModel

import case_deck as CD
import config as C
import crop as Cr
import framing as Fr
import naming as N
import ppt_reader as Rd
import ppt_writer as W
import registration_teeth as Reg   # 차수 간 정합 = 치아 중심점
import uninstaller as Un
import updater as Up
import storage as S
import template as T
from classify import load_classifier
from coords import (EMU_PER_CM, EditorState, Placement, WindowCm, apply_cover_clamp,
                    cover_base_ext_cm, cover_fit_placement, editor_to_placement,
                    emu_to_cm, flip_editor_v, placement_from_photo_affine)

BACKEND_DIR = Path(__file__).resolve().parent
FRONTEND_DIR = BACKEND_DIR.parent / "frontend"

cfg = C.load_config()
classifier = load_classifier(cfg)
# 자동 프레이밍 모델. 없으면 None 이고 초진은 종전대로 cover-fit 으로 놓인다.
framer = Fr.load_framer(cfg)
TEMPLATE_PRS = T.load_presentation(cfg.resolve(cfg.paths.template_pptx))
SLOT_WINDOWS: dict[str, WindowCm] = T.slot_windows(TEMPLATE_PRS.slides[0], cfg.ppt.slot_names)
Rd.set_slot_windows(SLOT_WINDOWS)

# 정합 모델을 기동 때 미리 올린다 — 첫 분류 요청만 수 초 느려지는 걸 막는다.
# 실패해도 서버는 뜬다: 정합을 건너뛰고 프레이밍 모델이 받는다.
_REG_READY = Reg.warmup()


def _imread(path) -> np.ndarray | None:
    """한글·공백·특수문자 경로에 안전한 imread.

    Windows 의 `cv2.imread` 는 경로를 ANSI 로 열어서 비ASCII 경로(한글 환자
    폴더명 등)에서 조용히 None 을 돌려준다. 바이트를 파이썬이 읽고 cv2 는
    디코드만 하게 하면 경로 인코딩과 무관해진다.
    """
    try:
        return cv2.imdecode(np.fromfile(str(path), np.uint8), cv2.IMREAD_COLOR)
    except Exception:                                             # noqa: BLE001
        return None


def _windows_json(windows: dict[str, WindowCm]) -> dict:
    return {k: {"x": v.x, "y": v.y, "w": v.w, "h": v.h} for k, v in windows.items()}


def _layout_from_ppt(prs) -> dict[str, WindowCm]:
    """
    기존 PPT가 실제로 쓰던 슬롯 레이아웃. 없는 슬롯은 템플릿 좌표로 메운다.

    가장 마지막 차수 슬라이드를 기준으로 삼는다 — 레이아웃이 도중에 바뀐 PPT라면
    최신 것이 이어붙일 슬라이드와 가장 가깝다.
    """
    found: dict[str, WindowCm] = {}
    ctr = (emu_to_cm(prs.slide_width) / 2, emu_to_cm(prs.slide_height) / 2)
    for slide in prs.slides:
        got = W.read_slot_windows(slide, cfg.ppt.slot_names)
        # 사진이 실제로 들어간 슬라이드만 기준으로 삼는다(빈 앵커뿐인 템플릿 슬라이드 제외)
        if got and any(sh.name.startswith(W.PHOTO_NAME_PREFIX) for sh in slide.shapes):
            found = got
        else:
            est = _estimate_slot_windows(slide, ctr)   # 수제 십자뷰 — 사진 bbox 상속
            if est:
                found = est
    return {**SLOT_WINDOWS, **found}


def _estimate_slot_windows(slide, slide_ctr: tuple[float, float]) -> dict[str, WindowCm]:
    """수제 십자뷰(가로 8cm 이상 사진 5장)에서 슬롯 창 = 각 사진의 bbox.

    무게중심에 가장 가까운 사진이 정면, 나머지는 정면 기준 상/하/좌/우 —
    ppt_reader 의 기준영상 추정과 같은 규칙이라 두 결과가 어긋나지 않는다.
    """
    pics = [sh for sh in slide.shapes
            if getattr(sh, "shape_type", None) == 13
            and emu_to_cm(sh.width) >= 8.0]
    if len(pics) < 5:
        return {}
    ctr = [(emu_to_cm(p.left) + emu_to_cm(p.width) / 2,
            emu_to_cm(p.top) + emu_to_cm(p.height) / 2) for p in pics]
    gx, gy = slide_ctr        # 정면 = 슬라이드 중앙에 가장 가까운 사진
    fi = min(range(len(pics)),
             key=lambda k: (ctr[k][0] - gx) ** 2 + (ctr[k][1] - gy) ** 2)
    fx, fy = ctr[fi]
    picked = {"SLOT_FRONT": (fi, 0.0)}
    for k in range(len(pics)):
        if k == fi:
            continue
        dx, dy = ctr[k][0] - fx, ctr[k][1] - fy
        slot = (("SLOT_UPPER" if dy < 0 else "SLOT_LOWER") if abs(dy) >= abs(dx)
                else ("SLOT_LEFT" if dx < 0 else "SLOT_RIGHT"))
        d2 = dx * dx + dy * dy
        if slot not in picked or d2 < picked[slot][1]:
            picked[slot] = (k, d2)
    return {slot: WindowCm(x=emu_to_cm(pics[k].left), y=emu_to_cm(pics[k].top),
                           w=emu_to_cm(pics[k].width), h=emu_to_cm(pics[k].height))
            for slot, (k, _d) in picked.items()}


def _case_deck_ready() -> bool:
    """케이스 양식을 쓸 수 있는가. 설정이 꺼져 있거나 파일이 없으면 종전 방식."""
    cd = cfg.case_deck
    if not (cd and cd.enabled and cfg.paths.case_template_pptx):
        return False
    return Path(cfg.resolve(cfg.paths.case_template_pptx)).exists()


def _new_first_visit_ppt(stage_ppt: Path):
    """
    초진 PPT를 만든다. 반환: (프레젠테이션, 이번 차수를 기록할 십자뷰 슬라이드).

    케이스 양식을 쓸 수 있으면 17장 덱(1~16 + 십자뷰)을 만들고, 아니면 종전처럼
    십자뷰 한 장짜리를 만든다. 어느 쪽이든 사진은 '십자뷰 슬라이드'에 들어간다.
    """
    if _case_deck_ready():
        cd = cfg.case_deck
        prs = CD.build_first_visit_deck(
            cfg.resolve(cfg.paths.case_template_pptx),
            cfg.resolve(cfg.paths.template_pptx),
            stage_ppt,
            keep_slides=cd.keep_slides,
            note_slide_no=cd.note_slide_no,
            # 노트 슬라이드가 이미 날짜 박스를 들고 있으므로 겹치는 INFO_BOX는 뺀다
            cross_drop_shapes={cfg.ppt.info_box_name},
        )
        return prs, prs.slides[CD.cross_slide_index(prs)]

    prs = W.new_ppt_from_template(cfg.resolve(cfg.paths.template_pptx), stage_ppt)
    return prs, prs.slides[0]


# ── 케이스 덱의 얼굴 사진 자리 ────────────────────────────────────────────────
# 자리 이름은 "<슬라이드번호><위치>" — 4L, 4R, 7C, 10BIG.
# 좌표는 양식에서 한 번 읽어 캐시한다(양식은 실행 중에 바뀌지 않는다).
def _cell_key(slide_no: int, pos: str) -> str:
    return f"{slide_no}{pos}"


def _slide_overlays(prs, slide_nos) -> dict[int, list[dict]]:
    """자리 사진 **위에** 겹쳐 있는 양식 도형의 위치 (cm).

    양식의 얼굴 슬라이드는 사진이 슬라이드 아래끝까지 내려오고, 그 위로 캡션 띠가
    덮인다. 검수 판에서 이 도형을 빼고 그리면 사진이 다 보이는 것처럼 착각하게
    되는데 — 실제로는 아래쪽이 가려지므로 얼굴을 그만큼 위로 올려 잡아야 한다.
    화면에도 같이 그려서 그 사실이 보이게 한다.

    **단색으로 채워진 도형만** 센다. 분석 슬라이드(10·11)에는 얼굴 위에 계측선과
    표가 얹히는데 그건 비쳐 보이는 것들이라 가림이 아니다. 채우기가 있는 사각형만
    실제로 사진을 덮는다.
    """
    AUTO_SHAPE, SOLID = 1, 1
    out: dict[int, list[dict]] = {}
    for n in slide_nos:
        if not (0 < n <= len(prs.slides)):
            continue
        rects = []
        for sh in prs.slides[n - 1].shapes:
            if sh.shape_type != AUTO_SHAPE:
                continue
            try:
                if sh.fill.type != SOLID:
                    continue
                x, y = emu_to_cm(sh.left), emu_to_cm(sh.top)
                w, h = emu_to_cm(sh.width), emu_to_cm(sh.height)
            except (TypeError, AttributeError, ValueError):
                continue
            if w > 0 and h > 0:
                rects.append({"x": round(x, 3), "y": round(y, 3),
                              "w": round(w, 3), "h": round(h, 3)})
        if rects:
            out[n] = rects
    return out


_LINE_TYPE = 9      # MSO_SHAPE_TYPE.LINE (직선 연결선)


def _slide_lines(prs, slide_nos) -> dict[int, list[dict]]:
    """양식에 그려진 계측선의 끝점 (cm). 슬라이드 번호 -> 선 목록.

    선은 도형 상자(off/ext)와 뒤집기(flipH/flipV)로 방향이 정해진다 — 상자만
    보면 대각선의 두 끝을 구분할 수 없다. 화면에서 끌어 옮길 수 있어야 하므로
    이름을 열쇠로 준다(한 슬라이드 안에서는 이름이 겹치지 않는다).
    """
    out: dict[int, list[dict]] = {}
    for n in slide_nos:
        if not (0 < n <= len(prs.slides)):
            continue
        items = []
        for sh in prs.slides[n - 1].shapes:
            if sh.shape_type != _LINE_TYPE:
                continue
            try:
                x, y = emu_to_cm(sh.left), emu_to_cm(sh.top)
                w, h = emu_to_cm(sh.width), emu_to_cm(sh.height)
                xf = sh._element.spPr.xfrm
                fh = str(xf.get("flipH", "")).lower() in ("1", "true")
                fv = str(xf.get("flipV", "")).lower() in ("1", "true")
                pt = sh.line.width.pt if sh.line.width else None
            except (TypeError, AttributeError):
                continue
            items.append({
                "id": f"{n}:{sh.name}", "name": sh.name,
                "x1": round(x + (w if fh else 0), 3), "y1": round(y + (h if fv else 0), 3),
                "x2": round(x + (0 if fh else w), 3), "y2": round(y + (0 if fv else h), 3),
                "width_pt": pt,
            })
        if items:
            out[n] = items
    return out


def _load_case_anchors() -> tuple[dict[str, CD.Anchor], tuple[float, float], dict]:
    if not _case_deck_ready():
        return {}, (0.0, 0.0), {}
    cd = cfg.case_deck
    prs = T.load_presentation(cfg.resolve(cfg.paths.case_template_pptx))
    table = CD.read_deck_anchors(prs, cd.face_slides, cd.big_slides,
                                 list(cd.intraoral_slides.values()))
    size = (emu_to_cm(prs.slide_width), emu_to_cm(prs.slide_height))
    overlays = _slide_overlays(prs, cd.mask_slides)
    return {_cell_key(a.slide, a.pos): a for a in table.values()}, size, overlays


CASE_ANCHORS, CASE_SLIDE_CM, CASE_OVERLAYS = _load_case_anchors()


def _load_note_boxes() -> dict:
    """양식에서 노트 칸 4종의 자리를 읽어 둔다 (검수 화면 오버레이용).

    덱을 만들 때 `ensure_note_boxes` 가 쓰는 것과 같은 규칙·같은 순서로 찾으므로,
    화면에서 본 자리와 결과물의 자리가 어긋나지 않는다.
    """
    if not _case_deck_ready():
        return {}
    cd = cfg.case_deck
    prs = T.load_presentation(cfg.resolve(cfg.paths.case_template_pptx))
    return CD.note_box_windows(prs, cd.note_slide_no,
                               range(cd.keep_slides + 1, cd.note_slide_no))


NOTE_BOXES = _load_note_boxes()


def _load_info_boxes() -> dict:
    """양식 첫 장(환자정보)의 텍스트 박스 자리·글꼴·원래 글 (검수 화면용).

    이 장은 사진 자리가 없어 판이 새까맣게 비어 보인다. 화면에 실제 글을 얹어
    무엇이 적히는 장인지 보이게 하고, 거기서 바로 고칠 수 있게 하려는 것이다.
    """
    pi = cfg.patient_info
    if not (_case_deck_ready() and pi.enabled):
        return {}
    prs = T.load_presentation(cfg.resolve(cfg.paths.case_template_pptx))
    if not (1 <= pi.slide_no <= len(prs.slides._sldIdLst)):
        return {}
    out: dict[str, dict] = {}
    for sh in prs.slides[pi.slide_no - 1].shapes:
        if not sh.has_text_frame:
            continue
        try:
            x, y = emu_to_cm(sh.left), emu_to_cm(sh.top)
            w, h = emu_to_cm(sh.width), emu_to_cm(sh.height)
        except (TypeError, AttributeError):
            continue
        if w <= 0 or h <= 0:
            continue
        text = sh.text_frame.text
        out[sh.name] = {
            "x": round(x, 3), "y": round(y, 3), "w": round(w, 3), "h": round(h, 3),
            "font": CD._run_font(sh), "text": text,
            # 우리가 채우는 줄을 들고 있는 박스만 고칠 수 있다
            "editable": any(ln.strip().startswith(head)
                            for ln in text.split("\n") for head in pi.lines),
        }
    return out


INFO_BOXES = _load_info_boxes()


def _load_case_lines() -> dict[int, list[dict]]:
    if not _case_deck_ready():
        return {}
    prs = T.load_presentation(cfg.resolve(cfg.paths.case_template_pptx))
    return _slide_lines(prs, range(1, cfg.case_deck.keep_slides + 1))


CASE_LINES = _load_case_lines()
# 사용자가 직접 채우는 얼굴 자리 (구내 슬라이드와 파생 자리는 제외)
FACE_CELLS: list[str] = [k for k, a in CASE_ANCHORS.items()
                         if a.slide in set(cfg.case_deck.face_slides)]
# 파생 자리: 슬라이드 4 좌측 사진을 그대로 다시 쓴다. 사용자가 따로 고르지 않는다.
MIRROR_SOURCE = _cell_key(4, "L")
MIRROR_CELLS: list[str] = [k for k, a in CASE_ANCHORS.items()
                           if a.slide in set(cfg.case_deck.big_slides)]


def _face_layout_json() -> dict:
    def cell(k):
        a = CASE_ANCHORS[k]
        return {"cell": k, "slide": a.slide, "pos": a.pos,
                "label": cfg.case_deck.slide_labels.get(a.slide, ""),
                "x": round(a.window.x, 3), "y": round(a.window.y, 3),
                "w": round(a.window.w, 3), "h": round(a.window.h, 3)}
    return {
        "enabled": bool(CASE_ANCHORS),
        "slide_w": round(CASE_SLIDE_CM[0], 3), "slide_h": round(CASE_SLIDE_CM[1], 3),
        "cells": [cell(k) for k in sorted(FACE_CELLS, key=lambda k: (CASE_ANCHORS[k].slide, k))],
        "mirrors": [{**cell(k), "from": MIRROR_SOURCE}
                    for k in sorted(MIRROR_CELLS, key=lambda k: CASE_ANCHORS[k].slide)],
        # 사진 위를 덮는 양식 도형(캡션 띠 등). 슬라이드 번호 -> 사각형 목록.
        "overlays": {str(n): r for n, r in CASE_OVERLAYS.items()},
        # 양식에 그려진 계측선. 화면에서 끌어 옮길 수 있다.
        "lines": {str(n): v for n, v in CASE_LINES.items()},
        # 구내 한 장짜리 슬라이드(12~16). 여기는 사람이 배정하는 자리가 아니라
        # **구내 검수에서 정한 것이 그대로 오는** 자리다. 그래도 그려 줘야 한다 —
        # 안 그리면 그 장을 열었을 때 판이 새까맣고, 사진이 안 들어간 것처럼 보인다.
        "intraoral": [
            {"slot": slot, "slide": no,
             "x": round(a.window.x, 3), "y": round(a.window.y, 3),
             "w": round(a.window.w, 3), "h": round(a.window.h, 3)}
            for slot, no in cfg.case_deck.intraoral_slides.items()
            if (a := CASE_ANCHORS.get(_cell_key(no, "FULL"))) is not None
        ],
        # 훑어볼 수 있는 슬라이드 전체(양식에서 그대로 가져오는 앞쪽 장들).
        "slides": list(range(1, cfg.case_deck.keep_slides + 1)) if CASE_ANCHORS else [],
        "labels": {str(k): v for k, v in cfg.case_deck.slide_labels.items()},
    }


def _place_faces(prs, s: "Session") -> dict[str, Path]:
    """
    배정된 얼굴 사진을 케이스 덱의 자리에 넣는다.

    양식의 자리잡이 사진(플레이스홀더)은 지우고, 검수 화면에서 잡은 구도
    (session.face_editors)대로 놓는다. 잡아 둔 것이 없으면 cover-fit.
    파생 자리(10·11)는 슬라이드 4 좌측 사진을 그대로 다시 쓴다.

    반환: `{사진 id: 구워 넣은 파일}`. 환자 폴더에 **슬라이드와 같은 그림**을
    저장하려면 호출부가 이게 필요하다. 한 사진이 여러 자리에 들어가면(10·11)
    첫 자리 것을 준다 — 같은 구도를 크기만 달리해 놓은 것이라 어느 쪽이든 같은
    그림이고, 자리 순서가 정해져 있으니 결과가 흔들리지 않는다.
    """
    bakes: dict[str, Path] = {}
    for cell, pid in sorted(_face_slots_json(s).items()):
        anchor = CASE_ANCHORS.get(cell)
        if anchor is None:
            continue
        slide = prs.slides[anchor.slide - 1]
        # 양식에 놓여 있던 자리잡이 사진을 걷어낸다 — 남겨두면 뒤에 비쳐 보인다.
        for sh in list(slide.shapes):
            if sh.shape_type == 13 and sh.name == anchor.shape_name:
                sh._element.getparent().remove(sh._element)
        photo = _photo(s, pid)
        win = anchor.window
        st = _face_editor(s, cell)
        # 얼굴 슬라이드에는 가려 줄 마스크가 없다 — 구워 넣어야 옆 사진을 안 덮는다
        baked, bwh = _bake_window(photo, win, st, False, s.tmp / f"bake_face_{cell}.jpg")
        if baked:
            W.place_photo_in_window(slide, f"FACE_{cell}", win, baked, bwh,
                                    placement=_exact_placement(win),
                                    letterbox_color=_letterbox_color())
            bakes.setdefault(pid, baked)
        else:
            bw, bh = cover_base_ext_cm(photo.w, photo.h, win)
            pl = editor_to_placement(st, win, bw, bh, PPC)
            W.place_photo_in_window(slide, f"FACE_{cell}", win, photo.path,
                                    (photo.w, photo.h), placement=pl,
                                    letterbox_color=_letterbox_color())
    return bakes


def _place_intraoral(prs, s: "Session") -> list[int]:
    """구내 사진을 **한 장씩 크게** 싣는 슬라이드(12~16)를 채운다.

        FRONT → 12    RIGHT → 13    LEFT → 14    UPPER → 15    LOWER → 16

    십자뷰에 들어가는 것과 **같은 사진, 같은 구도**다. 창만 3배 크다.

    구도를 옮겨 오는 방법은 슬라이드 10·11 이 4L 을 물려받는 것과 같다 — 창의
    비율이 같을 때만 성립하므로(십자뷰 슬롯 1.3333 vs 여기 1.3340) 그 검사를 하고,
    어긋나면 cover-fit 으로 물러난다. dx·dy 는 cm 로 환산되는 절대량이라 창 폭
    비율만큼 같이 키워야 잘린 영역이 같아진다.

    환자 폴더로 가는 사본은 여기서 만들지 않는다 — 십자뷰 쪽 사본을 쓴다. 재진에는
    이 슬라이드가 없으니, 그래야 초진·재진이 같은 파일을 남긴다.
    """
    filled: list[int] = []
    for slot, slide_no in cfg.case_deck.intraoral_slides.items():
        anchor = CASE_ANCHORS.get(_cell_key(slide_no, "FULL"))
        members = s.bins.get(slot) or []
        if anchor is None or not members:
            continue
        photo = _photo(s, members[0])
        slide = prs.slides[anchor.slide - 1]
        # 양식의 빈 내용 개체 틀을 걷어낸다 — 두면 사진 위에 안내문이 얹힌다.
        for sh in list(slide.shapes):
            if sh.is_placeholder and sh.name == anchor.shape_name:
                sh._element.getparent().remove(sh._element)

        win = anchor.window
        src = s.slot_windows[slot]
        st = photo.editor
        if abs(src.w / src.h - win.w / win.h) < 1e-3:
            k = win.w / src.w
            st = EditorState(st.dx_px * k, st.dy_px * k, st.scale, st.angle_deg)
        else:
            st = EditorState(0.0, 0.0, 1.0, st.angle_deg)

        baked, bwh = _bake_window(photo, win, st, photo.flip_v,
                                  s.tmp / f"bake_io_{slot}.jpg")
        if baked:
            W.place_photo_in_window(slide, f"IO_{slot}", win, baked, bwh,
                                    placement=_exact_placement(win),
                                    letterbox_color=_letterbox_color())
        else:
            bw, bh = cover_base_ext_cm(photo.w, photo.h, win)
            pl = editor_to_placement(st, win, bw, bh, PPC)
            W.place_photo_in_window(slide, f"IO_{slot}", win, photo.path,
                                    (photo.w, photo.h), placement=pl,
                                    letterbox_color=_letterbox_color(),
                                    flip_v=photo.flip_v)
        filled.append(slide_no)
    return filled


def _exact_placement(win: WindowCm) -> Placement:
    """창과 **똑같은** 자리·크기의 배치.

    창에 맞춰 구운 사진은 창이 곧 제 크기다. 그런데 배치를 안 주면 cover-fit 이
    구운 파일의 **정수 픽셀 비율**로 크기를 다시 셈하는데, 그 값이 창 비율과 아주
    조금 어긋난다. 0.002cm 남짓이지만 하필 소수 둘째 자리 경계를 넘으면
    PowerPoint 에 8.38 이 8.39 로 보인다 — 직전 차수 사진과 크기가 달라 보인다.
    """
    e = EMU_PER_CM
    return Placement(off_x=int(round(win.x * e)), off_y=int(round(win.y * e)),
                     ext_cx=int(round(win.w * e)), ext_cy=int(round(win.h * e)),
                     rot=0)


def _bake_window(photo, win: WindowCm, st: EditorState, flip_v: bool, dst: Path):
    """창에 보이는 그림만 잘라 dst 에 굽는다. (경로, 크기) 또는 (None, None).

    이렇게 넣으면 사진이 창과 정확히 같은 크기가 되어 옆 슬롯을 침범할 수 없다.
    회전·상하반전도 픽셀에 들어가므로 PPT 도형은 회전 0, flipV 없음이 된다.
    원본 파일은 환자 폴더에 따로 저장되므로 여기서 굽는 것은 사본뿐이다.
    """
    ppcm = cfg.geometry.export_px_per_cm
    if not ppcm:
        return None, None                    # 굽기 꺼짐 — 종전 방식
    arr = _imread(photo.path)
    if arr is None:
        return None, None                    # 못 읽으면 종전 방식으로 물러난다
    out = Cr.render_window(arr, win, st, flip_v, ppcm, PPC,
                           Cr.hex_to_bgr(_letterbox_color()))
    # EXIF 를 원본에서 옮겨 심는다 — cv2.imwrite 는 EXIF 를 쓸 줄 모른다. 스테이징
    # 원본은 회전이 이미 픽셀에 구워져 있어(orientation ≤1) 그대로 옮겨도 뷰어가
    # 이중 회전하지 않는다. 화질 규약은 종전과 동일 (q95, subsampling=0 = 4:4:4).
    try:
        from PIL import Image as _Im                              # noqa: PLC0415
        with _Im.open(photo.path) as _src:
            exif = _src.info.get("exif")
        _Im.fromarray(out[:, :, ::-1]).save(
            dst, "JPEG", quality=95, subsampling=0,
            **({"exif": exif} if exif else {}))
    except Exception:                                             # noqa: BLE001
        return None, None
    # 촬영시각을 옮겨 심는다. 이 사본이 환자 폴더로 갈 수 있고, 그때 mtime 이
    # '구운 시각'이면 차수 날짜가 오늘로 보인다 — 목록이 그걸로 날짜를 만든다.
    try:
        st = photo.path.stat()
        os.utime(dst, (st.st_atime, st.st_mtime))
    except OSError:
        pass
    return dst, (out.shape[1], out.shape[0])


def _face_editor(s: "Session", cell: str) -> EditorState:
    """자리의 편집기 값. 파생 자리는 원본 자리 값을 창 크기에 맞춰 옮겨 온다.

    파생 자리(10·11)는 슬라이드 4 좌측과 **비율이 같고 크기만 다르다**. dx·dy 는
    cm 로 환산되는 절대량이라 그대로 쓰면 큰 창에서 구도가 밀린다. 창 폭 비율만큼
    같이 키우면 잘린 영역이 정확히 같아진다. 비율이 어긋나는 양식이 오면 그 환산이
    성립하지 않으므로 cover-fit 으로 물러난다.
    """
    if cell in s.face_editors:
        return s.face_editors[cell]
    src = MIRROR_SOURCE if cell in MIRROR_CELLS else None
    if src and src in s.face_editors:
        a, b = CASE_ANCHORS.get(src), CASE_ANCHORS.get(cell)
        if a and b and abs(a.window.w / a.window.h - b.window.w / b.window.h) < 1e-3:
            k = b.window.w / a.window.w
            st = s.face_editors[src]
            return EditorState(st.dx_px * k, st.dy_px * k, st.scale, st.angle_deg)
    return EditorState()


# ── 차수 노트 ─────────────────────────────────────────────────────────────────
_FIELD_RE = re.compile(r"\{([a-z_]+)\}")


def _render_note_line(line: str, fields: dict[str, str]) -> str | None:
    """
    한 줄을 채운다. 그 줄이 쓰는 칸이 모두 비어 있으면 줄 자체를 버린다.
    'U: ' 같은 빈 껍데기를 슬라이드에 남기지 않기 위한 것이다.
    """
    keys = _FIELD_RE.findall(line)
    if keys and not any((fields.get(k) or "").strip() for k in keys):
        return None
    return _FIELD_RE.sub(lambda m: (fields.get(m.group(1)) or "").strip(), line)


def _render_notes(fields: dict[str, str], templates: dict[str, str] | None = None) -> dict[str, str]:
    """칸 → 박스별 최종 텍스트.

    **빈 줄은 서식에 적힌 그대로 남긴다.** 양식의 박스는 빈 문단으로 글의 시작
    높이를 잡아 두는데(예: 우상단은 Tx. Period 앞에 빈 줄 세 개), 앞뒤를 다듬어
    버리면 슬라이드에서 글이 위로 밀려 양식과 어긋난다. 사라지는 것은 '칸을
    쓰는데 그 칸이 빈' 줄뿐이다(_render_note_line).
    """
    out: dict[str, str] = {}
    for box, tpl in (templates or cfg.notes.boxes).items():
        lines = [_render_note_line(ln, fields) for ln in tpl.split("\n")]
        out[box] = "\n".join(ln for ln in lines if ln is not None)
    return out


# `9 month(24.09.26, MARPE)` 처럼 괄호에 장치명이 따라붙기도 한다 — 날짜만 뽑는다.
_PERIOD_DATE = re.compile(r"\((\d{2}\.\d{2}\.\d{2})")
PERIOD_KEYS = {"tx": "Tx. Period", "rx": "Rx. Period", "app": "App. Period"}


def _period_line(text: str, label: str) -> str | None:
    """상태 칸에서 그 기간의 줄 하나. 없으면 None."""
    for line in (text or "").splitlines():
        if line.strip().startswith(label):
            return line.strip()
    return None


def _period_history(visits, label: str) -> dict:
    """이전 차수들에서 그 기간의 **기준일 이력**과 **직전 값**.

    기준일은 `Rx. Period: 0 month (24.08.12)` 처럼 괄호에 적힌다 — PPT 가 차수 이력의
    단일 출처이므로 별도 저장소를 두지 않는다. 실측 기록을 보면 기준일은 **여러 번
    바뀐다**(재진단·치료 단계 전환). 그래서 하나를 고르는 대신 목록으로 낸다.

    반환: {"dates": [최신순], "last": "직전 차수에 적힌 줄 그대로"}
    """
    dates, last = [], None
    for v in sorted(visits, key=lambda z: (z.date or "", z.slide_index)):
        line = _period_line(getattr(v, "status_text", ""), label)
        if line is None:
            continue
        last = line
        m = _PERIOD_DATE.search(line)
        if m and m.group(1) not in dates:
            dates.append(m.group(1))
    return {"dates": list(reversed(dates)), "last": last}


def _months_since(start: str | None, now: datetime, unit: str = "int"):
    """기준일부터 지금까지의 개월 — **가까운 단위로 반올림** (2026-08-12 결정).

    실측 기록이 `12.5 month`·`13.5 month` 처럼 반달을 쓴다. 어느 쪽을 쓸지는
    사람마다 다르므로 설정에서 고르게 하고, 기본은 정수다.

    예전에는 찬 개월만 셌다(내림) — 한 달에서 하루 모자라도 이전 값이라 기록
    관행(가까운 쪽으로 적는다)과 어긋났다. 지금은:

        정수:  n개월 +  0~14일 → n     · 15일~ → n+1
        0.5:   n개월 +  0~ 7일 → n     ·  8~22일 → n.5  · 23일~ → n+1
    """
    if not start:
        return None
    try:
        d = datetime.strptime(start.strip().rstrip("."), "%y.%m.%d")
    except ValueError:
        return None
    n = (now.year - d.year) * 12 + (now.month - d.month)
    rest = now.day - d.day
    if rest < 0:
        n -= 1
        rest += 30                      # 대략. 반올림 판정에만 쓴다
    if unit == "half":
        return max(0.0, n + (1.0 if rest >= 23 else 0.5 if rest >= 8 else 0.0))
    return max(0, n + (1 if rest >= 15 else 0))


def _fmt_months(v) -> str:
    """`16` → "16" · `16.5` → "16.5" · `16.0` → "16" (정수는 소수점 없이)."""
    if v is None:
        return ""
    return str(int(v)) if float(v) == int(v) else str(v)


def _first_visit_date(visits) -> str | None:
    """기존 PPT의 슬라이드들에서 초진 날짜를 고른다.

    '초진'이라고 적힌 슬라이드를 먼저 믿고, 없으면 가장 이른 날짜를 쓴다
    (초진 슬라이드의 라벨이 지워졌거나 차수 표기가 흐트러진 PPT 대비).
    """
    dated = [v for v in visits if v.date]
    if not dated:
        return None
    first = [v for v in dated if v.kind == "first"]
    return (first[0] if first else min(dated, key=lambda v: v.date)).date


def _months_between(a: str | None, b: datetime) -> int | None:
    """'YY.MM.DD' 부터 b 까지의 개월 수. 날짜를 못 읽으면 None."""
    if not a:
        return None
    try:
        d = datetime.strptime(a.strip().rstrip("."), "%y.%m.%d")
    except ValueError:
        return None
    n = (b.year - d.year) * 12 + (b.month - d.month)
    if b.day < d.day:        # 같은 달 안에서 아직 날짜가 안 찼으면 한 달 덜 센다
        n -= 1
    return max(0, n)


def _photo_date(s: "Session") -> datetime | None:
    """이 차수의 날짜 = **사진 촬영일** (EXIF, 최빈값).

    노트의 날짜는 PPT 를 만드는 작업일이 아니라 사진을 찍은 날이어야 한다 —
    찍고 며칠 뒤에 정리하는 일이 흔하다. 최빈값을 쓰는 이유: 카메라 시계가 틀린
    사진이 섞여도 하루가 다수결로 이긴다. EXIF 가 하나도 없으면 None.
    """
    days = [p.taken_at.date() for p in s.photos if p.taken_at]
    if not days:
        return None
    day = Counter(days).most_common(1)[0][0]
    return datetime(day.year, day.month, day.day)


def _note_auto(s: "Session") -> dict:
    """서식이 쓸 수 있는 자동 계산값.

    사용자가 손으로 고칠 수 있어야 하므로 여기서는 **기본값만** 만든다 —
    한 번 고친 칸은 s.note_fields 가 이기고, 이 값들은 다시 덮어쓰지 않는다.

    "지금"은 촬영일(`_photo_date`)이다 — 날짜 칸과 개월 계산이 같은 날을 봐야
    `0 month (날짜)` 줄이 안 어긋난다. EXIF 가 없으면 작업일로 물러난다.
    """
    now = _photo_date(s) or datetime.now()
    today = now.strftime(cfg.ppt.info_date_format)
    months = _months_between(s.first_date, now)
    unit = _months_unit()

    def _per(k):
        """(개월 문자열, 기준일, 괄호, 직전값 그대로 쓸지)."""
        hist = s.period_hist.get(k) or {}
        chosen = s.period_start.get(k) or ""
        if chosen == "none":                # 체크를 다 푼 상태 — 기준일 없이 쓴다
            return "0", "", ""
        # 기준일은 사람이 고른 것 > 이전 차수에 적힌 것 순이다. 없으면 없는 대로
        # 둔다 — Rx 도 마찬가지다(초진일을 자동으로 넣지 않는다, 2026-08-14 결정).
        start = chosen or (hist.get("dates") or [None])[0]
        m = _months_since(start, now, unit)
        return (_fmt_months(m) if m is not None else "0",
                start or "", f" ({start})" if start else "")

    per = {k: _per(k) for k in PERIOD_KEYS}
    return {
        "first_date": s.first_date or "",
        "today": today,
        "months": "" if months is None else str(months),
        # 기간 셋. 기준일은 사람이 목록에서 고르고, 기본은 가장 최근 것이다.
        # 괄호를 통째로 내주는 이유: 서식 문자열은 조건을 못 써서, 날짜만 비우면
        # `0 month ()` 처럼 빈 괄호가 남는다.
        "months_tx": per["tx"][0], "tx_date": per["tx"][1], "tx_paren": per["tx"][2],
        "months_rx": per["rx"][0], "rx_date": per["rx"][1], "rx_paren": per["rx"][2],
        "months_app": per["app"][0], "app_date": per["app"][1], "app_paren": per["app"][2],
        "visit": s.visit or "",
        "visit_label": _render_label(today, s.visit or "",
                                     getattr(s, "label_fp", None)),
    }


def _period_options(s: "Session", k: str) -> list[str]:
    """기준일로 고를 만한 날짜 — 최신순.

    기간 줄에 적혀 있던 기준일 + **덱에서 파싱한 차수 날짜 전부** + 이번 차수.
    기준일은 대개 "어느 차수부터" 라서, 이력에 적힌 것만 주면 처음 정하는 기간은
    고를 것이 하나도 없다.
    """
    hist = (s.period_hist.get(k) or {}).get("dates") or []
    today = (_photo_date(s) or datetime.now()).strftime(cfg.ppt.info_date_format)
    # 날짜 → 차수 글자. 어느 차수의 날짜인지 보여야 "그 차수부터" 를 고를 수 있다.
    letter = {v["date"]: v["visit"] for v in (s.visit_dates or []) if v.get("date")}
    seen = {d for d in (*hist, *letter, today) if d}
    return [{"date": d, "visit": ("이번 차수" if d == today and d not in letter
                                  else letter.get(d, ""))}
            for d in sorted(seen, key=lambda d: Rd._date_key(d) or (0, 0, 0),
                            reverse=True)]


def _note_defaults(s: "Session") -> dict:
    """자동값을 끼운 칸별 기본값. 빈 칸에만 쓰인다."""
    auto = _note_auto(s)
    # "이전 차수 값 그대로" — 개월을 다시 세지 않고 직전 줄을 그대로 쓴다.
    # 실측에서 App. Period 가 여러 차수에 걸쳐 `9 month(24.09.26, MARPE)` 로
    # 멈춰 있던 경우가 이것이다 (장치를 뗀 시점에 고정).
    keep_line = {}
    for k, label in PERIOD_KEYS.items():
        if not s.period_keep.get(k):
            continue
        last = (s.period_hist.get(k) or {}).get("last") or ""
        if last.startswith(label):
            keep_line[f"{k}_period"] = last[len(label):].lstrip(": ").strip()
    out = {}
    for f in cfg.notes.fields:
        if f.key in keep_line:
            out[f.key] = keep_line[f.key]
            continue
        tmpl = getattr(f, "default", "") or ""
        if not tmpl:
            continue
        try:
            v = tmpl.format(**auto)
        except (KeyError, IndexError):
            continue
        if v.strip():
            out[f.key] = v
    return out


def _note_values(s: "Session") -> dict:
    """실제로 서식에 들어갈 값 = 기본값 위에 사용자가 고친 값을 덮은 것."""
    v = _note_defaults(s)
    v.update({k: t for k, t in s.note_fields.items()})
    return v


def _note_templates(s: "Session") -> dict[str, str]:
    """이 차수에 쓸 박스별 서식 = 설정값 위에 세션에서 고친 것을 덮은 것."""
    t = dict(cfg.notes.boxes)
    t.update(s.note_templates)
    return t


def _note_text(s: "Session") -> dict[str, str]:
    """박스별 최종 텍스트 = 서식 결과 위에 사람이 통째로 고쳐 쓴 것을 덮은 것.

    날짜 칸(NOTE_DATE)만 서식이 아니라 차수 라벨에서 온다 — 원래 커밋할 때
    `_write_visit_label` 이 쓰던 값이다. 화면에도 같이 내보내서 네 칸을 한 자리에서
    보고 고칠 수 있게 한다.
    """
    out = _render_notes(_note_values(s), _note_templates(s))
    out.setdefault(CD.NOTE_DATE, _note_auto(s)["visit_label"])
    out.update({k: v for k, v in s.note_overrides.items()})
    return out


def _notes_json(s: "Session") -> dict:
    return {
        # 기간별 기준일 목록·직전값·현재 선택. 화면이 체크박스들을 그린다.
        # today = 이 차수의 날짜(촬영일) — 기준일 후보 목록의 "이 차수" 항목이다.
        "periods": {
            k: {"dates": (s.period_hist.get(k) or {}).get("dates", []),
                "options": _period_options(s, k),
                "last": (s.period_hist.get(k) or {}).get("last") or "",
                "start": s.period_start.get(k, ""),
                "keep": bool(s.period_keep.get(k)),
                "today": (_photo_date(s) or datetime.now())
                         .strftime(cfg.ppt.info_date_format)}
            for k in PERIOD_KEYS},
        "fields": [f.model_dump() for f in cfg.notes.fields],
        "values": _note_values(s),
        "edited": dict(s.note_fields),      # 사람이 직접 고친 칸만
        "auto": _note_auto(s),
        "defaults": _note_defaults(s),
        "order": CD.NOTE_ORDER,             # 화면에 붙는 ①~⑤ 번호 순서
        "layout": NOTE_BOXES,               # 양식에서 읽은 칸 위치(cm)+글꼴
        "overrides": dict(s.note_overrides),
        "templates": _note_templates(s),    # 지금 쓰이는 서식
        "templates_default": dict(cfg.notes.boxes),
        "slide": {"w": CASE_SLIDE_CM[0], "h": CASE_SLIDE_CM[1]},
        # 직전 차수 슬라이드의 선 — 설정이 '복사 안 함' 이면 따라오지 않으므로 안 낸다
        "prev_lines": ([] if _copy_shapes() == "none"
                       else list(getattr(s, "prev_lines", []) or [])),
        "preview": _note_text(s),
        # 줄 끝 괄호(날짜)를 작게 쓰는 크기. 화면 오버레이도 같은 규칙을 써야
        # 슬라이드에서 보일 모습과 어긋나지 않는다.
        "date_pt": cfg.notes.date_pt,
        "date_pt_except": list(cfg.notes.date_pt_except),
        # 박스별로 그 박스를 채우는 칸 이름 — 판 위 박스를 누르면 이 칸들만
        # 오른쪽에 띄운다. 서식({칸이름})을 사람이 읽을 필요가 없게 하려는 것이다.
        "box_fields": {box: _FIELD_RE.findall(tpl)
                       for box, tpl in _note_templates(s).items()},
        # 양식 첫 장(환자정보). 초진 덱을 만들 때만 쓴다.
        "patient_info": {
            "enabled": bool(cfg.patient_info.enabled and cfg.patient_info.lines
                            and s.mode == "first" and CASE_ANCHORS),
            "slide": cfg.patient_info.slide_no,
            "boxes": INFO_BOXES,              # 양식에서 읽은 자리(cm)+글꼴
            "preview": _info_preview(s),      # 확정 뒤에 적힐 글
            "fields": [f.key for f in cfg.notes.fields if f.group == "patient"],
        },
    }


# ── 환자정보 슬라이드 (양식 첫 장) ────────────────────────────────────────────
_PAREN_GROUP = re.compile(r"\(([^()]*)\)")


def _render_info_line(tpl: str, fields: dict[str, str]) -> str:
    """머리말 뒤에 올 글을 만든다. 빈 칸이 만든 껍데기는 걷어낸다.

    "( {sex}/{age} )" 처럼 괄호로 묶인 곳은 그 안의 칸이 **모두** 비면 괄호째
    사라지고, 한쪽만 비면 남은 구분자(`/`)만 정리한다 — "(M/)" 같은 자국을
    슬라이드에 남기지 않기 위한 것이다.
    """
    def group(m):
        inner = _FIELD_RE.sub(lambda x: (fields.get(x.group(1)) or "").strip(), m.group(1))
        inner = re.sub(r"\s*/\s*", "/", inner).strip(" /")
        return f"({inner})" if inner else ""

    out = _PAREN_GROUP.sub(group, tpl)
    out = _FIELD_RE.sub(lambda m: (fields.get(m.group(1)) or "").strip(), out)
    return re.sub(r"[ \t]{2,}", " ", out).rstrip()


def _info_fields(s: "Session") -> dict[str, str]:
    """환자정보 줄이 쓸 수 있는 값 = 노트 칸 + 환자를 고를 때 이미 받은 식별자.

    폴더명에서 못 얻은 식별자(빈 값)는 덮어쓰지 않는다 — 병원번호가 폴더명에
    없는 환자는 노트의 병원번호 칸에 사람이 적은 값이 슬라이드로 간다.
    """
    ids = {"name": s.ids.name, "hospital_id": s.ids.hospital_id,
           "ortho_id": s.ids.ortho_id}
    return {**_note_values(s), **{k: v for k, v in ids.items() if v}}


def _info_preview(s: "Session") -> dict[str, str]:
    """확정 뒤 첫 장의 각 텍스트 박스가 갖게 될 글. 화면에 그대로 얹는다.

    `case_deck.fill_prefixed_lines` 와 **같은 규칙**으로 만든다 — 머리말이 처음
    나온 줄만, 그 줄이 쓰는 칸이 다 비면 손대지 않고. 규칙이 갈라지면 화면에서
    본 것과 슬라이드에 적히는 것이 달라진다.
    """
    fields = _info_fields(s)
    out: dict[str, str] = {}
    for name, box in INFO_BOXES.items():
        used: set[str] = set()
        lines = []
        for ln in box["text"].split("\n"):
            head = next((h for h in cfg.patient_info.lines
                         if h not in used and ln.strip().startswith(h)), None)
            tpl = cfg.patient_info.lines.get(head) if head else None
            keys = _FIELD_RE.findall(tpl) if tpl else []
            if head is None or (keys and not any((fields.get(k) or "").strip() for k in keys)):
                lines.append(ln)                     # 양식 그대로
            else:
                used.add(head)
                lines.append(head + _render_info_line(tpl, fields))
        out[name] = "\n".join(lines)
    return out


def _fill_patient_info(prs, s: "Session") -> dict[str, str]:
    """양식 첫 장(환자정보)을 이 환자의 값으로 채운다. 반환: 채운 줄."""
    pi = cfg.patient_info
    if not (pi.enabled and pi.lines):
        return {}
    if not (1 <= pi.slide_no <= len(prs.slides._sldIdLst)):
        return {}
    fields = _info_fields(s)
    filled = {}
    for head, tpl in pi.lines.items():
        # 그 줄이 쓰는 칸이 모두 비면 손대지 않는다 — 양식의 안내 문구를 지우면
        # 사람이 나중에 무엇을 적어야 할 자리인지 알 수 없게 된다.
        keys = _FIELD_RE.findall(tpl)
        if keys and not any((fields.get(k) or "").strip() for k in keys):
            continue
        filled[head] = _render_info_line(tpl, fields)
    if not filled:
        return {}
    return CD.fill_prefixed_lines(prs.slides[pi.slide_no - 1], filled)


def _letterbox_color() -> str:
    """회전·축소로 드러나는 빈 자리를 채울 색 (RGB hex, '#' 없이).

    화면 캔버스·저장되는 JPG·PPT 배경판이 **같은 값**을 봐야 한다. 셋이 갈리면
    검수 화면에서 본 모습과 결과물이 달라진다. 설정이 없으면 config 기본값.
    """
    try:
        v = json.loads(SETTINGS_FILE.read_text(encoding="utf-8")).get("letterbox_color")
        if isinstance(v, str) and re.fullmatch(r"[0-9A-Fa-f]{6}", v):
            return v.upper()
    except Exception:                                   # noqa: BLE001
        pass
    return cfg.geometry.letterbox_color


def _small_pt(box: str) -> float | None:
    """그 박스에서 줄 끝 괄호를 작게 쓸 크기. 줄이지 않는 박스면 None.

    날짜 칸의 "(초진 A)" 는 날짜가 아니라 차수 표시다 — 작게 쓰면 안 된다.
    """
    if box in cfg.notes.date_pt_except:
        return None
    return cfg.notes.date_pt or None


def _write_visit_label(slide, text: str, style: dict | None = None,
                       pin: bool = True) -> None:
    """
    이번 차수를 슬라이드에 적는다.

    십자뷰 한 장짜리 PPT는 INFO_BOX에, 케이스 덱은 양식의 날짜 박스(NOTE_DATE)에
    적힌다. 어느 쪽이 있든 한 곳에만 쓴다.

    위치는 (0.40, 0.40)cm 로 고정한다 — 사람이 만든 PPT 의 박스는 자리가 제각각이라
    초진(케이스 덱 상수)과 재진(여기)이 어긋나 보였다.
    """
    for name in (cfg.ppt.info_box_name, CD.NOTE_DATE):
        if CD.set_note_text(slide, name, text):
            if style:   # 수제 PPT 상속 — 원본 라벨 박스의 자리·폰트를 따른다
                CD.style_note_box(slide, name, style)
            elif pin:   # 통째 복사 상속이면 원본 자리를 지킨다 — 고정하지 않는다
                CD.pin_shape(slide, name, 0.40, 0.40)
                # 템플릿 라벨은 색이 테마(bg1)라 수제 테마에서 검정이 될 수 있다
                CD.force_run_color(slide, name, "FFFFFF")
            return


PPC = cfg.geometry.render_px_per_cm


@contextlib.asynccontextmanager
async def _lifespan(_app: FastAPI):
    """청소 스레드는 서버 수명과 함께 산다. TestClient를 컨텍스트 없이 쓰면 뜨지 않는다."""
    _log_framer()
    _apply_note_sizes()                        # 설정된 노트 글자 크기 반영
    stop = threading.Event()
    threading.Thread(target=_sweeper_loop, args=(stop,), daemon=True).start()
    yield
    stop.set()


def _log_framer() -> None:
    """자동 프레이밍 모델 상태를 기동 로그에 한 번 찍는다."""
    if framer is None:
        print("[프레이밍] 모델 없음 — 초진은 cover-fit(회전 0, 중심)으로 놓입니다.")
        return
    n = {len(v) for v in framer.files.values()}
    kind = "배포본" if n == {1} else f"fold {max(n)}개 앙상블"
    print(f"[프레이밍] {framer.meta.get('tag')} / {kind} / "
          f"{framer.iw}x{framer.ih} / 클래스 {len(framer.files)}개")
    ph = framer.placeholder
    if ph:
        print(f"  ⚠ 임시 대역 모델입니다 — {ph.get('reason')}")
        print(f"    {ph.get('risk')}")
        print(f"    교체: {ph.get('replace_with')}")


app = FastAPI(title="교정과 사진 자동화", lifespan=_lifespan)

SESSIONS: dict[str, "Session"] = {}
# 환자 루트는 실행 중에 바뀔 수 있다(Setup의 '저장 위치 변경').
# config.yaml은 기본값, settings.json이 있으면 그쪽이 이긴다 — 주석 달린 YAML을
# 프로그램이 다시 쓰면 형식이 망가지므로 사용자 선택은 별도 파일에 둔다.
# **설정은 프로그램 폴더 안**, **환자 자료는 밖**. 둘은 성격이 다르다:
#
#     settings.json   프로그램 설정. 잃어도 저장 위치를 한 번 다시 고르면 끝이고,
#                     안에 있으면 폴더를 옮길 때 따라간다. `.gitignore` 에 있어
#                     `git pull` 이 안 건드린다
#     환자 자료        의료 기록. 저장소에 실수로 들어가면 이력에서 못 지우고,
#                     `git clean -xfd` 한 줄에 사라진다. 밖이어야 한다
PROGRAM_DIR = BACKEND_DIR.parents[1]
SETTINGS_FILE = PROGRAM_DIR / "settings.json"


def _patient_files(d: Path) -> list[Path]:
    """환자 폴더의 파일들 — 두 단계 하위 폴더까지.

    사진 파일을 열지는 않는다. 이 목록의 용도는 둘뿐이다: 폴더 내용을 화면에
    그대로 보여주는 것, 그리고 하위 폴더에 있어도 PPT 를 찾아내는 것.
    """

    def walk(base: Path, depth: int, out: list[Path]) -> None:
        for f in sorted(base.iterdir(), key=lambda x: x.name.lower()):
            if f.name.startswith("."):
                continue
            if f.is_dir():
                if depth < 2:
                    walk(f, depth + 1, out)
            elif f.is_file():
                out.append(f)

    out: list[Path] = []
    walk(d, 0, out)
    return out


LEGACY_SETTINGS = Path.home() / "ortho-webapp" / "settings.json"


def _months_unit() -> str:
    """개월 표기 단위 — "int" 또는 "half".

    실측 기록이 `12.5 month` 처럼 반달을 쓴다. 사람마다 관행이 달라 고르게 하되,
    **`settings.json`(설치본 공용)** 에 둔다 — 브라우저마다 다르면 한 환자 PPT 안에서
    표기가 섞인다.
    """
    try:
        v = json.loads(SETTINGS_FILE.read_text(encoding="utf-8")).get("months_unit")
        return v if v in ("int", "half") else "int"
    except Exception:                                   # noqa: BLE001
        return "int"


def _save_raw() -> bool:
    """원본을 함께 남길까. 기본은 **아니오**.

    환자 폴더에는 슬라이드에 실린 그대로(잘린 사진)가 간다 — 폴더와 PPT 가 다른
    그림이면 나중에 어느 쪽이 진짜인지 다투게 된다.

    켜면 원본이 `_raw` 이름으로 함께 남는다. **끄면 원본은 어디에도 남지 않는다** —
    확정과 함께 업로드 임시본이 지워지고, 잘라낸 영역은 되돌릴 수 없다.
    """
    try:
        return bool(json.loads(SETTINGS_FILE.read_text(encoding="utf-8")).get("save_raw"))
    except Exception:                                   # noqa: BLE001
        return False


def _default_root() -> Path:
    r"""첫 실행에서 제시할 저장 위치.

    **프로그램 폴더 옆**이 기본이다 — 눈에 띄고, 프로그램을 어디에 깔든 말이 된다
    (`D:\tools\ortho` → `D:\tools\ortho_data`). 다만 `Program Files` 처럼 쓸 수
    없는 자리면 홈 폴더로 물러난다. 사용자가 첫 화면에서 확인·변경한다.

    독스트링에 Windows 경로가 들어가므로 **raw 문자열**이어야 한다. 안 그러면
    `\t` 가 탭이 되고 `\o` 는 잘못된 이스케이프라 파이썬 3.12+ 에서 경고가 난다.
    """
    sib = PROGRAM_DIR.parent / (PROGRAM_DIR.name + "_data")
    try:
        sib.mkdir(parents=True, exist_ok=True)
        probe = sib / ".write_test"
        probe.write_text("", encoding="utf-8")
        probe.unlink()
        return sib
    except OSError:
        return Path.home() / "ortho-webapp" / "data"


def _saved_root_str() -> str:
    """설정에 적힌 저장 위치 — **있는지는 보지 않는다.**

    외장 드라이브에 환자 폴더를 둔 설치본이 있다. 드라이브를 빼면 그 경로는
    잠깐 사라질 뿐 잘못된 설정이 아니다. "아직 안 골랐다" 와 "골라 뒀는데 지금
    닿지 않는다" 를 가르려면 존재 여부를 빼고 읽는 길이 하나 필요하다.
    """
    try:
        return json.loads(SETTINGS_FILE.read_text(encoding="utf-8")).get("root") or ""
    except Exception:                                   # noqa: BLE001
        return ""


def _saved_root() -> Path | None:
    """사용자가 고른 저장 위치. 없으면 config.yaml 기본값이 쓰인다.

    예전 설치본은 `webapp/settings.json` 에 뒀다 — 있으면 한 번만 옮긴다.
    """
    # 예전 설치본은 홈 폴더에 뒀다 — 있으면 한 번만 옮긴다
    if not SETTINGS_FILE.exists() and LEGACY_SETTINGS.exists():
        try:
            SETTINGS_FILE.parent.mkdir(parents=True, exist_ok=True)
            SETTINGS_FILE.write_text(LEGACY_SETTINGS.read_text(encoding="utf-8"),
                                     encoding="utf-8")
            LEGACY_SETTINGS.rename(LEGACY_SETTINGS.with_suffix(".json.migrated"))
        except Exception:                              # noqa: BLE001
            pass
    try:
        p = json.loads(SETTINGS_FILE.read_text(encoding="utf-8")).get("root")
        return Path(p).expanduser().resolve() if p and Path(p).expanduser().is_dir() else None
    except Exception:
        return None


# 사용자가 고른 값 > config.yaml 이 비었으면 계산된 기본값 > config.yaml
ROOT = _saved_root() or (_default_root() if not cfg.paths.root
                         else Path(cfg.resolve(cfg.paths.root)).resolve())
# 이력 로그는 환자 루트와 함께 둔다 — 기록물이라 같이 백업돼야 한다
LOG_FILE = (Path(cfg.resolve(cfg.paths.log_file)) if cfg.paths.log_file
            else ROOT / "_audit_log.jsonl")
SESS_ROOT = ROOT / "_sessions_tmp"
SESS_ROOT.mkdir(parents=True, exist_ok=True)

# ── 임시 업로드 보관 기한 ─────────────────────────────────────────────────────
# 확정 전 업로드 원본은 _sessions_tmp/<세션id>/ 에만 있다. 브라우저를 그냥 닫거나
# 서버가 재시작되면 이 폴더는 주인 없이 남는다. 정책은 하나로 둔다:
# "확정되지 않은 임시 업로드는 최대 48시간 보관한다."
SESSION_TTL = 48 * 3600      # 마지막 요청 이후 이만큼 지나면 세션을 버린다
SWEEP_INTERVAL = 600         # 청소 주기(초)


# ── 세션 모델 ─────────────────────────────────────────────────────────────────
class Photo:
    def __init__(self, pid, path, w, h):
        self.id = pid
        self.path = Path(path)
        self.w, self.h = w, h
        self.orig_name = None
        self.label = None
        self.confidence = 0.0
        self.probs = {}
        self.slot = None           # 배정된 슬롯(SLOT_*) 또는 'FACE' 또는 None
        # 교합면은 사용자가 상하반전된 화면에서 본다(config.flip_v_slots). 원본
        # 파일과 PPT 안의 이미지는 그대로 두고 표시만 뒤집으므로, editor 값도
        # **뒤집힌 화면 기준**으로 들고 있다. 슬롯이 바뀌면 _put 이 환산한다.
        self.flip_v = False
        self.editor = EditorState()  # 현재 배치(편집기 값, flip_v면 반전 화면 기준)
        self.ref_visit = None      # 정합에 채택된 기준 차수
        self.badge = "ok"          # ok | low | manual | missing
        self.taken_at = None       # EXIF 촬영시각, 서브초 있으면 microsecond까지
        self.exif_seq = None       # EXIF ImageNumber (대개 없다)
        # 이 배치가 어디서 왔는가: model=프레이밍 예측, registration=차수 간 정합,
        # cover=예측을 기각하고 cover-fit, None=모델 없음(종전 동작)
        self.framing = None
        self.framing_note = None   # 기각 사유 등 (검수화면 툴팁용)


class Session:
    def __init__(self, mode, ids, visit):
        self.id = uuid.uuid4().hex[:12]
        self.mode = mode           # 'first' | 'revisit'
        self.ids = ids             # naming.Identifiers
        self.visit = visit         # 이번 차수 알파벳
        self.tmp = SESS_ROOT / self.id
        self.tmp.mkdir(parents=True, exist_ok=True)
        self.touched = time.time()   # 마지막 요청 시각 — 청소 기준

        self.photos: list[Photo] = []
        # 상자 = 순서 있는 목록. 0번이 대표(슬라이드에 들어가고 (1)~(5) 이름을 받음).
        # 나머지는 같은 자리의 추가 촬영본으로 파일만 저장된다.
        self.bins: dict[str, list[str]] = {}   # 'SLOT_*'|'FACE' -> [photo_id]
        # 어느 슬롯을 **어느 사진으로** 정합/프레이밍했나. 검수로 다시 들어올 때
        # 대표가 그대로면 건너뛴다 — 좌·우를 고치고 돌아오면 그 두 칸만 다시 돈다.
        self.framed: dict[str, str] = {}       # 'SLOT_*' -> photo_id
        self.patient_dir: Path | None = None
        self.ppt_path: Path | None = None  # 재진: 기존 PPT 경로
        self.references: dict[str, dict[str, np.ndarray]] = {}  # slot -> {visit: img}
        # 이 세션이 쓸 슬롯 창. 재진이면 기존 PPT의 레이아웃으로 덮어쓴다 —
        # 템플릿 좌표는 이상적이지만 환자 PPT는 여백·간격이 조금씩 다를 수 있고,
        # 새 차수만 템플릿대로 넣으면 그 슬라이드만 어긋난다.
        # 에디터 화면·정합·최종 삽입이 모두 이 창을 함께 쓴다.
        self.slot_windows: dict[str, WindowCm] = dict(SLOT_WINDOWS)
        # 케이스 덱의 얼굴 자리 배정. 자리이름 -> photo_id.
        # 분류기가 정면/45도/측면을 구분하지 못하므로 사람이 직접 고른다.
        self.face_slots: dict[str, str] = {}
        # 얼굴 자리별 편집기 값. 사진이 아니라 **자리**에 매단다 — 같은 사진이
        # 슬라이드 4 좌측과 10·11(분석)에 함께 쓰이는데, 창 크기가 달라 구도를
        # 따로 잡을 수 있어야 하기 때문이다. 없는 자리는 cover-fit.
        self.face_editors: dict[str, EditorState] = {}
        # 자리별 구도 근거: 'model'(프레이밍 예측) | 'cover'(cover-fit)
        self.face_framing: dict[str, str] = {}
        # 사진별 프레이밍 예측 캐시. 예측은 **사진에만** 달렸고 자리(창)에는
        # 달리지 않아서, 자리를 바꿔 배치해도 다시 추론할 필요가 없다.
        self.face_frames: dict[str, object] = {}
        # 차수 노트의 입력 칸. 서식에 끼워 십자뷰 텍스트 박스로 나간다.
        self.note_fields: dict[str, str] = {}
        # 초진 날짜("YY.MM.DD"). 재진이면 기존 PPT에서 읽고, 초진(A)이면 오늘.
        # Rx. Period 의 경과 개월을 여기서 센다.
        self.first_date: str | None = None
        # 치료·장치 시작일("YY.MM.DD"). 초진보다 늦을 수 있어 따로 센다.
        # 이전 차수 슬라이드의 `Tx./App. Period` 괄호에서 읽고, 없으면 None —
        # 임상의가 "이 차수부터 시작" 을 체크하면 그 차수 날짜로 정해진다.
        # 기간별 이력 {"tx": {"dates":[최신순], "last": "직전 줄"}}
        self.period_hist: dict[str, dict] = {}
        # 덱에서 읽은 차수와 날짜 — 기준일 후보로 화면에 낸다
        self.visit_dates: list[dict] = []
        # 직전 차수 슬라이드에 그려져 있던 선 — 검수 판에 그대로 겹쳐 보여준다
        self.prev_lines: list[dict] = []
        # 사람이 고른 기준일. 안 고르면 이력의 최신값이 쓰인다.
        self.period_start: dict[str, str] = {}
        # "이전 차수 값 그대로" — 개월을 다시 안 세고 직전 줄을 그대로 쓴다.
        # 실측에서 App. Period 가 9 month 에 멈춰 있던 경우가 이것이다.
        self.period_keep: dict[str, bool] = {}
        # 박스 통째로 고쳐 쓴 것. 서식 결과를 이긴다.
        self.note_overrides: dict[str, str] = {}
        # 박스별 서식을 이 차수에서만 바꿔 쓴 것 (config.notes.boxes 를 덮는다).
        self.note_templates: dict[str, str] = {}
        # 계측선을 끌어 옮긴 양 (cm). "슬라이드:도형이름" -> [dx, dy].
        # 선의 길이·방향은 그대로 두고 자리만 옮긴다.
        self.line_moves: dict[str, list[float]] = {}

    @property
    def slots(self) -> dict[str, str]:
        """슬롯별 대표 사진. 읽기 전용 — 쓰기는 bins로 한다."""
        return {k: v[0] for k, v in self.bins.items() if k != "FACE" and v}

    @property
    def face(self) -> list[str]:
        return self.bins.get("FACE", [])


def get_session(sid) -> Session:
    # 410은 "있었는데 사라졌다" — 프론트가 경로/사진 없음(404)과 구분해 안내한다.
    s = SESSIONS.get(sid)
    if not s:
        raise HTTPException(410, "세션이 만료되었거나 존재하지 않습니다")
    s.touched = time.time()
    return s


# ── 임시 폴더 청소 ────────────────────────────────────────────────────────────
def discard_session(s: "Session") -> None:
    """세션을 버리고 임시 업로드도 함께 지운다.

    확정을 마친 세션에도 쓴다 — 그때는 원본이 이미 환자 폴더로 복사된 뒤라
    임시본을 남길 이유가 없다.
    """
    SESSIONS.pop(s.id, None)
    shutil.rmtree(s.tmp, ignore_errors=True)


def sweep_sessions(now: float | None = None) -> int:
    """기한 지난 세션과 고아 폴더를 지운다. 반환값은 지운 폴더 수."""
    now = time.time() if now is None else now
    n = 0
    for s in list(SESSIONS.values()):
        if now - s.touched > SESSION_TTL:
            discard_session(s)
            n += 1
    # 서버가 재시작되면 SESSIONS는 비고 폴더만 남는다. 이 고아들도 같은 기한으로
    # 정리하되, 하루이틀은 남겨 두어 "작업하던 게 날아갔다"는 문의에 원본을
    # 수동으로 건져줄 여지를 둔다.
    live = {s.id for s in SESSIONS.values()}
    if not SESS_ROOT.is_dir():
        return n
    for d in SESS_ROOT.iterdir():
        if not d.is_dir() or d.name in live:
            continue
        try:
            if now - d.stat().st_mtime > SESSION_TTL:
                shutil.rmtree(d, ignore_errors=True)
                n += 1
        except OSError:
            pass
    return n


def _sweeper_loop(stop: threading.Event) -> None:
    while not stop.wait(SWEEP_INTERVAL):
        try:
            sweep_sessions()
        except Exception:
            pass          # 청소가 실패해도 서버는 계속 돌아야 한다


# ── 정합 결과 → 배치 환산 ─────────────────────────────────────────────────────
def _win_px_to_cm_affine(win: WindowCm) -> np.ndarray:
    return np.array([[1.0 / PPC, 0.0, win.x], [0.0, 1.0 / PPC, win.y]], np.float32)


def _clamp(st: EditorState, win: WindowCm, bw: float, bh: float) -> EditorState:
    """
    geometry.allow_letterbox가 false일 때만 배율 하한(cover clamp)을 건다.
    true면 사람이 손으로 자른 것과 같이 빈 공간(레터박스)이 생기도록 허용하고,
    그 자리는 PPT에서 geometry.letterbox_color로 칠해진다.
    """
    if cfg.geometry.allow_letterbox:
        return st
    return apply_cover_clamp(st, win, bw, bh)


def registration_to_editor(M_new_to_winpx, win: WindowCm, pw, ph) -> EditorState:
    """정합 유사변환(new_px→창_px) → 편집기 상태(EditorState)."""
    Twin = _win_px_to_cm_affine(win)
    M3 = np.vstack([np.array(M_new_to_winpx, np.float32), [0, 0, 1]])
    T3 = np.vstack([Twin, [0, 0, 1]])
    A = (T3 @ M3)[:2, :]                       # new_px → cm
    pl = placement_from_photo_affine(A.tolist(), pw, ph)
    from coords import placement_to_editor
    bw, bh = cover_base_ext_cm(pw, ph, win)
    st = placement_to_editor(pl, win, bw, bh, PPC)
    return _clamp(st, win, bw, bh)


def framing_to_editor(res: "Fr.FramingResult", win: WindowCm, pw, ph) -> EditorState:
    """프레이밍 예측(raw→canonical) → 편집기 상태.

    모델이 주는 T 는 "raw 픽셀 → 잘라낸 결과물" 이다. 결과물이 슬롯을 꽉 채우도록
    canonical→창픽셀 배율 하나만 앞에 붙이면 정합과 똑같은 모양이 되어, 아래는
    `registration_to_editor` 를 그대로 재사용한다 (문서 §12.1 — 새 좌표 개념이 없다).

    min() + 가운데 정렬로 두는 이유: crop 종횡비(4:3)와 슬롯 종횡비(8.4×6.3=4:3)는
    지금 같지만, 어느 한쪽이 바뀌어도 사진이 잘려나가는 대신 레터박스가 생기게 된다.
    """
    cw, ch = res.canon_wh
    Wpx, Hpx = win.w * PPC, win.h * PPC
    k = min(Wpx / cw, Hpx / ch)
    C_can_to_winpx = np.array([[k, 0.0, (Wpx - k * cw) / 2.0],
                               [0.0, k, (Hpx - k * ch) / 2.0]], np.float64)
    T3 = np.vstack([np.array(res.matrix, np.float64), [0, 0, 1]])
    M = (np.vstack([C_can_to_winpx, [0, 0, 1]]) @ T3)[:2, :]   # raw px → 창 px
    return registration_to_editor(M, win, pw, ph)


# ── 스키마 ────────────────────────────────────────────────────────────────────
class FirstReq(BaseModel):
    name: str
    hospital_id: str = ""      # 병원번호는 요구하지 않는다 — 있으면 형식만 검사
    ortho_id: str


class RevisitReq(BaseModel):
    ppt_path: str


class AssignReq(BaseModel):
    session_id: str
    photo_id: str
    slot: str | None   # SLOT_* | 'FACE' | None(제거)


class AdjustReq(BaseModel):
    session_id: str
    slot: str
    dx: float
    dy: float
    scale: float
    angle: float


# ── 정적/루트 ─────────────────────────────────────────────────────────────────
#
# ### 화면 파일은 **캐시하지 않는다**
#
# `no-cache` 는 '저장하지 마라'가 아니라 '**쓰기 전에 반드시 물어보라**'는 뜻이다.
# 브라우저는 여전히 사본을 들고 있다가 `ETag` 로 확인하고, 안 바뀌었으면 304 를
# 받는다. localhost 라 그 왕복은 사실상 공짜다.
#
# 이걸 안 걸면 브라우저가 **제멋대로 정한 기간**(heuristic freshness) 동안 물어보지도
# 않고 옛 사본을 쓴다. 실제로 그렇게 당했다 — 서버는 새 코드인데 Edge 는 아주 오래된
# `app.js` 를 캐시에서 꺼내 쓰고 있었다. 서버 로그에 `GET /` 은 있는데
# `GET /static/app.js` 가 아예 없어서 알았다.
#
# 이 앱은 `git pull` 로 자기를 갱신한다. 그때마다 새 백엔드 + 옛 화면이라는 조합이
# 생길 수 있고, 그건 있지도 않은 버그처럼 보인다. 갱신되는 프로그램에서 화면 파일
# 캐시는 상시 위험이다.
NO_CACHE = {"Cache-Control": "no-cache"}


@app.middleware("http")
async def _no_store_api(request, call_next):
    """API 응답은 **저장하지 않는다.**

    화면 파일과 달리 API 는 조회할 때마다 답이 달라진다 — 방금 확정한 차수, 방금
    생긴 파일이 그 답이다. 캐시 지시가 없으면 브라우저가 제 나름의 기준으로 옛
    응답을 다시 쓸 수 있고, 그러면 사람이 새로고침해야 저장 결과가 보인다.
    """
    resp = await call_next(request)
    if request.url.path.startswith("/api/"):
        resp.headers["Cache-Control"] = "no-store"
    return resp


@app.get("/")
def index():
    return FileResponse(FRONTEND_DIR / "index.html", headers=NO_CACHE)


class _NoCacheStatic(StaticFiles):
    """`/static/*` 도 매번 확인시킨다. 바뀐 게 없으면 304 로 끝난다."""

    def file_response(self, *args, **kwargs):
        r = super().file_response(*args, **kwargs)
        r.headers.setdefault("Cache-Control", "no-cache")
        return r


@app.get("/api/health")
def health():
    fr = {"loaded": False}
    if framer is not None:
        fr = {"loaded": True, "tag": framer.meta.get("tag"),
              "input": [framer.iw, framer.ih],
              "classes": sorted(framer.files),
              "models_per_class": {c: len(v) for c, v in framer.files.items()},
              # 임시 대역이면 여기 사유가 실린다 — 배포 전 교체하라는 신호.
              "placeholder": framer.placeholder}
    return {"ok": True, "classifier": type(classifier).__name__,
            # 첫 실행이면 화면이 저장 위치를 묻는다. **파일 존재가 아니라 `root`
            # 가 적혔는지**로 본다 — 다른 설정(개월 표기 등)이 먼저 저장되면
            # 파일은 생기지만 저장 위치는 여전히 안 고른 상태다.
            # 아직 안 고른 경우에만 첫 실행 화면을 띄운다. 골라 뒀는데 지금
            # 닿지 않는 것(외장 드라이브 분리)은 아래 root_missing 으로 구분한다 —
            # 첫 실행처럼 물으면 사용자가 임시 위치를 확정해 버리고, 그 순간
            # 설정의 외장 경로가 덮여 사라진다.
            "needs_setup": not _saved_root() and not _saved_root_str(),
            "root_missing": ("" if _saved_root() else _saved_root_str()),
            "root": str(ROOT), "program_dir": str(PROGRAM_DIR),
            "framing": fr,
            "classes": cfg.classes, "slots": cfg.ppt.slot_names,
            "px_per_cm": PPC,
            "rotation_range_deg": cfg.geometry.rotation_range_deg,
            "windows": {k: {"x": v.x, "y": v.y, "w": v.w, "h": v.h}
                        for k, v in SLOT_WINDOWS.items()}}


# ── 파일 브라우저 (재진: 환자 PPT를 경로로 선택) ──────────────────────────────
# 폴더 선택 창. 브라우저는 폴더를 고르게는 해도 **절대경로를 안 알려준다**
# (`webkitdirectory`·`showDirectoryPicker` 둘 다 보안상 이름만 준다). 서버가 그
# 폴더에 환자 자료를 써야 하므로 경로가 필요하고, 서버가 사용자와 같은 PC 에서
# 도니까 직접 띄운다.
# **맨 앞으로 띄워야 한다.** 서버가 백그라운드 프로세스라 그냥 열면 창이 다른 창
# 뒤에 가려지고, 사용자에게는 "눌러도 반응 없음" 으로 보인다. TopMost 인 빈 폼을
# 주인으로 넘겨 앞으로 끌어낸다.
_PS_PICK = (
    # 고른 경로를 **UTF-8 로** 내보내게 못 박는다. 안 그러면 PowerShell 은 콘솔
    # 출력 코드페이지(한국어 Windows 는 CP949)로 쓴다 — `C:\사용자\바탕화면` 처럼
    # 한글이 든 경로가 깨져 돌아오고, 없는 폴더를 저장 위치로 잡게 된다.
    "[Console]::OutputEncoding = [Text.Encoding]::UTF8; "
    "Add-Type -AssemblyName System.Windows.Forms; "
    "$o = New-Object System.Windows.Forms.Form; "
    "$o.TopMost = $true; $o.ShowInTaskbar = $false; "
    "$o.StartPosition = 'CenterScreen'; $o.Size = '1,1'; $o.Show(); "
    "$d = New-Object System.Windows.Forms.FolderBrowserDialog; "
    "$d.Description = '환자 자료를 저장할 폴더'; "
    "$d.ShowNewFolderButton = $true; "
    "if ($env:ACF_START) { $d.SelectedPath = $env:ACF_START }; "
    "$r = $d.ShowDialog($o); $o.Close(); "
    "if ($r -eq 'OK') { [Console]::Out.Write($d.SelectedPath) }"
)

# 창이 없는 리눅스·맥용. Windows 쪽 경로가 안 통하는 환경에서만 쓴다.
_TK_PICK = r"""
import sys, tkinter as tk
from tkinter import filedialog
r = tk.Tk(); r.withdraw(); r.attributes('-topmost', True)
print(filedialog.askdirectory(title='환자 자료를 저장할 폴더',
                              initialdir=sys.argv[1] or None) or '')
"""


def _powershell() -> str | None:
    """Windows 의 powershell. WSL 에서도 interop 으로 그대로 부를 수 있다.

    WSL 에서 tkinter 를 쓰면 창이 WSLg 쪽에 떠서, 브라우저를 Windows 에서 보는
    사용자에게는 보이지 않는다 — 눌러도 반응이 없는 것처럼 된다. `powershell.exe`
    를 부르면 **진짜 Windows 탐색기 창**이 뜨고 경로도 Windows 경로로 돌아온다.
    """
    if os.name == "nt":
        return "powershell"
    return shutil.which("powershell.exe")


def _to_local(path: str) -> str:
    """Windows 경로 → 이 프로세스가 쓸 수 있는 경로 (WSL 이면 /mnt/c/...)."""
    if os.name == "nt" or not path:
        return path
    try:
        r = subprocess.run(["wslpath", "-u", path], capture_output=True,
                           encoding="utf-8", errors="replace", timeout=10)
        if r.returncode == 0 and r.stdout.strip():
            return r.stdout.strip()
    except (OSError, subprocess.SubprocessError):
        pass
    return path


@app.get("/api/pick-folder")
def pick_folder(start: str = ""):
    """운영체제의 **폴더 선택 창**을 띄우고 고른 경로를 돌려준다.

    Windows·WSL 은 `powershell.exe`(탐색기 창), 그 밖은 tkinter 를 쓴다. 창을
    못 띄우면 실패를 돌려주고 화면이 앱 안 폴더 트리로 물러난다.
    """
    ps = _powershell()
    try:
        if ps:
            env = {**os.environ}
            if start:
                w = subprocess.run(["wslpath", "-w", start], capture_output=True,
                                   encoding="utf-8", errors="replace",
                                   timeout=10) if os.name != "nt" else None
                env["ACF_START"] = (w.stdout.strip() if w and w.returncode == 0
                                    else start)
            r = subprocess.run([ps, "-NoProfile", "-STA", "-Command", _PS_PICK],
                               capture_output=True, encoding="utf-8",
                               errors="replace", timeout=600, env=env)
            path = (r.stdout or "").strip().replace("\r", "").strip()
            if r.returncode != 0 and not path:
                return {"ok": False, "detail": (r.stderr or "").strip()[-200:]
                        or "창을 띄우지 못했습니다"}
            if not path:
                return {"ok": False, "cancelled": True}
            return {"ok": True, "path": _to_local(path), "shown": path}
        r = subprocess.run([sys.executable, "-c", _TK_PICK, start],
                           capture_output=True, encoding="utf-8",
                           errors="replace", timeout=600,
                           env={**os.environ, "PYTHONIOENCODING": "utf-8"})
        if r.returncode != 0:
            return {"ok": False, "detail": (r.stderr or "").strip()[-200:]
                    or "창을 띄우지 못했습니다"}
        path = (r.stdout or "").strip()
        return ({"ok": True, "path": path, "shown": path} if path
                else {"ok": False, "cancelled": True})
    except Exception as e:                                        # noqa: BLE001
        return {"ok": False, "detail": f"{type(e).__name__}"}


@app.get("/api/browse")
def browse(path: str = ""):
    root = ROOT
    base = (root / path).resolve() if path else root
    if not str(base).startswith(str(root.resolve())):
        raise HTTPException(400, "루트 밖 접근 불가")
    if not base.exists():
        raise HTTPException(404, "경로 없음")
    dirs, ppts = [], []
    for p in sorted(base.iterdir()):
        rel = str(p.relative_to(root))
        if p.is_dir() and not p.name.startswith("_"):
            dirs.append({"name": p.name, "path": rel})
        elif p.suffix.lower() == ".pptx":
            ppts.append({"name": p.name, "path": rel})
    parent = str(base.parent.relative_to(root)) if base != root else None
    return {"cwd": str(base.relative_to(root)) or ".", "parent": parent,
            "dirs": dirs, "ppts": ppts}


# ── 세션 생성: 초진 ───────────────────────────────────────────────────────────
@app.post("/api/session/first")
def session_first(req: FirstReq):
    try:
        ids = N.validate_identifiers(
            req.name, req.hospital_id, req.ortho_id,
            hospital_digits=cfg.identifiers.hospital_id.digits,
            ortho_digits=cfg.identifiers.ortho_id.digits,
            name_regex=cfg.identifiers.name.allow_regex,
            require_hospital=False)
    except N.NamingError as e:
        raise HTTPException(400, str(e))
    s = Session("first", ids, "A")
    folder = _gen_folder_name(ids)
    s.patient_dir = ROOT / folder
    SESSIONS[s.id] = s
    return {"session_id": s.id, "visit": "A", "folder": folder,
            "folder_exists": s.patient_dir.exists()}


# ── 세션 생성: 재진 ───────────────────────────────────────────────────────────
@app.post("/api/session/revisit")
def session_revisit(req: RevisitReq):
    root = ROOT
    ppt_path = (root / req.ppt_path).resolve()
    if not str(ppt_path).startswith(str(root.resolve())) or not ppt_path.exists():
        raise HTTPException(400, "PPT 경로 오류")
    try:
        ids = _parse_ppt_name(ppt_path.name)
    except N.NamingError as e:
        raise HTTPException(400, str(e))

    patient_dir = ppt_path.parent
    s = Session("revisit", ids, "A")
    s.patient_dir = patient_dir
    s.ppt_path = ppt_path

    # 기준영상 복원
    try:
        prs = T.load_presentation(ppt_path)
    except PermissionError:
        raise HTTPException(409, "PPT가 다른 프로그램(PowerPoint)에서 열려 "
                                 "있습니다 — 닫은 뒤 다시 시도해 주세요")
    visits = Rd.read_all_visits(prs, cfg, PPC)
    s.references = Rd.references_for_registration(visits)
    # 차수 이력의 진실은 PPT 뿐이다 — 사진 파일명은 보지 않는다.
    letters = sorted({vs.visit for vs in visits if vs.visit}, key=N.letter_to_num)
    s.visit = N.next_visit_letter(letters)
    SESSIONS[s.id] = s
    return {"session_id": s.id, "visit": s.visit, "prev_visits": letters,
            "ids": {"name": ids.name, "hospital_id": ids.hospital_id, "ortho_id": ids.ortho_id},
            "patient_dir": str(patient_dir)}


# ── 환자 목록 / 세션 열기 (Setup) ─────────────────────────────────────────────
# 루트는 평평하다고 가정한다: 환자 폴더가 루트 바로 아래 나열된다.
# 초진/재진은 사용자가 선언하지 않는다 — 폴더 안에 PPT가 있느냐로 서버가 판정한다.
IMG_EXT = (".jpg", ".jpeg", ".png")


# 환자 목록이 뜰 때마다 모든 PPT 를 여는 건 느리다 — (경로, mtime)로 캐시.
_VISITS_CACHE: dict[str, tuple[float, dict]] = {}


_EMPTY_SCAN = {"visits": [], "excluded": [], "fallback": False, "slides": 0}


def _ppt_visit_letters(path: Path) -> dict:
    """PPT 라벨 스캔 결과(차수 장부) — 잠금·손상이면 빈 장부."""
    try:
        mt = path.stat().st_mtime
    except OSError:
        return _EMPTY_SCAN
    hit = _VISITS_CACHE.get(str(path))
    if hit and hit[0] == mt:
        return hit[1]
    try:
        out = Rd.scan_ppt_visits(T.load_presentation(path), cfg)
    except Exception:                                   # noqa: BLE001 — 잠금·손상
        out = _EMPTY_SCAN
    if len(_VISITS_CACHE) > 256:
        _VISITS_CACHE.clear()
    _VISITS_CACHE[str(path)] = (mt, out)
    return out


def _ppt_reject_reason(name: str, ids) -> dict:
    """이 파일을 왜 그 환자의 PPT 로 안 봤는가 — 화면에 보일 한 줄 + 해결 여부.

    구형 `.ppt` 는 **같은 이름을 .pptx 로 저장했을 때 인식되는지**까지 본다.
    이름 자체가 양식과 어긋나 있으면 변환만으로는 안 되는데, 그걸 알려주지 않으면
    시키는 대로 저장하고도 여전히 안 붙는 이유를 알 수 없다.
    """
    if name.startswith("~$"):
        return {"why": "PowerPoint 임시 파일입니다"}
    if name.lower().endswith(".ppt"):
        try:
            got = _parse_ppt_name(name + "x")
            ok = got.ortho_id == ids.ortho_id
        except N.NamingError:
            ok = False
        return {"why": (".pptx 로 다시 저장하면 인식됩니다" if ok
                        else ".pptx 로 저장해도 등록된 이름 양식과 맞지 않습니다"),
                "convertible": ok}
    try:
        got = _parse_ppt_name(name)
    except N.NamingError:
        return {"why": "등록된 PPT 이름 양식과 맞지 않습니다"}
    if got.ortho_id != ids.ortho_id:
        return {"why": f"교정번호가 다릅니다 ({got.ortho_id})"}
    return {"why": "인식 가능"}


def _scan_patient(d: Path) -> dict | None:
    """환자 폴더 하나 → 목록 한 줄. 폴더명이 명명 규칙에 안 맞으면 None.

    차수 이력·날짜의 출처는 **PPT 라벨뿐**이다 — 사진 파일은 읽지 않는다.
    사진이 어떤 이름·폴더 구조로 저장돼 있는지 이 앱이 알 필요가 없다.
    """
    try:
        ids = _parse_folder(
            d.name,
            hospital_digits=cfg.identifiers.hospital_id.digits,
            ortho_digits=cfg.identifiers.ortho_id.digits,
            name_regex=cfg.identifiers.name.allow_regex, label="폴더명")
    except N.NamingError:
        return None
    entries = _patient_files(d)      # 하위 폴더 정리 환자 — 두 단계까지
    files = [e.name for e in entries]
    picked = _find_ppt(entries, d, ids)
    ppt_name = picked.relative_to(d).as_posix() if picked else None
    scan = _EMPTY_SCAN
    if picked is not None:
        scan = _ppt_visit_letters(picked)
    visits = sorted({v["visit"] for v in scan["visits"]}, key=N.letter_to_num)

    # 차수별 날짜 — 라벨의 "26.08.12" 를 목록 표기(2026.08.12)로 편다.
    pdates = {v["visit"]: v["date"] for v in scan["visits"] if v["date"]}

    def date_of(letter: str) -> str | None:
        if letter in pdates:
            nums = re.findall(r"\d+", pdates[letter])
            if len(nums) >= 3:
                return f"20{int(nums[0]):02d}.{int(nums[1]):02d}.{int(nums[2]):02d}"
        return None

    return {
        "folder": d.name, "name": ids.name,
        "hospital_id": ids.hospital_id, "ortho_id": ids.ortho_id,
        "visits": visits, "next_visit": N.next_visit_letter(visits),
        "visit_dates": {L: date_of(L) for L in visits},
        "ppt": ppt_name,
        # 확인 줄 재료 — 인식된 차수(슬라이드 번호 포함), 제외된 라벨 장, 폴백 여부.
        "visit_slides": scan["visits"],
        "label_excluded": scan["excluded"],
        "label_fallback": scan["fallback"],
        "ppt_slides": scan["slides"],
        # 새 슬라이드를 **어느 장 뒤에** 넣을지. 날짜가 가장 늦은 차수 슬라이드다 —
        # "몇 번째가 된다" 보다 "몇 번 뒤" 가 사람이 슬라이드를 세는 방식이다.
        "suggest_after": (max(scan["visits"],
                              key=lambda v: (N.letter_to_num(v["visit"]),
                                             v["slide_no"]))["slide_no"]
                          if scan["visits"] else None),
        # PPT 를 못 찾았을 때: 폴더의 PPT 파일마다 못 알아본 이유 한 줄.
        # 왜 안 붙는지 아무도 모르는 채 새 PPT 가 생기는 것을 막는다.
        "ppt_diag": ([] if ppt_name else
                     [{"name": f, **_ppt_reject_reason(f, ids)}
                      for f in files
                      if f.lower().endswith((".pptx", ".ppt"))][:5]),
        "updated": datetime.fromtimestamp(d.stat().st_mtime).strftime("%Y-%m-%d"),
    }


@app.get("/api/patients")
def patients():
    """Setup 화면의 환자 목록. 규칙에 안 맞는 폴더는 세지만 싣지 않는다."""
    root = ROOT
    root.mkdir(parents=True, exist_ok=True)
    found, skipped = [], []
    for d in sorted(root.iterdir()):
        if not d.is_dir() or d.name.startswith("_"):
            continue
        rec = _scan_patient(d)
        found.append(rec) if rec else skipped.append(d.name)
    return {"patients": found, "skipped": skipped,
            "root": str(root.resolve()),
            "rules": {"hospital_digits": cfg.identifiers.hospital_id.digits,
                      "ortho_digits": cfg.identifiers.ortho_id.digits,
                      # 새 환자 미리보기용 — 실제로 만들어질 모습. 인식 전용
                      # 블록(순번 등)과 병원번호는 생성 때 없는 셈 친다.
                      "folder_pattern": N.strip_recognition(
                          _folder_pattern(), {"hospital_id"}),
                      "photo_pattern": cfg.naming.photo_pattern,
                      "slots": len(cfg.ppt.slot_names)}}


# ── 시작 화면 미리보기 — PPT에 실린 그림 ────────────────────────────────────
# 완성본 JPG 의 (1)~(n) 번호에는 기대지 않는다 — 번호 규칙이 다른 옛 폴더에서
# 미리보기 부위가 어긋난다. 슬라이드에 실제로 실린 그림을 복원해 주면 번호와
# 무관하게 항상 맞다. PPT 파싱은 싸지 않으므로 (경로, mtime)로 캐시한다.
_PV_CACHE: dict[str, tuple[float, dict]] = {}


def _pv_jpeg(bgr: np.ndarray) -> str:
    """미리보기용 축소 JPEG 데이터 URI. 원본 해상도를 내보낼 이유가 없다."""
    h, w = bgr.shape[:2]
    if max(h, w) > 360:
        sc = 360 / max(h, w)
        bgr = cv2.resize(bgr, (max(1, int(w * sc)), max(1, int(h * sc))),
                         interpolation=cv2.INTER_AREA)
    ok, buf = cv2.imencode(".jpg", bgr, [cv2.IMWRITE_JPEG_QUALITY, 82])
    if not ok:
        raise ValueError("JPEG 인코딩 실패")
    return "data:image/jpeg;base64," + base64.b64encode(buf.tobytes()).decode()


@app.get("/api/ppt_preview")
def ppt_preview(folder: str):
    """첫 차수 십자 5장 + 얼굴 2장을 PPT에서 복원해 데이터 URI로 준다."""
    root = ROOT.resolve()
    d = (root / folder).resolve()
    if not str(d).startswith(str(root)) or not d.is_dir():
        raise HTTPException(404, "환자 폴더가 없습니다")
    try:
        ids = _parse_folder(d.name,
                            hospital_digits=cfg.identifiers.hospital_id.digits,
                            ortho_digits=cfg.identifiers.ortho_id.digits,
                            name_regex=cfg.identifiers.name.allow_regex,
                            label="폴더명")
    except N.NamingError:
        raise HTTPException(404, "폴더명이 규칙에 맞지 않습니다")
    path = _find_ppt(_patient_files(d), d, ids)
    if path is None:
        return {"slots": {}, "faces": []}
    key, mt = str(path), path.stat().st_mtime
    hit = _PV_CACHE.get(key)
    if hit and hit[0] == mt:
        return hit[1]
    try:
        prs = T.load_presentation(path)
    except OSError:                     # PowerPoint 잠금 등 — 미리보기만 조용히 생략
        return {"slots": {}, "faces": []}
    slots: dict[str, str] = {}
    pv_ctr = (emu_to_cm(prs.slide_width) / 2, emu_to_cm(prs.slide_height) / 2)
    for i, slide in enumerate(prs.slides):
        vs = Rd.read_visit_slide(slide, i, cfg, PPC, slide_ctr=pv_ctr)
        if vs.slots:                # 덱은 시간순 — 첫 십자 슬라이드가 초진
            slots = {k: _pv_jpeg(ref.image) for k, ref in vs.slots.items()}
            break
    faces: list[str] = []
    from PIL import Image as _Im                                # noqa: PLC0415

    def _face_uri(im) -> str:
        im.thumbnail((360, 360))
        out = io.BytesIO()
        im.convert("RGB").save(out, "JPEG", quality=82)
        return ("data:image/jpeg;base64,"
                + base64.b64encode(out.getvalue()).decode())

    def _pic_image(sh):
        try:
            return getattr(sh, "image", None)
        except Exception:           # noqa: BLE001 — 외부 참조 그림 등
            return None

    # 앱이 만든 덱 — 커밋이 붙인 이름(PHOTO_FACE_*)이 정확하다. 파생 자리
    # (big_slides)에 같은 사진이 또 실리므로 sha1로 거른다.
    seen: set[str] = set()
    for slide in prs.slides:
        if len(faces) >= 2:
            break
        for sh in slide.shapes:
            if len(faces) >= 2:
                break
            if not str(getattr(sh, "name", "")).startswith("PHOTO_FACE_"):
                continue
            img = _pic_image(sh)
            if img is None or img.sha1 in seen:
                continue
            seen.add(img.sha1)
            try:
                with _Im.open(io.BytesIO(img.blob)) as im:
                    faces.append(_face_uri(im))
            except Exception:       # noqa: BLE001 — 깨진 그림
                continue
    # 수제 PPT 폴백 — 이름 규약이 없다. 앞에서부터 훑어 **너비 12cm 이상 사진이
    # 2장 이상**인 첫 슬라이드를 얼굴 슬라이드로 보고 그 슬라이드의 사진만 쓴다
    # (십자 사진은 8cm 안팎이라 안 걸리고, 큰 사진이 한 장뿐인 슬라이드는
    # 엑스레이 등일 수 있다). 뒤 슬라이드는 보지 않는다.
    if not faces:
        for slide in prs.slides:
            big = [sh for sh in slide.shapes
                   if _pic_image(sh) is not None
                   and emu_to_cm(sh.width) >= 12.0]
            if len(big) < 2:
                continue
            for sh in big[:2]:
                try:
                    with _Im.open(io.BytesIO(_pic_image(sh).blob)) as im:
                        faces.append(_face_uri(im))
                except Exception:   # noqa: BLE001
                    continue
            break
    res = {"slots": slots, "faces": faces}
    if len(_PV_CACHE) > 64:         # 시작 화면에서 훑는 규모면 충분한 상한
        _PV_CACHE.clear()
    _PV_CACHE[key] = (mt, res)
    return res


# ── 저장 위치(루트) 고르기 ────────────────────────────────────────────────────
# 브라우저는 서버의 절대 경로를 알려줄 수 없으므로 탐색을 서버가 대신한다.
# 로컬호스트 전용 오프라인 앱이라 서버 파일시스템을 훑는 것이 허용된다.
def _fs_roots() -> list[str]:
    if os.name == "nt":
        return [f"{d}:\\" for d in string.ascii_uppercase if Path(f"{d}:\\").exists()]
    return ["/"]


@app.get("/api/fs")
def fs(path: str = ""):
    base = Path(path).expanduser() if path else ROOT
    if not base.is_dir():
        raise HTTPException(404, f"폴더가 없습니다: {base}")
    base = base.resolve()
    try:
        dirs = [{"name": p.name, "path": str(p)}
                for p in sorted(base.iterdir(), key=lambda x: x.name.lower())
                if p.is_dir() and not p.name.startswith(".")]
    except PermissionError:
        raise HTTPException(403, f"접근 권한이 없습니다: {base}")
    return {"path": str(base),
            "parent": str(base.parent) if base.parent != base else None,
            "dirs": dirs, "drives": _fs_roots(), "current_root": str(ROOT)}


class MkdirReq(BaseModel):
    path: str
    name: str


@app.post("/api/fs/mkdir")
def fs_mkdir(req: MkdirReq):
    """저장 위치 픽커에서 새 폴더 만들기 — 탐색기로 나갔다 올 필요가 없게."""
    base = Path(req.path).expanduser()
    if not base.is_dir():
        raise HTTPException(404, f"폴더가 없습니다: {base}")
    name = req.name.strip().rstrip(".")
    if not name or set(chr(92) + '/:*?"<>|') & set(name):
        raise HTTPException(400, '폴더 이름이 비었거나 \\ / : * ? " < > | 가 들어 있습니다')
    p = (base / name).resolve()
    if p.exists():
        if p.is_dir():
            return {"path": str(p), "existed": True}   # 이미 있으면 그리로 들어간다
        raise HTTPException(400, f"같은 이름의 파일이 있습니다: {name}")
    try:
        p.mkdir()
    except OSError as e:
        raise HTTPException(400, f"만들지 못했습니다: {e}")
    return {"path": str(p), "existed": False}


@app.post("/api/root/recheck")
def root_recheck():
    """사라졌던 저장 위치가 돌아왔는지 다시 본다 — 드라이브를 꽂고 누른다.

    설정은 건드리지 않는다. 경로가 살아 있으면 그리로 되돌아가고, 아니면 어디에
    무엇이 없는지 그대로 알려준다.
    """
    global ROOT, SESS_ROOT
    saved = _saved_root()
    if saved is None:
        return {"ok": False, "path": _saved_root_str(), "root": str(ROOT)}
    for s in list(SESSIONS.values()):     # 임시 업로드는 옛 루트에 있다
        discard_session(s)
    ROOT = saved
    SESS_ROOT = ROOT / "_sessions_tmp"
    SESS_ROOT.mkdir(parents=True, exist_ok=True)
    _VISITS_CACHE.clear()
    _PV_CACHE.clear()
    return {"ok": True, "root": str(ROOT)}


class RootReq(BaseModel):
    path: str


@app.post("/api/root")
def set_root(req: RootReq):
    """저장 위치 변경. 열려 있던 세션은 옛 경로를 가리키므로 함께 버린다."""
    global ROOT, SESS_ROOT
    p = Path(req.path).expanduser()
    if not p.is_dir():
        # 첫 실행에서는 아직 없는 폴더를 고를 수 있어야 한다. 다만 **부모는 있어야**
        # 만든다 — 드라이브 문자나 상위 폴더 오타로 엉뚱한 트리가 생기는 걸 막는다.
        if not p.parent.is_dir():
            raise HTTPException(400, f"상위 폴더가 없습니다: {p.parent}")
        try:
            p.mkdir(parents=False)
        except OSError as e:
            raise HTTPException(400, f"폴더를 만들 수 없습니다: {p} ({e.strerror})")
    # 열려 있던 세션은 옛 루트를 가리키므로 임시 업로드까지 함께 버린다.
    # 루트를 바꾼 뒤에는 옛 _sessions_tmp가 청소 대상 밖으로 나가 영영 남는다.
    for s in list(SESSIONS.values()):
        discard_session(s)
    ROOT = p.resolve()
    SESS_ROOT = ROOT / "_sessions_tmp"
    SESS_ROOT.mkdir(parents=True, exist_ok=True)
    # **있던 설정을 지우지 않는다.** 예전에는 {"root": ...} 하나로 통째 덮어써서,
    # 저장 위치를 한 번 바꾸면 이름 양식·개월 표기·여백 색·복사 설정·PPT 기억이
    # 전부 사라졌다. 고친 항목만 얹는다 (prefs_set 과 같은 방식).
    try:
        d = json.loads(SETTINGS_FILE.read_text(encoding="utf-8"))
        if not isinstance(d, dict):
            d = {}
    except Exception:                                   # noqa: BLE001
        d = {}
    d["root"] = str(ROOT)
    SETTINGS_FILE.parent.mkdir(parents=True, exist_ok=True)
    SETTINGS_FILE.write_text(json.dumps(d, ensure_ascii=False, indent=2),
                             encoding="utf-8")
    return {"root": str(ROOT)}


@app.get("/api/folder")
def folder_contents(folder: str):
    """환자 폴더 안을 그대로 보여준다 — 무엇이 이미 있는지 보고 사진을 넣게.

    파일 이름·종류·크기만 나열한다. 차수 이력은 PPT 라벨이 진실이라 여기서
    사진 파일명을 해석하지 않는다.
    """
    d = (ROOT / folder).resolve()
    if not str(d).startswith(str(ROOT)) or not d.is_dir():
        raise HTTPException(404, f"폴더를 찾을 수 없습니다: {folder}")
    try:
        ids = _parse_folder(
            d.name,
            hospital_digits=cfg.identifiers.hospital_id.digits,
            ortho_digits=cfg.identifiers.ortho_id.digits,
            name_regex=cfg.identifiers.name.allow_regex, label="폴더명")
    except N.NamingError:
        ids = None

    entries = _patient_files(d)
    items = [{
        "name": f.relative_to(d).as_posix(), "size": f.stat().st_size,
        "kind": ("ppt" if f.suffix.lower() == ".pptx"
                 else "photo" if f.suffix.lower() in IMG_EXT else "other"),
    } for f in entries]
    # 이 환자의 PPT 로 자동 선택된 파일 — 화면이 "선택됨"으로 표시한다.
    picked = _find_ppt(entries, d, ids) if ids else None
    sel = picked.relative_to(d).as_posix() if picked else None
    if picked is not None:
        for it, e in zip(items, entries):
            if e == picked:
                it["selected"] = True
                break
    return {"folder": d.name, "path": str(d), "items": items, "ppt": sel}


class PptPickReq(BaseModel):
    folder: str
    ppt: str            # 환자 폴더 기준 상대경로


@app.post("/api/folder/ppt")
def folder_pick_ppt(req: PptPickReq):
    """이어붙일 덱을 사람이 고른다.

    고르는 규칙(바로 아래 우선 · 마지막에 쓴 것 기억)이 대부분 맞지만, 한 번
    어긋나면 되돌릴 길이 있어야 한다 — 잘못 기억된 채로 두면 이후 모든 차수가
    엉뚱한 덱으로 가고, 한 환자의 기록이 둘로 갈린다.
    """
    root = ROOT.resolve()
    d = (root / req.folder).resolve()
    if not str(d).startswith(str(root)) or not d.is_dir():
        raise HTTPException(404, f"환자 폴더를 찾을 수 없습니다: {req.folder}")
    p = (d / req.ppt).resolve()
    if d not in p.parents or p.suffix.lower() != ".pptx" or not p.is_file():
        raise HTTPException(400, "그 폴더 안의 .pptx 파일이어야 합니다")
    rel = p.relative_to(d).as_posix()
    _remember_ppt(d.name, rel)
    return {"ok": True, "ppt": rel}


class OpenReq(BaseModel):
    folder: str | None = None          # 기존 환자: 폴더명
    name: str | None = None            # 새 환자: 식별자 3종
    hospital_id: str | None = None
    ortho_id: str | None = None
    # 확인 줄에서 고친 값 — 이번 차수 글자, 새 슬라이드를 넣을 자리
    visit: str | None = None
    insert_after: int | None = None   # 이 번호의 슬라이드 **뒤**에 (0 = 맨 앞)


@app.post("/api/session")
def session_open(req: OpenReq):
    """
    세션 시작 (초진/재진 통합).
      - folder 지정   → 기존 환자 폴더를 연다
      - 식별자 지정   → 새 환자. 폴더가 이미 있으면 409로 막고 목록에서 고르게 한다
    차수와 모드는 폴더 내용에서 서버가 판정한다.
    """
    root = ROOT
    dig = dict(hospital_digits=cfg.identifiers.hospital_id.digits,
               ortho_digits=cfg.identifiers.ortho_id.digits,
               name_regex=cfg.identifiers.name.allow_regex)

    if req.folder:
        d = (root / req.folder).resolve()
        if not str(d).startswith(str(root.resolve())) or not d.is_dir():
            raise HTTPException(404, f"환자 폴더를 찾을 수 없습니다: {req.folder}")
        try:
            ids = _parse_folder(d.name, label="폴더명", **dig)
        except N.NamingError as e:
            raise HTTPException(400, str(e))
    else:
        try:
            # 병원번호는 요구하지 않는다 — 값이 있으면 형식만 검사하고, 비어
            # 있으면 이름을 만들 때 그 블록을 뺀다. 나중에 검수 화면의
            # 환자정보 칸으로 채울 수 있다.
            ids = N.validate_identifiers(req.name, req.hospital_id, req.ortho_id,
                                         require_hospital=False, **dig)
        except N.NamingError as e:
            raise HTTPException(400, str(e))
        d = root / _gen_folder_name(ids)
        if d.exists():
            return JSONResponse(status_code=409, content={
                "error": "patient_exists", "folder": d.name,
                "detail": f"이미 등록된 환자입니다: {d.name} — 목록에서 선택해 이어서 진행하세요"})

    visits: list[str] = []            # 차수 이력의 진실은 PPT 뿐이다
    ppt_path = d / _gen_ppt_name(ids)
    if d.is_dir():
        # 옛 형식 이름이거나 하위 폴더에 있어도 찾아낸다. 기억해 둔 파일이 있으면
        # 생성 이름보다 그쪽이 이긴다 — 지난번에 이어붙인 덱으로 계속 가야 한다.
        hit = _find_ppt(_patient_files(d), d, ids)
        if hit is not None:
            ppt_path = hit
    has_ppt = ppt_path.exists()

    # mode = PPT를 어떻게 만들 것인가. 기존 PPT가 있어야만 'revisit'(이어붙이기).
    s = Session("revisit" if has_ppt else "first", ids, N.next_visit_letter(visits))
    s.patient_dir = d
    if has_ppt:
        s.ppt_path = ppt_path
        try:
            prs = T.load_presentation(ppt_path)
        except PermissionError:
            raise HTTPException(409, "PPT가 다른 프로그램(PowerPoint)에서 열려 "
                                     "있습니다 — 닫은 뒤 다시 시도해 주세요")
        s.slot_windows = _layout_from_ppt(prs)
        # 기준영상도 그 PPT의 창으로 복원해야 정합이 같은 좌표계에서 이뤄진다.
        seen = Rd.read_all_visits(prs, cfg, PPC, s.slot_windows)
        s.references = Rd.references_for_registration(seen)
        if not s.references:
            # PPT 는 찾았는데 기준영상이 0 — 라벨(차수)을 못 읽었거나 십자뷰가
            # 인식 조건에 안 맞는 것이다. 정합이 왜 안 됐는지 나중에 추적하게 남긴다.
            S.append_audit(LOG_FILE, {
                "event": "references_empty", "patient": d.name,
                "ppt": ppt_path.name, "slides": len(prs.slides._sldIdLst),
                "cross_slides": len(seen)})
        # 차수 장부는 가벼운 스캔(십자뷰 단일 기준 + 폴백)과 같은 규칙으로 센다 —
        # 목록 화면과 세션이 다른 차수를 말하면 확인 줄이 거짓말이 된다.
        scan = Rd.scan_ppt_visits(prs, cfg)
        ppt_letters = [v["visit"] for v in scan["visits"]]
        if ppt_letters:
            visits = sorted({*visits, *ppt_letters}, key=N.letter_to_num)
            s.visit = N.next_visit_letter(visits)
        # 재진 라벨은 이 덱의 표기(마지막 십자뷰 라벨의 지문)를 따른다
        with_label = [vs for vs in seen if getattr(vs, "label_text", "")]
        s.label_fp = (Rd.label_fingerprint(with_label[-1].label_text)
                      if with_label else None)
        s.first_date = _first_visit_date(seen)
        try:
            # 수제 라벨 위치·폰트 상속 — 있으면 좋은 것이라, 어떤 실패도
            # 세션 열기를 막으면 안 된다 (스타일 없이 기본 규칙으로 간다).
            s.label_style = Rd.last_label_style(prs, cfg)
            s.status_style = Rd.last_status_style(prs, cfg)
            # 통째 복사 상속용 원본 상자 XML — 있으면 속성 상속 대신 이걸 쓴다
            s.inherit_sp = Rd.last_label_status_xml(prs, cfg)
        except Exception:                                # noqa: BLE001
            s.label_style = None
            s.status_style = None
            s.inherit_sp = {}
        s.period_hist = {k: _period_history(seen, lab) for k, lab in PERIOD_KEYS.items()}
        # 기준일은 대개 "어느 차수부터" 다 — 파싱한 차수 날짜를 후보로 남긴다
        s.visit_dates = [{"visit": v["visit"], "date": v["date"]}
                         for v in scan["visits"] if v.get("date")]
        # 새 장이 이어붙을 자리 = 차수 글자가 가장 큰 장. 거기 그려진 선을 읽어
        # 검수 판에 겹쳐 보여준다 — 확정하면 그 자리로 따라오기 때문이다.
        if scan["visits"]:
            src_no = max(scan["visits"],
                         key=lambda v: (N.letter_to_num(v["visit"]),
                                        v["slide_no"]))["slide_no"]
            s.prev_lines = _slide_lines(prs, [src_no]).get(src_no, [])
    elif s.visit == "A":
        # 오늘이 초진이다 — 경과 개월은 0.
        s.first_date = datetime.now().strftime(cfg.ppt.info_date_format)

    # 확인 줄에서 고친 값 — 서버가 한 번 더 검증한다 (의료 기록물, 이중 잠금)
    if req.visit:
        v = req.visit.strip().upper()
        if not re.fullmatch(r"[A-Z]{1,2}", v):
            raise HTTPException(400, "차수는 영문 대문자 1~2글자여야 합니다")
        if v in visits:
            raise HTTPException(400, f"차수 {v} 는 이미 이 PPT에 있습니다")
        s.visit = v
    if req.insert_after is not None:
        last = scan["slides"] if has_ppt else 0
        if not 0 <= req.insert_after <= last:
            raise HTTPException(
                400, f"슬라이드 위치는 0~{last} 사이여야 합니다 (0 = 맨 앞)")
        s.insert_after = int(req.insert_after)
    SESSIONS[s.id] = s

    return {"session_id": s.id, "mode": s.mode, "visit": s.visit,
            "prev_visits": visits, "folder": d.name,
            "folder_exists": d.is_dir(), "ppt_exists": has_ppt,
            "windows": _windows_json(s.slot_windows),
            "ids": {"name": ids.name, "hospital_id": ids.hospital_id,
                    "ortho_id": ids.ortho_id}}


# ── 업로드 + 분류 (+ 재진 정합) ───────────────────────────────────────────────
def _choose_refs(refs_by_visit: dict[str, np.ndarray],
                 cur_visit: str) -> dict[str, np.ndarray]:
    """기준영상 후보를 **먼저 쓸 것부터** 담아 돌려준다.

    기본은 직전 차수다 — 바로 앞 회차와 같은 구도로 놓는 것이 이 도구의 목적이고,
    사람이 기대하는 결과이기도 하다. 초진(A)은 직전 정합이 실패했을 때만 쓰는
    안전망으로 남긴다: 직전만 연쇄로 따라가면 이전 차수의 배치 오차가 그대로
    물려 내려가는데, 초진이라는 고정점이 그 사슬을 끊는다.
    """
    if not refs_by_visit:
        return {}
    letters = sorted(refs_by_visit, key=N.letter_to_num)
    prev = [l for l in letters if N.letter_to_num(l) < N.letter_to_num(cur_visit)]
    chosen = {}
    if prev:
        chosen[f"직전({prev[-1]})"] = refs_by_visit[prev[-1]]
    if "A" in refs_by_visit and "A" not in [p for p in prev[-1:]]:
        chosen[f"초진(A)"] = refs_by_visit["A"]
    return chosen or {f"기준({letters[-1]})": refs_by_visit[letters[-1]]}


# 사진 투입(Setup)과 자동 분류(Pre-processing)를 분리한다.
# 붙어 있으면 "사진 몇 장 더 넣고 분류만 다시" 같은 게 불가능하다.
# JPEG이고 EXIF 회전이 없으면 원본 바이트를 그대로 둔다(무손상).
# 회전이 있으면 그때만 픽셀에 굽는다 — PowerPoint는 EXIF 회전을 무시하므로
# 원본을 그대로 넣으면 화면에선 똑바른 사진이 PPT에선 누워 버린다.
EXIF_ORIENT, EXIF_SUB_IFD = 274, 0x8769
EXIF_DT_TAGS = (36867, 36868)      # DateTimeOriginal, DateTimeDigitized
# 서브초(SubSecTime*). EXIF 의 촬영시각은 초까지만이라, 한 자리를 연달아 찍는
# 구내·안모 촬영에서는 같은 초에 여러 장이 몰려 순서가 안 갈린다. 기종에 따라
# 없을 수 있고, 있으면 1/10~1/1000초까지 들어온다.
EXIF_SUBSEC = {36867: 37521, 36868: 37522, 306: 37520}
EXIF_IMAGE_NUMBER = 37393          # ImageNumber — 기록하지 않는 기종이 더 많다
LOSSLESS_EXT = {".jpg", ".jpeg"}
# python-pptx 가 PPT 에 넣어 주는 포맷. 여기 없는 것은 반드시 다시 인코딩해야
# 한다 — 예를 들어 일부 카메라의 .jpg 는 실제로 MPO(다중 프레임 JPEG)라서,
# 확장자만 믿고 원본을 넘기면 확정 단계에서 통째로 실패한다.
PPTX_IMAGE_FORMATS = {"BMP", "GIF", "JPEG", "PNG", "TIFF", "WMF"}
_DIGITS = re.compile(r"\d+")


def _subsec_us(raw) -> int:
    """SubSecTime('83', '0421' …) → microsecond. 소수점 이하 자릿수를 그대로 해석."""
    digits = "".join(ch for ch in str(raw or "") if ch.isdigit())[:6]
    return int(digits.ljust(6, "0")) if digits else 0


def _exif_facts(im) -> tuple[int, datetime | None, int | None]:
    """(회전 플래그, 촬영시각, 일련번호). 읽을 수 없으면 (1, None, None).

    촬영시각은 서브초가 있으면 microsecond 까지 채운다 — 초 단위로만 두면
    연달아 찍은 사진의 순서가 동률이 되어 정렬이 무의미해진다(_shot_order_key).
    """
    try:
        ex = im.getexif()
    except Exception:
        return 1, None, None
    orient = ex.get(EXIF_ORIENT) or 1
    sub = {}
    try:
        sub = ex.get_ifd(EXIF_SUB_IFD) or {}
    except Exception:
        pass
    raw_seq = sub.get(EXIF_IMAGE_NUMBER, ex.get(EXIF_IMAGE_NUMBER))
    try:
        seq = int(raw_seq) if raw_seq is not None else None
    except (TypeError, ValueError):
        seq = None
    for tag, src in ((EXIF_DT_TAGS[0], sub), (EXIF_DT_TAGS[1], sub), (306, ex)):
        raw = src.get(tag)
        if not raw:
            continue
        try:
            dt = datetime.strptime(str(raw).strip(), "%Y:%m:%d %H:%M:%S")
        except ValueError:
            continue
        us = _subsec_us(src.get(EXIF_SUBSEC[tag], sub.get(EXIF_SUBSEC[tag])))
        return orient, dt.replace(microsecond=us), seq
    return orient, None, seq


def _name_seq(name: str | None) -> int:
    """파일명 안의 마지막 숫자 뭉치 (IMG_1234.JPG → 1234). 없으면 −1."""
    nums = _DIGITS.findall(Path(name or "").stem)
    return int(nums[-1]) if nums else -1


def _shot_order_key(p: "Photo", idx: int) -> tuple:
    """촬영 순서 정렬 키.

    EXIF 가 부실한 기종을 위해 단계적으로 물러난다:
      1) 촬영시각(서브초까지) 2) EXIF 일련번호 3) 파일명 끝 숫자 4) 업로드 순서

    카메라는 대개 파일명에도 촬영 순서대로 번호를 매기므로 3)만으로도 실무에서는
    거의 맞는다. 맨 뒤에 업로드 순서가 있어 어떤 경우에도 결과가 결정적이다.
    """
    return (0 if p.taken_at else 1,          # 시각을 아는 사진이 앞
            p.taken_at or datetime.min,
            p.exif_seq if p.exif_seq is not None else -1,
            _name_seq(p.orig_name),
            idx)


async def _stage_photos(s: "Session", files: list[UploadFile]) -> list[Photo]:
    """세션 임시폴더에 저장만 한다. 분류는 하지 않는다."""
    from PIL import Image as _Im, ImageOps as _Ops
    staged = []
    for uf in files:
        data = await uf.read()
        try:
            with _Im.open(io.BytesIO(data)) as im:
                im.load()
                orient, taken, seq = _exif_facts(im)
                pid = uuid.uuid4().hex[:10]
                dst = s.tmp / f"{pid}.jpg"
                ext = Path(uf.filename or "").suffix.lower()
                # 확장자가 아니라 **실제 포맷**을 본다. .jpg 로 저장된 MPO 를
                # 그대로 넘기면 여기서는 통과하고 확정에서 터진다.
                if (ext in LOSSLESS_EXT and im.format == "JPEG"
                        and orient in (0, 1)):
                    dst.write_bytes(data)          # 원본 그대로 — EXIF·화질 보존
                    pw, ph = im.size
                else:
                    fixed = _Ops.exif_transpose(im)
                    # exif_transpose 가 Orientation 태그를 지우므로 EXIF(촬영시각 등)를
                    # 실어도 이중 회전이 없다. 재인코딩 때 EXIF 가 통째로 사라지던
                    # 것을 보존한다.
                    fixed.convert("RGB").save(dst, "JPEG", quality=95, subsampling=0,
                                              exif=fixed.getexif().tobytes())
                    pw, ph = fixed.size
        except Exception:
            continue                                # 이미지가 아니거나 깨진 파일
        photo = Photo(pid, dst, pw, ph)
        photo.orig_name = uf.filename or dst.name
        photo.taken_at = taken
        photo.exif_seq = seq
        # 촬영시각을 파일 시각으로 새겨 둔다. 확정 때 copy2로 그대로 따라가서
        # 환자 폴더의 사진이 "찍은 날"을 갖게 된다.
        if taken:
            ts = taken.timestamp()
            os.utime(dst, (ts, ts))
        s.photos.append(photo)
        staged.append(photo)
    # 초진(A)의 초진일은 세션을 만든 날이 아니라 **사진을 찍은 날**이다. 세션 생성
    # 시점에는 사진이 없어 작업일을 임시로 넣는데, EXIF 가 도착하는 여기서 바로잡는다.
    # (재진은 이전 PPT 에서 읽은 초진일이 진실이므로 건드리지 않는다.)
    if s.visit == "A":
        pd = _photo_date(s)
        if pd:
            s.first_date = pd.strftime(cfg.ppt.info_date_format)
    return staged


@app.post("/api/photos/{sid}")
async def add_photos(sid: str, files: list[UploadFile] = File(...)):
    s = get_session(sid)
    staged = await _stage_photos(s, files)
    return {"added": len(staged),
            "photos": [_photo_json(s, p) for p in s.photos]}


@app.delete("/api/photos/{sid}/{pid}")
def drop_photo(sid: str, pid: str):
    s = get_session(sid)
    photo = _photo(s, pid)
    _unassign(s, pid)
    s.photos = [p for p in s.photos if p.id != pid]
    photo.path.unlink(missing_ok=True)
    return {"photos": [_photo_json(s, p) for p in s.photos]}


@app.post("/api/classify/{sid}")
def classify_session(sid: str):
    """투입된 사진을 분류하고 슬롯을 자동 배정한다.

    **정합·프레이밍은 안 한다** — `/api/register` 가 한다. 이유는 `_frame` 참고.
    """
    s = get_session(sid)
    _classify(s, s.photos)
    return {"photos": [_photo_json(s, p) for p in s.photos], "review": _review_json(s)}


class RegisterReq(BaseModel):
    slots: list[str] | None = None    # None 이면 배정된 슬롯 전부
    force: bool = False               # 이미 계산한 자리도 다시


@app.post("/api/register/{sid}")
def register_session(sid: str, req: RegisterReq = Body(default=RegisterReq())):
    """배정이 확정된 뒤 도는 무거운 단계. 화면은 슬롯을 하나씩 부르며 진행을 보인다."""
    s = get_session(sid)
    done = _frame(s, req.slots, force=req.force)
    return {"done": done, "photos": [_photo_json(s, p) for p in s.photos],
            "review": _review_json(s)}


def _classify(s: "Session", targets: list[Photo]) -> None:
    """라벨을 붙이고 상자에 넣는다. **여기까지가 가볍다.**"""
    thr = cfg.thresholds
    from PIL import Image as _Im
    for photo in targets:
        with _Im.open(photo.path) as _im:
            pred = classifier.predict(_im.copy(), filename=photo.orig_name)
        photo.label, photo.confidence, photo.probs = pred.label, pred.confidence, pred.probs

    # 같은 클래스가 여러 장이면 경쟁시켜 버리지 않고 한 상자에 쌓는다.
    # 상자 안은 신뢰도 내림차순 — 맨 위가 대표가 된다.
    slot_by_class = cfg.slot_by_class
    for photo in targets:
        if photo.label in slot_by_class:
            _put(s, photo, slot_by_class[photo.label])
        elif photo.label in cfg.face.classes:
            _put(s, photo, "FACE")
        # 그 외(OTHERS 등)는 미배정으로 남긴다 — 검수화면에서 수동 배정/제외 가능
        if photo.confidence < thr.classify_confidence:
            photo.badge = "low"
    for key, lst in s.bins.items():
        if key != "FACE":
            lst.sort(key=lambda pid: -_conf(s, pid))

    # 얼굴도 촬영 순서대로 한 번 놓아 둔다 — 사람이 빈 판에서 시작하지 않게.
    # 상자부터 촬영순으로 세운다(/api/sort 와 같은 규칙) — 업로드 순서가 뒤섞여
    # 있어도 첫 화면이 곧 촬영 흐름이 된다. 이미 고른 자리가 있으면 건드리지
    # 않는다(재분류로 지워지면 안 된다).
    if not s.face_slots:
        lst = s.bins.get("FACE", [])
        seen = {pid: i for i, pid in enumerate(lst)}
        lst.sort(key=lambda pid: _shot_order_key(_photo(s, pid), seen[pid]))
        _auto_assign_faces(s)


def _frame(s: "Session", slots: list[str] | None = None, *,
           force: bool = False) -> list[str]:
    """슬롯의 초기 구도를 잡는다 — 재진이면 **정합**, 초진이면 **프레이밍 모델**.

    ### 왜 분류와 떨어져 있나

    정합은 "이 사진이 **어느 자리**에 들어가는가"에 딸린 계산이다. 기준영상이 자리마다
    다르므로, 자리가 정해지기 전에 계산하면 값이 틀린다.

    예전에는 `_classify` 가 분류 직후 여기까지 했다. 그래서 분류기가 좌·우를 바꿔
    넣은 걸 사람이 고쳐도 **틀린 자리의 기준영상에 맞춘 배치**가 그대로 남았고,
    OTHERS 로 빠졌다가 수동 배정된 사진은 정합이 아예 없었다. 대표를 다른 장으로
    바꾸면 옛 대표의 배치를 물려받았다.

    이제 화면의 `검수·조정으로` 버튼이 이걸 부른다 — 그 시점에 배정이 확정된다.

    ### 다시 불러도 싸다

    `s.framed` 에 **어느 사진으로 계산했는지**를 남긴다. 대표가 그대로면 건너뛴다.
    좌·우를 고치고 돌아오면 그 두 칸만 돈다. `Reg.centers` 가 이미지 해시로 캐시하니
    기준영상은 캐시 적중이고 새 사진 한 장만 추론한다.

    반환: 실제로 계산한 슬롯 이름들.
    """
    reg_thr = cfg.thresholds.registration.model_dump()
    # 전체를 도는 길이면 **비워진 자리의 기록도 함께 훑는다** — 안 그러면 사진을
    # 뺀 자리에 옛 기록이 남고, 나중에 다른 사진이 오면 '이미 했다'로 건너뛴다.
    want = list(slots) if slots is not None else list(s.slots) + [
        k for k in s.framed if k not in s.slots]
    done: list[str] = []
    for slot in want:
        pid = s.slots.get(slot)
        if pid is None:
            s.framed.pop(slot, None)      # 비워진 자리 — 기록도 지운다
            continue
        if not force and s.framed.get(slot) == pid:
            continue
        photo = _photo(s, pid)
        win = s.slot_windows[slot]

        if s.mode != "revisit":
            _auto_frame(s, photo, win)
            s.framed[slot] = pid
            done.append(slot)
            continue

        arr = _imread(photo.path)
        refs = _choose_refs(s.references.get(slot, {}), s.visit)
        if not refs:
            # 기준 차수가 없다 — 정합할 대상이 없으니 초진과 같은 처지다.
            photo.ref_visit = None
            S.append_audit(LOG_FILE, {
                "event": "register_skipped", "reason": "no_reference_visits",
                "patient": s.patient_dir.name if s.patient_dir else "",
                "visit": s.visit, "slot": slot})
            _auto_frame(s, photo, win, fallback_badge="manual", bgr=arr)
        else:
            # 기준영상은 "이전 차수 PPT에서 보였던 그림"이라 교합면이면 이미
            # 뒤집혀 있다. 특징점 매칭도 유사변환(det>0)도 거울상은 다루지
            # 못하므로 신규 사진을 같은 방향으로 맞춰서 넣는다. 그러면 결과
            # 변환이 곧 반전 프레임 기준이라 photo.editor 로 그대로 들어간다.
            arr_reg = cv2.flip(arr, 0) if photo.flip_v else arr
            audit_base = {"patient": s.patient_dir.name if s.patient_dir else "",
                          "visit": s.visit, "slot": slot}
            # 정합이 실패하면 아래 else 로 가서 프레이밍 모델이 받는다.
            #
            # pseudo crop: 분할기에 학습 분포(완성본 모습)의 그림을 주고, 매칭도
            # 그 좌표에서 한다 (registration_teeth 독스트링 "좌표계"). 프레이밍
            # 추론은 무반전 원본으로 하고(_auto_frame 과 같은 규약), raw→pseudo 는
            # register 가 마지막에 합성하므로 crop 오차는 전파되지 않는다.
            try:
                pw = Reg.pseudo_frame(arr, framer, photo.label, flip_v=photo.flip_v)
                # 후보를 담긴 순서대로 시도하고 **되는 첫 기준에서 멈춘다**.
                # 점수로 고르면 직전이 멀쩡한데도 초진이 채택될 수 있는데,
                # 사람이 기대하는 것은 "바로 앞 회차와 같은 구도"다.
                best, res = None, None
                for _name, _ref in refs.items():
                    b, r, _ = Reg.register_best(arr_reg, {_name: _ref},
                                                thresholds=reg_thr, prewarp=pw)
                    if res is None or r.ok:
                        best, res = b, r
                    if r.ok:
                        break
            except Exception as e:                              # noqa: BLE001
                # 정합 오류가 검수 진입을 막으면 안 된다 — 남기고 프레이밍으로.
                S.append_audit(LOG_FILE, {
                    "event": "register_error", **audit_base,
                    "error": f"{type(e).__name__}: {e}"[:300]})
                photo.ref_visit = None
                _auto_frame(s, photo, win, fallback_badge="manual", bgr=arr)
                s.framed[slot] = pid
                done.append(slot)
                continue
            if res.ok:
                photo.editor = registration_to_editor(res.matrix, win, photo.w, photo.h)
                photo.ref_visit = best
                photo.badge = "ok"
                photo.framing = "registration"
            else:
                # 정합 실패 → cover-fit 보다 프레이밍 모델이 낫다. 배지는 '수동'을
                # 유지한다 — 차수 간 정렬이 안 됐다는 사실 자체는 변하지 않는다.
                # 어느 기준에 몇 개가 대응됐고 잔차가 얼마였는지 수치를 남긴다.
                S.append_audit(LOG_FILE, {
                    "event": "register_rejected", **audit_base, "ref": best,
                    "n_matches": res.n_matches, "n_inliers": res.n_inliers,
                    "reproj_error_px": round(res.reproj_error_px, 2),
                    "score": round(res.score, 4)})
                photo.ref_visit = best
                _auto_frame(s, photo, win, fallback_badge="manual", bgr=arr)
        s.framed[slot] = pid
        done.append(slot)
    return done


# 투입·분류·정합을 한 번에 하는 엔드포인트. 화면은 세 단계로 나눠 부르지만
# (사진을 더 넣고 분류만 다시 돌리는 흐름, 그리고 배정을 확정한 뒤에야 정합),
# 이 한 방 경로는 API 테스트와 외부 자동화가 쓴다 — 그쪽엔 고칠 사람이 없으므로
# 자동 배정을 그대로 확정으로 본다.
@app.post("/api/upload/{sid}")
async def upload(sid: str, files: list[UploadFile] = File(...)):
    s = get_session(sid)
    staged = await _stage_photos(s, files)
    _classify(s, staged)
    _frame(s)
    return {"photos": [_photo_json(s, p) for p in s.photos], "review": _review_json(s)}


def cover_fit_editor(photo, slot) -> EditorState:
    return EditorState()  # scale=1, 중심, 무회전 = cover-fit


def _auto_frame(s: "Session", photo, win: WindowCm, fallback_badge: str | None = None,
                bgr=None):
    """프레이밍 모델로 초기 배치를 잡는다. 못 쓰면 cover-fit 으로 물러난다.

    5개 구내 슬롯은 창이 전부 같은 4:3이라 결과가 슬롯에 의존하지 않는다. 그래도 창을
    받아 두는 이유는 나중에 슬롯 크기가 갈라져도 이 함수가 그대로 맞기 때문이다.

    배지 규약: **모델이 있는데 예측을 기각했을 때만** '수동'을 붙인다. 모델 자체가
    없는 설치본에서는 전부 cover-fit 이고 그건 종전 동작이므로 경고할 일이 아니다.
    """
    if framer is None or not framer.has(photo.label or ""):
        photo.editor = cover_fit_editor(photo, None)
        photo.framing = None
        if fallback_badge:
            photo.badge = fallback_badge
        return
    # 재진 폴백에서는 정합이 이미 디코드한 배열을 넘겨받는다 (6000x4000 JPEG 디코드가
    # 장당 ~114ms라 두 번 읽으면 그대로 두 배가 된다).
    # 프레이밍 모델은 **반전 없는 원본**으로 학습됐다. 교합면은 상하 비대칭이
    # 커서 뒤집어 넣으면 예외 없이 정확도만 조용히 떨어지므로, 추론은 항상
    # 원본으로 하고 결과 좌표만 반전 프레임으로 옮긴다.
    arr = bgr if bgr is not None else _imread(photo.path)
    res = framer.predict(arr, photo.label)
    photo.framing_note = res.method
    if res.ok:
        st = framing_to_editor(res, win, photo.w, photo.h)
        photo.editor = flip_editor_v(st) if photo.flip_v else st
        photo.framing = "model"
        if fallback_badge:
            photo.badge = fallback_badge
    else:
        photo.editor = cover_fit_editor(photo, None)
        photo.framing = "cover"
        photo.badge = fallback_badge or "manual"


# ── 배정 조작 ─────────────────────────────────────────────────────────────────
def _photo(s, pid) -> Photo:
    for p in s.photos:
        if p.id == pid:
            return p
    raise HTTPException(404, "사진 없음")


def _conf(s, pid): return _photo(s, pid).confidence


def _bin(s, key) -> list[str]:
    return s.bins.setdefault(key, [])


def _detach(s, pid) -> None:
    """어느 상자에 있든 빼낸다."""
    for lst in s.bins.values():
        if pid in lst:
            lst.remove(pid)
    # 얼굴 상자를 떠나면 케이스 덱에 잡아 둔 자리도 놓아준다.
    for cell, held in list(s.face_slots.items()):
        if held == pid:
            del s.face_slots[cell]
    _photo(s, pid).slot = None


def _sync_flip(photo) -> None:
    """슬롯에 맞춰 반전 여부를 갱신하고, 바뀌면 편집기 값을 그 프레임으로 옮긴다.

    `flip_v` 는 "editor 값이 어느 프레임 기준인가"를 뜻한다. 반전 프레임과 원본
    프레임의 (dx,dy,scale,angle) 은 dy·angle 의 부호만 다르고 **가리키는 사진
    영역은 같으므로**, 잘못 분류된 교합면을 다른 슬롯으로 옮겨도 사용자가 잡아
    둔 구도가 그대로 남는다.

    상자에서 빼기만 할 때(_detach)는 일부러 건드리지 않는다 — 값이 어느 프레임인지는
    변하지 않았고, 다시 어딘가에 넣을 때 여기서 한 번만 맞추면 된다.
    """
    want = photo.slot in cfg.flip_v_slots
    if want != photo.flip_v:
        photo.flip_v = want
        photo.editor = flip_editor_v(photo.editor)


def _put(s, photo, key, at=None) -> None:
    """상자에 넣는다. at=0이면 대표 자리, None이면 맨 뒤."""
    old = next(((k, v.index(photo.id)) for k, v in s.bins.items() if photo.id in v), None)
    _detach(s, photo.id)
    lst = _bin(s, key)
    idx = len(lst) if at is None else max(0, min(int(at), len(lst) + 1))
    # 같은 상자 안에서 아래로 옮길 때, 빼내면서 뒤쪽 인덱스가 한 칸 당겨진다
    if old and old[0] == key and old[1] < idx:
        idx -= 1
    lst.insert(min(idx, len(lst)), photo.id)
    photo.slot = key
    _sync_flip(photo)


def _unassign(s, pid): _detach(s, pid)


class AssignReq2(BaseModel):
    session_id: str
    photo_id: str
    slot: str | None          # 'SLOT_*' | 'FACE' | None(=OTHERS로 빼기)
    at: int | None = None     # 0이면 대표 자리로


@app.post("/api/assign")
def assign(req: AssignReq2):
    """상자 사이 이동. at=0이면 대표(슬라이드에 들어갈 사진)로 올린다."""
    s = get_session(req.session_id)
    photo = _photo(s, req.photo_id)
    if req.slot:
        _put(s, photo, req.slot, at=req.at)
    else:
        _detach(s, photo.id)
    return {"review": _review_json(s), "photos": [_photo_json(s, p) for p in s.photos]}


class SortReq(BaseModel):
    session_id: str
    slot: str                 # 'FACE' | 'SLOT_*'


@app.post("/api/sort")
def sort_bin(req: SortReq):
    """상자 안을 촬영 순서로 세운다.

    얼굴 사진은 분류기가 정면/45도/측면을 가르지 못해 사람이 자리를 직접 고른다.
    업로드 순서가 뒤섞여 있으면 촬영 흐름을 눈으로 되짚어야 하는데, 촬영 순서로
    세워 두면 그 흐름이 그대로 보인다. 정렬 근거는 _shot_order_key 참고.
    """
    s = get_session(req.session_id)
    lst = s.bins.get(req.slot)
    if lst is None:
        raise HTTPException(400, f"상자 '{req.slot}' 가 없습니다")
    before = list(lst)
    seen = {pid: i for i, pid in enumerate(lst)}
    lst.sort(key=lambda pid: _shot_order_key(_photo(s, pid), seen[pid]))
    # 촬영시각을 하나도 모르면 파일명·업로드 순서로만 세운 것이라 UI가 알려준다.
    known = sum(1 for pid in lst if _photo(s, pid).taken_at)
    # 얼굴 자리 배치는 "정렬한 순서"에 매인 값이다. 순서를 다시 세웠으면 배치도
    # 따라가야 화면과 규약이 어긋나지 않는다.
    placed = _auto_assign_faces(s) if req.slot == "FACE" else 0
    return {"changed": lst != before, "n": len(lst), "with_time": known,
            "placed": placed,
            "review": _review_json(s), "photos": [_photo_json(s, p) for p in s.photos]}


def _auto_assign_faces(s: "Session") -> int:
    """FACE 상자에 세워진 순서대로 얼굴 자리를 채운다. 돌려주는 값은 채운 자리 수.

    분류기는 정면/45도/측면을 가르지 못하지만, 촬영 루틴은 늘 같은 순서다.
    그래서 "촬영순으로 세운 n번째 사진 → n번째 자리"(config.face_auto_order)
    라는 고정 대응만으로 자리가 정해진다. 사람이 손으로 고치는 것은 그대로
    가능하고, 다시 정렬하면 이 함수가 그 위를 덮어쓴다.

    슬라이드 10·11 같은 파생 자리는 여기서 다루지 않는다 — _face_slots_json 이
    4L 을 따라가게 만든다.
    """
    order = [c for c in cfg.case_deck.face_auto_order if c in CASE_ANCHORS]
    if not order:
        return 0
    pool = list(s.bins.get("FACE", []))
    s.face_slots.clear()
    # 자리에 다른 사진이 들어오므로 잡아 둔 구도도 같이 버린다
    s.face_editors.clear()
    s.face_framing.clear()
    for cell, pid in zip(order, pool):
        s.face_slots[cell] = pid
    for cell in list(s.face_slots):
        s.face_framing[cell] = _frame_face_cell(s, cell)
    return min(len(order), len(pool))


def _face_frame_result(s: "Session", photo):
    """사진 한 장의 프레이밍 예측(캐시). 모델이 없거나 실패하면 None."""
    if photo.id in s.face_frames:
        return s.face_frames[photo.id]
    res = None
    if framer is not None and framer.has("FACE"):
        arr = _imread(photo.path)
        if arr is not None:
            res = framer.predict(arr, "FACE")
    s.face_frames[photo.id] = res
    return res


def _frame_face_cell(s: "Session", cell: str) -> str:
    """자리 하나의 초기 구도를 프레이밍 모델로 잡는다. 근거('model'|'cover')를 돌려준다.

    모델의 FACE crop 은 3:4(0.75)인데 케이스 덱의 얼굴 자리는 0.725 로 조금 더
    길쭉하다. `framing_to_editor` 는 예측을 창 안에 min() 으로 앉히므로 그대로
    두면 위아래에 검은 띠가 남는다. 얼굴은 **항상** cover 로 끌어올린다
    (config.geometry.allow_letterbox 와 무관) — 사람 얼굴에 검은 띠가 남는 것보다
    좌우를 3% 남짓 더 자르는 편이 낫다. 회전이 섞여도 같은 함수가 막아 준다.
    """
    pid = s.face_slots.get(cell)
    anchor = CASE_ANCHORS.get(cell)
    if not pid or anchor is None:
        return "cover"
    photo = _photo(s, pid)
    res = _face_frame_result(s, photo)
    if res is None or not res.ok:
        return "cover"
    win = anchor.window
    st = framing_to_editor(res, win, photo.w, photo.h)
    bw, bh = cover_base_ext_cm(photo.w, photo.h, win)
    s.face_editors[cell] = apply_cover_clamp(st, win, bw, bh)
    return "model"


class FaceAutoReq(BaseModel):
    session_id: str


@app.post("/api/face/auto")
def face_auto(req: FaceAutoReq):
    """얼굴 자리를 촬영 순서대로 다시 배치한다(손으로 고친 것도 덮어쓴다)."""
    s = get_session(req.session_id)
    placed = _auto_assign_faces(s)
    framed = sum(1 for v in s.face_framing.values() if v == "model")
    return {"placed": placed, "cells": len(FACE_CELLS), "framed": framed,
            "review": _review_json(s), "photos": [_photo_json(s, p) for p in s.photos]}


@app.get("/api/notes/{sid}")
def notes_get(sid: str):
    """입력 칸 정의 + 지금까지 채운 값 + 박스별 미리보기."""
    return _notes_json(get_session(sid))


class NotesReq(BaseModel):
    session_id: str
    values: dict[str, str] = {}    # 바꿀 칸만 담아도 된다
    # 박스 통째로 고쳐 쓰기(검수 화면 오버레이). 서식 결과를 이긴다.
    # 빈 문자열이면 덮어쓰기를 걷어내고 서식으로 되돌린다.
    boxes: dict[str, str] | None = None
    # 박스별 서식 자체를 이 차수에서만 바꾸기. None 이면 설정값으로 되돌린다.
    templates: dict[str, str | None] | None = None
    # 기간별 선택: {"tx": {"start": "24.09.26", "keep": false}}
    #   start  기준일. 목록에서 고르거나 이 차수 날짜를 새로 넣는다
    #   keep   이전 차수 값을 그대로 쓴다 (개월을 다시 안 센다)
    period: dict[str, dict] | None = None


@app.post("/api/notes")
def notes_set(req: NotesReq):
    s = get_session(req.session_id)
    for k, sel in (req.period or {}).items():
        if k not in PERIOD_KEYS:
            continue
        if "start" in sel:
            s.period_start[k] = sel["start"] or ""
        if "keep" in sel:
            s.period_keep[k] = bool(sel["keep"])
    known = {f.key for f in cfg.notes.fields}
    unknown = set(req.values) - known
    if unknown:
        raise HTTPException(400, f"모르는 칸입니다: {', '.join(sorted(unknown))}")
    s.note_fields.update({k: (v or "") for k, v in req.values.items()})
    if req.boxes is not None:
        bad = set(req.boxes) - set(cfg.notes.boxes) - set(NOTE_BOXES)
        if bad:
            raise HTTPException(400, f"모르는 박스입니다: {', '.join(sorted(bad))}")
        for k, v in req.boxes.items():
            if v:
                s.note_overrides[k] = v
            else:
                s.note_overrides.pop(k, None)   # 서식으로 되돌리기
    if req.templates is not None:
        bad = set(req.templates) - set(cfg.notes.boxes)
        if bad:
            raise HTTPException(400, f"모르는 박스입니다: {', '.join(sorted(bad))}")
        for k, v in req.templates.items():
            if v is None:
                s.note_templates.pop(k, None)   # 설정값으로 되돌리기
            else:
                s.note_templates[k] = v
            # 서식을 손봤다는 것은 자동 채움을 다시 쓰겠다는 뜻이다 —
            # 통째로 고쳐 쓴 값이 남아 있으면 바꾼 서식이 화면에 안 보인다.
            s.note_overrides.pop(k, None)
    return _notes_json(s)


class LinesReq(BaseModel):
    session_id: str
    # "슬라이드:도형이름" -> [dx_cm, dy_cm]. 빈 목록이면 원래 자리로 되돌린다.
    moves: dict[str, list[float] | None]


@app.post("/api/lines")
def lines_set(req: LinesReq):
    """계측선을 끌어 옮긴 양을 저장한다. 확정할 때 그대로 반영된다."""
    s = get_session(req.session_id)
    known = {ln["id"] for v in CASE_LINES.values() for ln in v}
    bad = set(req.moves) - known
    if bad:
        raise HTTPException(400, f"모르는 선입니다: {', '.join(sorted(bad))}")
    for k, v in req.moves.items():
        if not v:
            s.line_moves.pop(k, None)          # 제자리로
        else:
            s.line_moves[k] = [float(v[0]), float(v[1])]
    return {"moves": dict(s.line_moves)}


def _apply_line_moves(prs, s: "Session") -> int:
    """저장해 둔 이동량을 덱의 실제 도형에 반영한다. 옮긴 선 수를 돌려준다."""
    moved = 0
    for key, (dx, dy) in s.line_moves.items():
        slide_no, _, name = key.partition(":")
        try:
            slide = prs.slides[int(slide_no) - 1]
        except (ValueError, IndexError):
            continue
        for sh in slide.shapes:
            if sh.name != name:
                continue
            sh.left = int(sh.left + dx * EMU_PER_CM)
            sh.top = int(sh.top + dy * EMU_PER_CM)
            moved += 1
            break
    return moved


@app.get("/api/case/layout")
def case_layout():
    """케이스 덱의 얼굴 자리 표. 양식에서 읽은 것이라 세션과 무관하다."""
    return _face_layout_json()


class FaceAssignReq(BaseModel):
    session_id: str
    cell: str                 # "4L" | "7C" ...
    photo_id: str | None      # None이면 그 자리를 비운다


@app.post("/api/face/assign")
def face_assign(req: FaceAssignReq):
    """
    얼굴 사진을 케이스 덱의 한 자리에 배정한다.

    한 사진은 한 자리에만 놓인다 — 다른 자리에 이미 있었다면 그쪽을 비운다.
    파생 자리(10·11)는 슬라이드 4 좌측을 따라가므로 직접 배정할 수 없다.
    """
    s = get_session(req.session_id)
    if req.cell not in FACE_CELLS:
        raise HTTPException(400, f"배정할 수 없는 자리입니다: {req.cell}")

    before = s.face_slots.get(req.cell)
    if req.photo_id is None:
        s.face_slots.pop(req.cell, None)
    else:
        photo = _photo(s, req.photo_id)
        if photo.id not in s.face:
            raise HTTPException(400, "얼굴 상자에 있는 사진만 배정할 수 있습니다")
        for k, pid in list(s.face_slots.items()):
            if pid == photo.id:
                # 사진이 떠난 자리의 구도는 남겨 둘 이유가 없다
                del s.face_slots[k]
                s.face_editors.pop(k, None)
                s.face_framing.pop(k, None)
        s.face_slots[req.cell] = photo.id
    # 자리의 사진이 바뀌면 그 자리에 잡아 둔 구도는 다른 사진 기준이라 무의미하다.
    # 새 사진에는 다시 프레이밍을 걸어 준다(예측은 사진 단위로 캐시돼 있다).
    if s.face_slots.get(req.cell) != before:
        s.face_editors.pop(req.cell, None)
        s.face_framing.pop(req.cell, None)
        if s.face_slots.get(req.cell):
            s.face_framing[req.cell] = _frame_face_cell(s, req.cell)
    return {"face_slots": _face_slots_json(s), "face_editors": _face_editors_json(s),
            "face_framing": dict(s.face_framing)}


class FaceAdjustReq(BaseModel):
    session_id: str
    cell: str
    dx: float
    dy: float
    scale: float
    angle: float


@app.post("/api/face/adjust")
def face_adjust(req: FaceAdjustReq):
    """얼굴 자리 하나의 구도를 저장한다 (구내의 /api/adjust 와 같은 규약)."""
    s = get_session(req.session_id)
    anchor = CASE_ANCHORS.get(req.cell)
    if anchor is None or req.cell in MIRROR_CELLS:
        raise HTTPException(400, f"조정할 수 없는 자리입니다: {req.cell}")
    pid = s.face_slots.get(req.cell)
    if not pid:
        raise HTTPException(400, "빈 자리입니다")
    photo = _photo(s, pid)
    win = anchor.window
    bw, bh = cover_base_ext_cm(photo.w, photo.h, win)
    st = _clamp(EditorState(req.dx, req.dy, req.scale, req.angle), win, bw, bh)
    s.face_editors[req.cell] = st
    return {"clamped_scale": st.scale}


def _face_slots_json(s: Session) -> dict:
    """자리 -> 사진. 파생 자리도 함께 채워 내려보낸다(화면이 계산하지 않게)."""
    out = {k: pid for k, pid in s.face_slots.items()}
    src = s.face_slots.get(MIRROR_SOURCE)
    for k in MIRROR_CELLS:
        if src:
            out[k] = src
    return out


def _face_editors_json(s: Session) -> dict:
    """자리 -> 편집기 값. 파생 자리는 환산된 값을 함께 내려보낸다."""
    out = {}
    for k in list(FACE_CELLS) + list(MIRROR_CELLS):
        if k not in _face_slots_json(s):
            continue
        st = _face_editor(s, k)
        out[k] = {"dx": round(st.dx_px, 2), "dy": round(st.dy_px, 2),
                  "scale": round(st.scale, 4), "angle": round(st.angle_deg, 3)}
    return out


@app.post("/api/adjust")
def adjust(req: AdjustReq):
    s = get_session(req.session_id)
    pid = s.slots.get(req.slot)
    if not pid:
        raise HTTPException(400, "슬롯이 비어있음")
    photo = _photo(s, pid)
    win = s.slot_windows[req.slot]
    bw, bh = cover_base_ext_cm(photo.w, photo.h, win)
    st = _clamp(EditorState(req.dx, req.dy, req.scale, req.angle), win, bw, bh)
    photo.editor = st
    pl = editor_to_placement(st, win, bw, bh, PPC)
    return {"placement": {"off_x": pl.off_x, "off_y": pl.off_y,
                          "ext_cx": pl.ext_cx, "ext_cy": pl.ext_cy, "rot": pl.rot},
            "clamped_scale": st.scale}


# ── 이미지 서빙 ───────────────────────────────────────────────────────────────
@app.get("/api/thumb/{sid}/{pid}")
def thumb(sid: str, pid: str):
    s = get_session(sid)
    p = _photo(s, pid)
    return FileResponse(p.path)


@app.get("/api/reference/{sid}/{slot}")
def reference(sid: str, slot: str, visit: str = ""):
    """기준영상 한 장. **이전 차수 슬라이드에 보이던 그림 그대로**다.

    창 기준으로 이미 맞춰져 있고 교합면이면 이미 뒤집혀 있다 — 그래서 화면에서
    창에 그대로 깔면 지금 편집 중인 구도와 바로 겹쳐 볼 수 있다(겹쳐보기).
    """
    s = get_session(sid)
    refs = s.references.get(slot, {})
    if not refs:
        raise HTTPException(404, "기준영상 없음")
    key = visit if visit in refs else sorted(refs, key=N.letter_to_num)[-1]
    ok, buf = cv2.imencode(".png", refs[key])
    if not ok:
        raise HTTPException(500, "기준영상을 만들지 못했습니다")
    return StreamingResponse(io.BytesIO(buf.tobytes()), media_type="image/png")


@app.get("/api/references/{sid}")
def reference_list(sid: str):
    """슬롯마다 **어느 차수를 겹쳐볼 수 있나**. 화면의 겹쳐보기 목록이 이걸 쓴다."""
    s = get_session(sid)
    return {slot: sorted(refs, key=N.letter_to_num)
            for slot, refs in s.references.items() if refs}


# ── 확정 (원자적 저장) ────────────────────────────────────────────────────────
def _build_plan(s) -> dict:
    """
    확정하면 **무엇이 어떤 이름으로 어디에 생기는지** 계산한다. 부수효과 없음.

    commit()이 이 결과를 그대로 쓴다. 미리보기와 실제 저장이 각자 이름을 계산하면
    언젠가 반드시 갈라지고, 그건 의료 기록물에서 나면 안 되는 종류의 버그다.
    """
    ids = s.ids
    index_by_class = cfg.index_by_class
    raw = _save_raw()
    # 새 사진은 차수마다 제 폴더로 간다 — 잘린 완성본은 '교정번호_차수/',
    # 원본 사본은 '교정번호_차수_raw/'. 기존 사진이 어디 있든 추적하지 않는다.
    ppre = f"{N.visit_dir(ids.ortho_id, s.visit)}/"
    rpre = f"{N.visit_raw_dir(ids.ortho_id, s.visit)}/"
    slots = []
    for slot in cfg.ppt.slot_names:
        members = s.bins.get(slot, [])
        if not members:
            slots.append({"slot": slot, "empty": True})
            continue
        cls = _slot_to_class(slot)
        idx = index_by_class[cls]
        base = N.photo_filename(ids.ortho_id, s.visit, idx, cfg.naming.photo_pattern)
        slots.append({
            "slot": slot, "empty": False, "cls": cls, "index": idx,
            "label": _photo(s, members[0]).label,
            "file": ppre + base,
            # 원본 사본. 추가 촬영본에는 없다 — 편집값이 없어 자를 것이 없고,
            # 그쪽은 지금도 원본 그대로 저장된다.
            "raw": (rpre + N.raw_filename(base, _photo(s, members[0]).path.name)
                    if raw else None),
            "extras": [
                {"label": _photo(s, pid).label,
                 "file": ppre + N.photo_extra_filename(ids.ortho_id, s.visit, idx, n,
                                                       cfg.naming.photo_extra_pattern)}
                for n, pid in enumerate(members[1:], start=2)],
        })
    faces, fidx = [], cfg.face.start_index
    for pid in s.face:
        base = N.photo_filename(ids.ortho_id, s.visit, fidx, cfg.naming.photo_pattern)
        faces.append({"label": _photo(s, pid).label, "file": ppre + base,
                      "raw": (rpre + N.raw_filename(base, _photo(s, pid).path.name)
                              if raw else None)})
        fidx += 1
    return {
        "patient_dir": str(s.patient_dir),
        "visit": s.visit,
        "mode": s.mode,
        # 기존 PPT 를 알아봤으면 **그 파일에** 이어 쓴다 — 이름이 옛 형식이거나
        # 하위 폴더에 있어도 원본 위치가 진실이다. 루트에 현재 양식 이름으로 또
        # 만들면 어느 쪽이 진짜인지 다투는 사본이 남는다. 새 PPT 만 생성 이름.
        "ppt": (Path(s.ppt_path).relative_to(s.patient_dir).as_posix()
                if s.ppt_path and Path(s.ppt_path).exists()
                else _gen_ppt_name(ids)),
        "ppt_exists": bool(s.ppt_path and Path(s.ppt_path).exists()),
        "slots": slots,
        "faces": faces,
        "missing": [slot for slot in cfg.ppt.slot_names if slot not in s.slots],
    }


@app.get("/api/plan/{sid}")
def plan(sid: str):
    """저장 직전 검토용 드라이런. 아무것도 쓰지 않는다."""
    return _build_plan(get_session(sid))


@app.post("/api/commit/{sid}")
def commit(sid: str, allow_missing: bool = False):
    s = get_session(sid)
    pl_plan = _build_plan(s)
    missing = pl_plan["missing"]
    if missing and not allow_missing:
        return JSONResponse(status_code=409, content={"error": "missing_slots", "missing": missing})

    ids = s.ids
    # PPT 에 기록되는 방문 날짜 — 확정을 누른 날이 아니라 사진을 찍은 날.
    date_str = (_photo_date(s) or datetime.now()).strftime(cfg.ppt.info_date_format)
    ppt_name = pl_plan["ppt"]

    # PowerPoint 가 열어 둔 PPT 는 덮어쓸 수 없다 — 무거운 작업 전에 먼저 알린다.
    if s.mode == "revisit" and s.ppt_path and Path(s.ppt_path).exists():
        try:
            with open(s.ppt_path, "rb+"):
                pass
        except PermissionError:
            return JSONResponse(status_code=409, content={
                "error": "ppt_locked",
                "detail": "PPT가 다른 프로그램(PowerPoint)에서 열려 있습니다 — "
                          "닫은 뒤 다시 확정해 주세요"})

    try:
        src_slide = None          # 도형을 물려받을 직전 차수 슬라이드 (재진에서만)
        with S.Transaction(s.patient_dir) as tx:
            # 1) PPT 준비
            if s.mode == "first":
                stage_ppt = s.tmp / ppt_name
                prs, slide = _new_first_visit_ppt(stage_ppt)
                # 사진은 있는데 PPT만 없는 폴더는 mode='first'로 새 PPT를 만들지만
                # 차수는 A가 아니다 — 그때 '(초진)'이라고 쓰면 기록이 틀린다.
                # 두 서식 모두 {visit} 를 쓴다 — "(초진 A)" 처럼 차수 글자가 붙는다
                info_text = _render_label(date_str, s.visit,
                                          getattr(s, "label_fp", None))
            else:
                stage_ppt = s.tmp / ppt_name
                shutil.copyfile(s.ppt_path, stage_ppt)
                prs = T.load_presentation(stage_ppt)
                # 확인 줄에서 고쳤으면 "그 번호의 장 뒤", 아니면 날짜순 규칙.
                # n 번 장 뒤 = 0-기반 삽입 위치 n (그래서 새 장은 n+1 번이 된다).
                pos = getattr(s, "insert_after", None)
                insert_idx = (min(max(pos, 0), len(prs.slides._sldIdLst))
                              if pos is not None else _revisit_insert_index(prs))
                # 도형을 물려받을 원본 = 새 장이 끼어들 **바로 앞 장**(직전 차수)
                src_slide = (prs.slides[insert_idx - 1] if insert_idx > 0 else None)
                slide = W.import_template_slide(prs, TEMPLATE_PRS, insert_idx)
                # 십자뷰 양식에는 노트 칸이 없다 — 그대로 두면 이번 차수에 적은
                # 노트가 갈 곳이 없어 조용히 사라진다. 화면 오버레이가 쓰는
                # 자리표(NOTE_BOXES) 그대로 만들어 화면과 결과물을 맞춘다.
                # 날짜 칸은 이 슬라이드의 INFO_BOX 가 맡는다(자리가 겹친다).
                if NOTE_BOXES:
                    CD.add_note_boxes_from_layout(
                        slide, NOTE_BOXES,
                        skip={CD.NOTE_DATE} if T.find_shape(slide, cfg.ppt.info_box_name) is not None else set())
                info_text = _render_label(date_str, s.visit,
                                          getattr(s, "label_fp", None))
            # 검은 마스크(MASK_*)는 **모든 경우** 제거한다 (2026-08-12 결정).
            # 근거: ① 사진을 창 크기로 구워 넣어 초과가 없고, ② 슬라이드 배경
            # 자체가 검정(000000)이라 시각적으로 동일하며, ③ 수제 레이아웃 상속
            # 시 템플릿 좌표의 마스크가 사진 가장자리를 가리는 문제도 사라진다.
            for sh in [x for x in slide.shapes if x.name.startswith("MASK_")]:
                sh._element.getparent().remove(sh._element)

            # 수제 PPT 상속 — 원본 라벨/상태 상자를 **통째로 복사**해 규약명으로
            # 얹는다 (템플릿 상자는 제거). 속성 개별 상속은 숨은 규칙(lstStyle
            # 기본값·endParaRPr 등)을 계속 놓쳐서 이 방식으로 바꿨다 (2026-08-12).
            inherit = (getattr(s, "inherit_sp", None) or {}) if s.mode == "revisit" else {}
            if inherit.get("label"):
                CD.replace_with_copied_box(slide, cfg.ppt.info_box_name,
                                           inherit["label"])
            if inherit.get("status"):
                CD.replace_with_copied_box(slide, CD.NOTE_STATUS,
                                           inherit["status"])
            # 나머지 노트 상자도 원본에 대응이 있으면 통째 복사 — 폰트 유지.
            # 없는 상자만 설정 기본 크기(add_note_boxes_from_layout)로 만들어진다.
            for key in (CD.NOTE_SOAP, CD.NOTE_LL, CD.NOTE_NEXT):
                if inherit.get(key):
                    CD.replace_with_copied_box(slide, key, inherit[key])
            # 사람이 그려 둔 선·화살표 등 — 설정에 따라 직전 장에서 가져온다
            if s.mode == "revisit":
                _inherit_shapes(src_slide, slide, _copy_shapes(),
                                emu_to_cm(prs.slide_width), set(inherit))

            # 2) 구내 슬롯 삽입 — 상자의 대표(0번)만 슬라이드에 들어간다
            #    파일명은 전부 _build_plan() 이 정한 것을 쓴다(미리보기와 동일 보장).
            for entry in pl_plan["slots"]:
                if entry["empty"]:
                    continue
                slot = entry["slot"]
                members = s.bins[slot]
                photo = _photo(s, members[0])
                win = s.slot_windows[slot]
                # 수제 레이아웃 상속 — place_photo_in_slot 은 **앵커 위치**에
                # 놓으므로, 새 슬라이드의 앵커·배경판을 세션 창(상속된 레이아웃)
                # 으로 먼저 옮겨야 사진이 사람 레이아웃 자리에 들어간다.
                for nm in (slot, W.backdrop_shape_name(slot)):
                    shp = T.find_shape(slide, nm)
                    if shp is not None:
                        shp.left = int(round(win.x * EMU_PER_CM))
                        shp.top = int(round(win.y * EMU_PER_CM))
                        shp.width = int(round(win.w * EMU_PER_CM))
                        shp.height = int(round(win.h * EMU_PER_CM))
                bw, bh = cover_base_ext_cm(photo.w, photo.h, win)
                pl = editor_to_placement(photo.editor, win, bw, bh, PPC)
                # 창에 보이는 만큼만 구워서 넣는다 — 그래야 이웃 슬롯을 침범하지
                # 않는다. 환자 폴더에도 **이것**이 간다: 폴더와 PPT 가 다른 그림이면
                # 나중에 어느 쪽이 진짜인지 다투게 된다.
                baked, bwh = _bake_window(photo, win, photo.editor, photo.flip_v,
                                          s.tmp / f"bake_{slot}.jpg")
                if baked:
                    # 창에 맞춰 구웠으므로 창 그대로 넣는다 — cover-fit 에 맡기면
                    # 구운 파일의 정수 픽셀 비율로 크기를 다시 셈해 직전 차수
                    # 사진과 0.01cm 어긋나 보인다.
                    W.place_photo_in_slot(slide, slot, baked, bwh,
                                          placement=_exact_placement(win),
                                          letterbox_color=_letterbox_color())
                    tx.stage_file(baked, entry["file"])
                else:
                    staged_img = tx.stage_file(photo.path, entry["file"])
                    W.place_photo_in_slot(slide, slot, staged_img, (photo.w, photo.h),
                                          placement=pl,
                                          letterbox_color=_letterbox_color(),
                                          flip_v=photo.flip_v)
                if entry.get("raw"):
                    tx.stage_file(photo.path, entry["raw"])
                # 같은 자리의 추가 촬영본: 파일로만 저장 (슬라이드는 대표 1장)
                for extra_pid, ex in zip(members[1:], entry["extras"]):
                    tx.stage_file(_photo(s, extra_pid).path, ex["file"])
            # 날짜 칸을 검수 화면에서 고쳐 썼으면 그쪽이 이긴다
            _write_visit_label(slide, s.note_overrides.get(CD.NOTE_DATE) or info_text,
                           None if inherit else getattr(s, "label_style", None),
                           pin=not inherit)
            # 차수 노트 — 채운 칸이 있는 박스만 건드린다(양식의 안내문을 함부로 지우지 않는다)
            # 자동 계산 기본값·오버레이 수정본이 함께 나간다 — 화면 미리보기와
            # 결과물이 같아야 한다
            for box, text in _note_text(s).items():
                # 빈 줄만 남은 박스는 '안 채운 것'이다. 서식의 빈 줄을 그대로
                # 살리면서(양식의 글 시작 높이) 손 안 댄 박스의 안내문은 지키려면
                # 공백을 걷어내고 판단해야 한다.
                if text.strip():
                    CD.set_note_text(slide, box, text, small_pt=_small_pt(box))
            # 통째 복사 상속이면 속성 조정이 필요 없다 — 빈 줄 크기만 바로잡는다
            if inherit:
                for key in (CD.NOTE_STATUS, cfg.ppt.info_box_name,
                            CD.NOTE_SOAP, CD.NOTE_LL, CD.NOTE_NEXT):
                    CD.fix_empty_para_sizes(slide, key)
            else:
                st_style = getattr(s, "status_style", None)
                if st_style:
                    CD.style_note_box(slide, CD.NOTE_STATUS, st_style)

            # 3) 얼굴. 케이스 덱이면 배정된 자리에 먼저 놓고, 거기서 구운 사본을
            #    파일로도 저장한다 — 슬라이드와 폴더가 같은 그림이어야 한다.
            face_bakes: dict[str, Path] = {}
            if s.mode == "first" and CASE_ANCHORS:
                face_bakes = _place_faces(prs, s)
                # 구내 한 장짜리 슬라이드(12~16)
                _place_intraoral(prs, s)
                # 검수 화면에서 끌어 옮긴 계측선을 실제 도형에 반영한다
                _apply_line_moves(prs, s)
                # 양식 첫 장(환자정보)을 이 환자의 값으로 채운다
                _fill_patient_info(prs, s)
            for pid, fe in zip(s.face, pl_plan["faces"]):
                photo = _photo(s, pid)
                tx.stage_file(face_bakes.get(pid) or photo.path, fe["file"])
                if fe.get("raw"):
                    tx.stage_file(photo.path, fe["raw"])

            # 4) PPT 저장 후 원자적 확정
            tx.stage_pptx(prs, ppt_name)
            moved = tx.commit()

        S.append_audit(LOG_FILE, {
            "event": "commit", "mode": s.mode, "visit": s.visit,
            "patient": s.patient_dir.name,
            # 이름이 아니라 환자 폴더 기준 상대경로 — 원본은 raw/ 하위로 간다
            "files": [p.relative_to(s.patient_dir).as_posix() for p in moved],
            "slots": {k: _photo(s, v).label for k, v in s.slots.items()},
        })
    except PermissionError as e:
        S.append_audit(LOG_FILE,
                       {"event": "commit_failed", "error": str(e)})
        raise HTTPException(409, "PPT가 다른 프로그램(PowerPoint)에서 열려 있어 "
                                 "저장하지 못했습니다(롤백됨) — 닫고 다시 확정해 "
                                 "주세요")
    except Exception as e:
        S.append_audit(LOG_FILE,
                       {"event": "commit_failed", "error": str(e)})
        raise HTTPException(500, f"확정 실패(롤백됨): {e}")

    # 다음에도 같은 덱으로 이어지도록 이번에 쓴 파일을 기억한다
    _remember_ppt(s.patient_dir.name, ppt_name)
    result = {"ok": True, "patient_dir": str(s.patient_dir),
              "ppt": ppt_name, "visit": s.visit,
              "files": [p.relative_to(s.patient_dir).as_posix() for p in moved]}
    discard_session(s)   # 업로드 원본은 환자 폴더로 복사됐다 — 임시본을 남기지 않는다
    return result


def _revisit_insert_index(prs) -> int:
    """새 차수 슬라이드가 들어갈 자리.

    ① **차수 글자가 가장 큰 십자뷰 슬라이드 바로 다음.** 예전에는 날짜가 가장
       늦은 장을 골랐는데, 라벨 날짜는 손으로 적다 보니 오타가 난다 — J 가 K 보다
       뒤 날짜인 덱이 실제로 있었고, 그때 새 장이 K 앞으로 들어갔다. 차수 글자는
       순서 그 자체라 그런 흔들림이 없다 (2026-08-14 결정).
    ② 글자를 하나도 못 읽으면: PHOTO_ 이름, 또는 유효 십자뷰(가로 8cm 이상
       사진 5장)인 마지막 장 다음.
    ③ 그것도 없으면 문서 맨 뒤.
    """
    scan = Rd.scan_ppt_visits(prs, cfg)
    if scan["visits"]:
        last = max(scan["visits"],
                   key=lambda v: (N.letter_to_num(v["visit"]), v["slide_no"]))
        return last["slide_no"]        # n 번 장 뒤 = 0-기반 삽입 위치 n
    last_io = -1
    for i, slide in enumerate(prs.slides):
        shapes = list(slide.shapes)
        big = sum(1 for sh in shapes
                  if getattr(sh, "shape_type", None) == 13
                  and emu_to_cm(sh.width) >= 8.0)
        has_photo = any(sh.name.startswith(W.PHOTO_NAME_PREFIX) for sh in shapes)
        if has_photo or big >= 5:      # 유효 십자뷰 기준 — 가로 8cm 이상 5장
            last_io = i
    return last_io + 1 if last_io >= 0 else len(prs.slides._sldIdLst)


# 새 슬라이드가 스스로 만드는 도형들 — 복사 대상에서 뺀다(겹쳐 두 벌이 된다).
_RESERVED_PREFIX = (W.PHOTO_NAME_PREFIX, "BACKDROP_", "MASK_")
_RESERVED_NAME = {cfg.ppt.info_box_name, CD.NOTE_DATE, CD.NOTE_STATUS,
                  CD.NOTE_SOAP, CD.NOTE_LL, CD.NOTE_NEXT}


def _copy_shapes() -> str:
    """직전 차수 슬라이드의 도형을 새 슬라이드로 가져올까 — "none"|"lines"|"all".

    기본은 "lines" 다. 정중선·교합평면 같은 기준선은 매 차수 같은 자리를 가리키므로
    따라오는 편이 맞고, 선은 글과 달리 옛 차수의 내용을 실어 나르지 않는다.
    """
    try:
        v = json.loads(SETTINGS_FILE.read_text(encoding="utf-8")).get("copy_shapes")
        return v if v in ("none", "lines", "all") else "lines"
    except Exception:                                   # noqa: BLE001
        return "lines"


# 글을 물려받을 수 있는 자유 기입 상자 — 날짜/차수·기간 상자는 매 차수 새로 쓴다.
_FREE_NOTE_BOXES = (CD.NOTE_SOAP, CD.NOTE_LL, CD.NOTE_NEXT)


def _inherit_shapes(src, dst, mode: str, slide_w_cm: float = 0.0,
                    inherited: set | None = None) -> int:
    """직전 차수 슬라이드의 도형을 새 슬라이드로 복사한다.

    "lines" 는 직선·연결선만 — 정중선·교합평면처럼 매 차수 같은 자리를 가리키는
    기준선이 이쪽이다. "all" 은 거기에 글상자와 **그 안의 글까지** 더한다: 지난
    차수 내용을 이어 고쳐 쓰는 방식이라, 그 차수에만 해당하는 주석도 따라온다.

    자유 기입 상자(좌상단 s/p·좌하단·우하단)의 글은 "all" 에서만 물려받는다.
    다른 모드에서는 상속으로 딸려 온 글을 지운다 — 수제 덱은 상자를 통째로
    복사해 오므로 지우지 않으면 지난 차수 글이 그대로 남는다.

    규약 상자 다섯은 이 경로로 오지 않는다. 이름이 있으면 이름으로, 이름이 없는
    수제 상자는 **역할로** 가려낸다(`Rd.note_role`) — 이름만 보면 이미 물려받은
    날짜/차수·Tx/Rx/App 이 한 벌 더 얹혀 글자가 겹친다.

    `inherited` 는 이번에 **실제로 물려받은 역할들**이다. 역할이 같다고 무조건
    빼면 안 된다: 앱이 만든 직전 슬라이드는 물려받는 것이 하나도 없는데, 거기에
    손으로 그려 둔 글상자가 노트 자리에 있다는 이유로 통째로 버려졌다.
    """
    if src is None or dst is None:
        return 0
    n = 0
    for key in _FREE_NOTE_BOXES:
        if mode == "all":
            sh = T.find_shape(src, key)
            t = (sh.text_frame.text if sh is not None
                 and getattr(sh, "has_text_frame", False) else "") or ""
            if t.strip() and CD.set_note_text(dst, key, t, small_pt=_small_pt(key)):
                n += 1
            continue
        sh = T.find_shape(dst, key)
        if sh is not None and getattr(sh, "has_text_frame", False) \
                and (sh.text_frame.text or "").strip():
            CD.set_note_text(dst, key, "")
    if mode == "none":
        return n
    # 규약 상자로 **실제로 물려받은 그 도형만** 건너뛴다. 역할이 같다고 다 버리면
    # 같은 구역에 글상자가 둘일 때 하나만 상속되고 나머지는 소리 없이 사라진다.
    # 어느 것이 상속되는지는 last_label_status_xml 과 같은 규칙(뒤엣것이 이긴다).
    taken: dict[str, object] = {}
    if slide_w_cm and inherited:
        for sh in src.shapes:
            role = Rd.note_role(sh, slide_w_cm)
            if role in inherited:
                taken[role] = sh
    skip = {id(sh._element) for sh in taken.values()}
    spTree = dst.shapes._spTree
    for sh in src.shapes:
        name = str(getattr(sh, "name", ""))
        if name.startswith(_RESERVED_PREFIX) or name in _RESERVED_NAME:
            continue
        if getattr(sh, "shape_type", None) == 13:       # PICTURE — 사진은 새로 넣는다
            continue
        if id(sh._element) in skip:
            continue                        # 이미 규약 상자로 물려받은 그 도형
        if mode == "lines" and getattr(sh, "shape_type", None) != _LINE_TYPE:
            continue
        new = copy.deepcopy(sh._element)
        ids = [int(e.get("id")) for e in spTree.iter(qn("p:cNvPr"))
               if (e.get("id") or "").isdigit()]
        for cnv in new.iter(qn("p:cNvPr")):
            cnv.set("id", str(max(ids, default=1) + 1))
            break                                       # 최상위 하나만 새 번호로
        spTree.append(new)
        n += 1
    return n


def _slot_to_class(slot):
    for cls, sl in cfg.slot_by_class.items():
        if sl == slot:
            return cls
    return None


# ── JSON 직렬화 ───────────────────────────────────────────────────────────────
def _photo_json(s, p: Photo):
    return {"id": p.id, "label": p.label, "confidence": round(p.confidence, 3),
            "slot": p.slot, "badge": p.badge, "ref_visit": p.ref_visit,
            "framing": p.framing, "framing_note": p.framing_note,
            "flip_v": p.flip_v,
            "taken_at": p.taken_at.isoformat(sep=" ", timespec="milliseconds") if p.taken_at else None,
            "thumb": f"/api/thumb/{s.id}/{p.id}",
            "editor": {"dx": round(p.editor.dx_px, 2), "dy": round(p.editor.dy_px, 2),
                       "scale": round(p.editor.scale, 4), "angle": round(p.editor.angle_deg, 3)}}


def _review_json(s):
    slots = {}
    for slot in cfg.ppt.slot_names:
        pid = s.slots.get(slot)
        slots[slot] = _photo_json(s, _photo(s, pid)) if pid else None
    face = [_photo_json(s, _photo(s, pid)) for pid in s.face]
    bins = {k: [_photo_json(s, _photo(s, pid)) for pid in s.bins.get(k, [])]
            for k in list(cfg.ppt.slot_names) + ["FACE"]}
    return {"mode": s.mode, "visit": s.visit, "slots": slots, "face": face, "bins": bins,
            "face_slots": _face_slots_json(s),
            "face_editors": _face_editors_json(s),
            "face_framing": dict(s.face_framing),
            "missing": [sl for sl in cfg.ppt.slot_names if sl not in s.slots]}


# 정적 파일 (프론트엔드)

# ── 업데이트 · 가중치 ─────────────────────────────────────────────────────────
# 터미널을 안 여는 사람이 쓴다. 버그를 고쳐도 상대 컴퓨터에 안 들어가면 고친 게 아니다.


def _busy() -> bool:
    """확정하지 않은 작업이 있나. 있으면 재시작이 그 작업을 날린다."""
    return any(getattr(s, "photos", None) for s in SESSIONS.values())


def _safe_check() -> Up.UpdateStatus:
    """확인이 어떻게 터지든 **500 을 내지 않는다.**

    한 번 이렇게 무너진 적이 있다. `git log` 출력을 CP949 로 디코딩하다 죽어
    500 이 났고, 화면의 `.catch(() => null)` 이 그걸 삼켜 배너가 아예 안 떴다.
    사용자에게는 '최신입니다'와 구분되지 않았다 — 업데이트가 조용히 멈춘 것이다.

    실패는 반드시 **사유를 달고** 화면까지 가야 한다.
    """
    try:
        return Up.check(busy=_busy())
    except Exception as e:                                        # noqa: BLE001
        st = Up.UpdateStatus()
        st.reason = f"확인 중 오류: {type(e).__name__}: {e}"[:300]
        return st


@app.get("/api/update/check")
def update_check():
    """새 버전이 있나. **네트워크를 쓴다** — 화면은 배너로만 쓰고 기다리게 하지 말 것."""
    return _safe_check().to_json()


class UpdateApplyReq(BaseModel):
    force: bool = False      # 직접 수정한 파일을 백업하고 강제 진행


@app.post("/api/update/apply")
def update_apply(req: UpdateApplyReq = Body(default=UpdateApplyReq())):
    st = _safe_check()
    if not st.has_update:
        return {"ok": False, "detail": st.reason or "이미 최신입니다"}
    # 확정하지 않은 작업(busy) 차단은 강제로도 못 넘는다 — 작업이 날아간다.
    if st.blocked and not (req.force and "직접 수정" in st.blocked):
        return {"ok": False, "detail": st.blocked}
    try:
        return Up.apply_update(force=req.force)
    except Exception as e:                                        # noqa: BLE001
        return {"ok": False, "detail": f"{type(e).__name__}: {e}"[:300]}


@app.post("/api/shortcut")
def make_shortcut():
    """바탕화면에 CRoCs 바로가기 — 설치 이후에도 설정에서 만들 수 있게.

    install.bat 의 같은 스텝은 설치/재설치 때만 돌아서, 이미 설치한 사람은
    이 버튼이 유일한 경로다. Windows 전용 (WSL 개발 환경의 powershell.exe 도
    허용 — 동작 검증용).
    """
    exe = "powershell" if os.name == "nt" else shutil.which("powershell.exe")
    if not exe:
        return {"ok": False, "detail": "Windows에서만 만들 수 있습니다"}
    repo = str(BACKEND_DIR.parent.parent)
    if os.name != "nt":                        # WSL — 경로만 Windows 식으로
        r = subprocess.run(["wslpath", "-w", repo], capture_output=True,
                           encoding="utf-8")
        if r.returncode == 0:
            repo = r.stdout.strip()
    ps = (
        "$sh=New-Object -ComObject WScript.Shell;"
        "$d=$sh.SpecialFolders.Item('Desktop');"
        "if(-not $d){$d=Join-Path $env:USERPROFILE 'Desktop'};"
        "$s=$sh.CreateShortcut((Join-Path $d 'CRoCs.lnk'));"
        f"$s.TargetPath='{repo}\\run.bat';"
        f"$s.WorkingDirectory='{repo}';"
        f"$s.IconLocation='{repo}\\assets\\crocs.ico,0';"
        "$s.Save(); Write-Output $d"
    )
    try:
        r = subprocess.run([exe, "-NoProfile", "-ExecutionPolicy", "Bypass",
                            "-Command", ps],
                           capture_output=True, encoding="utf-8",
                           errors="replace", timeout=30)
    except Exception as e:                                        # noqa: BLE001
        return {"ok": False, "detail": f"{type(e).__name__}: {e}"[:200]}
    if r.returncode != 0:
        return {"ok": False,
                "detail": ((r.stderr or r.stdout or "").strip()[:300]
                           or "만들지 못했습니다")}
    lines = [l for l in (r.stdout or "").splitlines() if l.strip()]
    return {"ok": True, "desktop": lines[-1] if lines else ""}


@app.post("/api/update/rollback")
def update_rollback():
    if _busy():
        return {"ok": False, "detail": "확정하지 않은 작업이 있습니다"}
    try:
        return Up.rollback()
    except Exception as e:                                        # noqa: BLE001
        return {"ok": False, "detail": f"{type(e).__name__}: {e}"[:300]}


@app.post("/api/update/restart")
def update_restart():
    """재시작 종료코드로 죽는다. `run.bat`/`run.command` 의 루프가 다시 띄운다.

    파이썬은 자기 자신을 바꿔치울 수 없다 — 이미 메모리에 올라간 모듈은 `git pull`
    뒤에도 그대로다. 프로세스를 새로 띄우는 수밖에 없다.
    """
    if _busy():
        return {"ok": False, "detail": "확정하지 않은 작업이 있습니다"}
    threading.Timer(0.5, Up.restart_now).start()      # 응답을 먼저 보내고 죽는다
    return {"ok": True}


class PrefsReq(BaseModel):
    months_unit: str | None = None
    save_raw: bool | None = None
    # 형식 **목록** — 첫 번째가 생성용, 전부가 인식용. 빈 목록 = 기본만.
    folder_patterns: list[str] | None = None
    ppt_patterns: list[str] | None = None
    # 새 PPT 라벨 표기 — "tight"(YY.MM.DD (초진 A)) | "spaced"(YY. MM. DD. (초진 A))
    label_format: str | None = None
    # 직전 차수 슬라이드의 도형 복사 — "none" | "lines" | "all"
    copy_shapes: str | None = None
    note_sizes: dict[str, float] | None = None   # 좌상단 s/p·좌하단·우하단 글자 크기
    letterbox_color: str | None = None           # 회전·축소로 드러나는 빈 자리 색


def _saved_patterns(key: str = "folder_patterns") -> list[str]:
    try:
        d = json.loads(SETTINGS_FILE.read_text(encoding="utf-8"))
        return [p for p in d.get(key, []) if isinstance(p, str) and p]
    except Exception:                                   # noqa: BLE001
        return []


def _ppt_pattern() -> str:
    """새 PPT 파일 이름 형식 — 목록의 첫 번째, 없으면 config 기본."""
    saved = _saved_patterns("ppt_patterns")
    return saved[0] if saved else cfg.naming.ppt_pattern


def _ppt_patterns() -> list[str]:
    """PPT 파일명 파싱에 시도할 형식들 — 등록 목록이 전부다.

    목록이 비어 있을 때만 config 기본형이 안전망으로 쓰인다. 예전에는 기본형을
    항상 뒤에 붙였는데, 그러면 병원번호가 든 기본형을 지울 길이 없다
    (2026-08-13 결정). 옛 형식 파일을 계속 읽으려면 그 형식을 목록에 남겨 둔다.
    """
    out = []
    for p in (_saved_patterns("ppt_patterns")
              or [cfg.naming.ppt_pattern, *cfg.naming.ppt_patterns_legacy]):
        # 원형 + 생성형 변형 — 순번 같은 인식 전용 블록이 든 형식으로 "없는 셈
        # 치고" 만든 파일이나 병원번호 없이 만든 파일도 같은 형식으로 읽힌다.
        for q in (p, N.strip_recognition(p),
                  N.strip_recognition(p, {"hospital_id"})):
            if q and q not in out:
                out.append(q)
    return out


# 새 PPT 라벨 표기 두 가지 — 설정에서 고른다. **새 PPT 를 만들 때만** 쓰인다.
_BASE_FP = {"paren": True, "paren_space": True, "letter_space": True,
            "has_letter": True}
_FMT_FP = {
    "tight":  {**_BASE_FP, "spaced": False, "trailing_dot": False},
    "spaced": {**_BASE_FP, "spaced": True,  "trailing_dot": True},
}


def _label_format() -> str:
    """새 PPT 라벨 표기 — "tight"(YY.MM.DD) 또는 "spaced"(YY. MM. DD.)."""
    try:
        v = json.loads(SETTINGS_FILE.read_text(encoding="utf-8")).get("label_format")
        return v if v in _FMT_FP else "tight"
    except Exception:                                   # noqa: BLE001
        return "tight"


def _render_label(date_str: str, visit: str, fp: dict | None = None) -> str:
    """차수 라벨 한 벌.

    새 PPT: 설정의 두 표기 중 하나 — **새 PPT 에만 적용**. 기존 PPT 이어쓰기:
    그 덱 마지막 십자뷰 라벨의 표기 지문(fp)을 따른다 — 한 덱 안에서 표기가
    섞이지 않게. 덱에 차수 글자가 없으면 글자 없이 쓴다 — 글자를 붙이면 기존
    무글자 슬라이드들이 "글자 차수 이전"이 되어 자동 부여에서 통째로 빠진다.
    지문이 없으면(라벨을 못 읽은 덱) 설정 표기로 물러난다.
    """
    fp = fp or _FMT_FP[_label_format()]
    nums = re.findall(r"\d+", date_str)
    if len(nums) >= 3:
        sep = ". " if fp.get("spaced") else "."
        date = (sep.join(f"{int(x):02d}" for x in nums[:3])
                + ("." if fp.get("trailing_dot") else ""))
    else:
        date = date_str
    kind = "초진" if visit == "A" else "재진"
    # 띄어쓰기까지 그 덱을 따른다 — "18(재진 C)" 인 덱에 "18 (재진 D)" 를 적으면
    # 같은 슬라이드 묶음 안에서 표기가 갈린다.
    inner = (f"{kind}{' ' if fp.get('letter_space', True) else ''}{visit}"
             if visit and fp.get("has_letter", True) else kind)
    gap = " " if fp.get("paren_space", True) else ""
    return (f"{date}{gap}({inner})" if fp.get("paren", True)
            else f"{date}{gap}{inner}")


def _parse_ppt_name(name: str):
    """PPT 파일명 → Identifiers. 등록된 형식을 차례로 시도한다."""
    last = None
    for pat in _ppt_patterns():
        try:
            return N.parse_ppt_filename(
                name, pat,
                hospital_digits=cfg.identifiers.hospital_id.digits,
                ortho_digits=cfg.identifiers.ortho_id.digits,
                name_regex=cfg.identifiers.name.allow_regex)
        except N.NamingError as e:
            last = e
    raise last


def _remembered_ppt(folder: str) -> str:
    """이 환자에게 마지막으로 이어붙인 PPT (환자 폴더 기준 상대경로).

    열쇠는 폴더 이름이다 — 맥에서 읽은 이름은 자모가 분해돼 있어(NFD) 그대로
    비교하면 방금 적어 둔 값도 못 찾는다. 넣을 때도 찾을 때도 조합형으로 맞춘다.
    """
    try:
        d = json.loads(SETTINGS_FILE.read_text(encoding="utf-8")).get("ppt_choice") or {}
        return str(d.get(N.nfc(folder)) or "")
    except Exception:                                   # noqa: BLE001
        return ""


def _remember_ppt(folder: str, rel: str) -> None:
    """확정한 뒤 그 파일을 기억해 둔다 — 다음에도 같은 덱으로 이어진다."""
    try:
        d = json.loads(SETTINGS_FILE.read_text(encoding="utf-8"))
    except Exception:                                   # noqa: BLE001
        d = {}
    folder, rel = N.nfc(folder), N.nfc(rel)
    choice = d.get("ppt_choice")
    if not isinstance(choice, dict):
        choice = {}
    if choice.get(folder) == rel:
        return
    choice[folder] = rel
    d["ppt_choice"] = choice
    try:
        SETTINGS_FILE.parent.mkdir(parents=True, exist_ok=True)
        SETTINGS_FILE.write_text(json.dumps(d, ensure_ascii=False, indent=2),
                                 encoding="utf-8")
    except OSError:
        pass


def _find_ppt(entries: list[Path], base: Path, ids) -> Path | None:
    """이 환자의 PPT 파일 고르기. 후보가 여럿이면 아래 순서로 하나를 정한다.

    ① **마지막에 이어붙인 그 파일** — 한 환자의 기록이 두 덱으로 갈리지 않게,
       확정할 때 기억해 둔 것을 가장 먼저 쓴다.
    ② **환자 폴더 바로 아래**에 있는 것 — 하위 폴더로 치운 것은 대개 보관본이다.
    ③ 생성 형식과 이름이 같은 것.
    ④ 그래도 남으면 이름 오름차순.

    후보는 `.pptx` 중 등록된 이름 형식으로 읽히고 교정번호가 이 환자인 것뿐이다.
    상대경로를 쓰므로 하위 폴더에 같은 이름이 있어도 서로 다른 후보로 다룬다.
    """
    remembered = N.nfc(_remembered_ppt(base.name))
    gen = N.nfc(_gen_ppt_name(ids))
    best = None
    for p in entries:
        if p.suffix.lower() != ".pptx" or p.name.startswith("~$"):
            continue
        try:
            got = _parse_ppt_name(p.name)
        except N.NamingError:
            continue
        if got.ortho_id != ids.ortho_id:
            continue
        rel = p.relative_to(base).as_posix()
        key = N.nfc(rel)
        rank = (0 if key == remembered else 1,      # ① 기억해 둔 것
                len(Path(rel).parts) - 1,           # ② 얕을수록 먼저
                0 if N.nfc(p.name) == gen else 1,   # ③ 생성 형식 이름
                key.lower())                        # ④ 이름순
        if best is None or rank < best[0]:
            best = (rank, p)
    return best[1] if best else None


def _folder_pattern() -> str:
    """새 폴더를 만들 때 쓰는 형식 — 목록의 첫 번째, 없으면 config 기본."""
    saved = _saved_patterns()
    return saved[0] if saved else cfg.naming.folder_pattern


def _folder_patterns() -> list[str]:
    """폴더명 파싱에 시도할 형식들 — 등록 목록이 전부다.

    목록이 비어 있을 때만 config 기본형이 안전망으로 쓰인다. 예전에는 기본형을
    항상 뒤에 붙였는데, 그러면 병원번호가 든 기본형을 지울 길이 없다
    (2026-08-13 결정). 옛 형식으로 만든 폴더를 계속 읽으려면 — 기존 폴더의
    이름을 고쳐 쓰는 선택지는 없다, 의료 기록이다 — 그 형식을 목록에 남겨 둔다.
    """
    out = []
    for p in _saved_patterns() or [cfg.naming.folder_pattern]:
        # 원형 + 생성형 변형 — 순번 없이 만든 폴더·병원번호 없는 폴더도 읽힌다.
        for q in (p, N.strip_recognition(p),
                  N.strip_recognition(p, {"hospital_id"})):
            if q and q not in out:
                out.append(q)
    return out


def _parse_folder(name: str, **kw):
    """폴더명 → Identifiers. 쓸 수 있는 형식을 차례로 시도한다."""
    last = None
    for pat in _folder_patterns():
        try:
            return N.parse_pattern(name, pat, **kw)
        except N.NamingError as e:
            last = e
    raise last


def _gen_folder_name(ids) -> str:
    """새 환자 폴더 이름 — ★ 형식에서 미리 정할 수 없는 블록(순번·*·자릿수
    범위)과, 비어 있으면 병원번호까지 뺀 **생성형**으로 만든다."""
    extra = set() if ids.hospital_id else {"hospital_id"}
    return N.folder_name(ids, N.strip_recognition(_folder_pattern(), extra))


def _gen_ppt_name(ids) -> str:
    """새 PPT 파일 이름 — 폴더와 같은 규칙의 생성형."""
    extra = set() if ids.hospital_id else {"hospital_id"}
    return N.ppt_filename(ids, N.strip_recognition(_ppt_pattern(), extra))


NOTE_SIZE_KEYS = (CD.NOTE_SOAP, CD.NOTE_LL, CD.NOTE_NEXT)


def _note_sizes() -> dict:
    try:
        d = json.loads(SETTINGS_FILE.read_text(encoding="utf-8"))
        return {k: float(v) for k, v in (d.get("note_sizes") or {}).items()
                if k in NOTE_SIZE_KEYS}
    except Exception:                                   # noqa: BLE001
        return {}


def _apply_note_sizes() -> None:
    """설정 크기를 case_deck 상수와 화면 레이아웃(NOTE_BOXES)에 반영한다."""
    sizes = _note_sizes()
    CD.apply_note_sizes(sizes)
    for k, v in sizes.items():
        if k in NOTE_BOXES and v:
            NOTE_BOXES[k]["size_pt"] = float(v)
            NOTE_BOXES[k].setdefault("font", {})["size_pt"] = float(v)


def _prefs_json() -> dict:
    return {"months_unit": _months_unit(), "save_raw": _save_raw(),
            "folder_patterns": _saved_patterns(),
            "folder_pattern_default": cfg.naming.folder_pattern,
            "ppt_patterns": _saved_patterns("ppt_patterns"),
            "ppt_pattern_default": cfg.naming.ppt_pattern,
            "label_format": _label_format(),
            "copy_shapes": _copy_shapes(),
            "note_sizes": {k: CD.NOTE_WINDOWS[k]["size_pt"]
                           for k in NOTE_SIZE_KEYS},
            "letterbox_color": _letterbox_color()}


@app.get("/api/prefs")
def prefs_get():
    return _prefs_json()


@app.post("/api/prefs")
def prefs_set(req: PrefsReq):
    """개인화 설정. **`settings.json`(설치본 공용)** 에 둔다 — 브라우저마다 다르면
    한 환자 PPT 안에서 표기가 섞이고, 저장되는 파일 구성도 그날그날 달라진다.

    보낸 항목만 바꾼다. 화면이 항상 전부를 보내야 한다면 한 항목을 고칠 때마다
    나머지를 실어 나르다가 언젠가 하나를 빠뜨린다.
    """
    if req.months_unit is not None and req.months_unit not in ("int", "half"):
        raise HTTPException(400, "int 또는 half")
    try:
        d = json.loads(SETTINGS_FILE.read_text(encoding="utf-8"))
    except Exception:                                   # noqa: BLE001
        d = {}
    if req.months_unit is not None:
        d["months_unit"] = req.months_unit
    if req.save_raw is not None:
        d["save_raw"] = bool(req.save_raw)
    if req.note_sizes is not None:
        clean = {}
        for k, v in req.note_sizes.items():
            if k not in NOTE_SIZE_KEYS:
                raise HTTPException(400, f"모르는 노트 칸: {k}")
            try:
                v = float(v)
            except (TypeError, ValueError):
                raise HTTPException(400, "숫자를 넣어주세요")
            if not 6 <= v <= 40:
                raise HTTPException(400, "글자 크기는 6~40pt 사이여야 합니다")
            clean[k] = v
        if clean:
            d["note_sizes"] = {**(d.get("note_sizes") or {}), **clean}
    if req.folder_patterns is not None:
        pats = []
        for pat in (p.strip() for p in req.folder_patterns):
            if not pat or pat in pats:
                continue
            if (set('\\/:*?"<>|') & set(pat)) or any(
                    k not in pat for k in ("{name}", "{ortho_id}")):
                raise HTTPException(400, '형식에는 {name} {ortho_id} 가 있어야 하고 '
                                         '(병원번호는 선택) \\ / : * ? " < > | 는 못 씁니다')
            # 역파싱 라운드트립 — 만들 수 있어도 못 읽으면 환자 목록이 깨진다.
            # 인식 전용 블록(순번·*·자릿수 범위)은 생성 때 빠지므로, 실제로
            # 만들어질 모습(생성형)으로 검사한다.
            gen = N.strip_recognition(pat)
            probe = N.Identifiers("홍길동", "1" * cfg.identifiers.hospital_id.digits,
                                  "2" * cfg.identifiers.ortho_id.digits)
            try:
                N.parse_pattern(N.folder_name(probe, gen), gen, label="폴더명",
                                hospital_digits=cfg.identifiers.hospital_id.digits,
                                ortho_digits=cfg.identifiers.ortho_id.digits,
                                name_regex=cfg.identifiers.name.allow_regex)
            except N.NamingError:
                raise HTTPException(400, f"'{pat}' 형식은 폴더명을 다시 읽어낼 수 "
                                         "없습니다 — 항목 사이에 구분 문자를 넣어주세요")
            pats.append(pat)
        if pats:
            d["folder_patterns"] = pats
        else:
            d.pop("folder_patterns", None)      # 빈 목록 = 기본 형식만
    if req.ppt_patterns is not None:
        pats = []
        probe = N.Identifiers("홍길동", "1" * cfg.identifiers.hospital_id.digits,
                              "2" * cfg.identifiers.ortho_id.digits)
        for pat in (p.strip() for p in req.ppt_patterns):
            if not pat:
                continue
            if not pat.lower().endswith(".pptx"):
                pat += ".pptx"                  # 확장자 블록은 없다 — 붙여 준다
            if pat in pats:
                continue
            if set(chr(92) + '/:?"<>|') & set(pat):
                raise HTTPException(400, 'PPT 이름에 \\ / : ? " < > | 는 못 씁니다')
            if "{ortho_id}" not in pat:
                raise HTTPException(400, "형식에는 {ortho_id} 가 있어야 합니다 "
                                         "(이름·병원번호는 선택)")
            gen = N.strip_recognition(pat)
            try:
                N.parse_ppt_filename(
                    N.ppt_filename(probe, gen), gen,
                    hospital_digits=cfg.identifiers.hospital_id.digits,
                    ortho_digits=cfg.identifiers.ortho_id.digits,
                    name_regex=cfg.identifiers.name.allow_regex)
            except (N.NamingError, KeyError):
                raise HTTPException(400, f"'{pat}' 형식은 파일명을 다시 읽어낼 "
                                         "수 없습니다 — 구분 문자를 넣어주세요")
            pats.append(pat)
        if pats:
            d["ppt_patterns"] = pats
        else:
            d.pop("ppt_patterns", None)
    if req.label_format is not None:
        if req.label_format not in _FMT_FP:
            raise HTTPException(400, "tight 또는 spaced")
        d["label_format"] = req.label_format
        d.pop("label_patterns", None)   # 옛 블록 양식 — 더는 쓰지 않는다
    if req.letterbox_color is not None:
        v = req.letterbox_color.strip().lstrip("#")
        if not re.fullmatch(r"[0-9A-Fa-f]{6}", v):
            raise HTTPException(400, "색은 RGB 6자리(예: 000000)여야 합니다")
        d["letterbox_color"] = v.upper()
    if req.copy_shapes is not None:
        if req.copy_shapes not in ("none", "lines", "all"):
            raise HTTPException(400, "none · lines · all 중 하나")
        d["copy_shapes"] = req.copy_shapes
    SETTINGS_FILE.parent.mkdir(parents=True, exist_ok=True)
    SETTINGS_FILE.write_text(json.dumps(d, ensure_ascii=False, indent=2), encoding="utf-8")
    _apply_note_sizes()                        # 노트 글자 크기도 즉시 반영
    return _prefs_json()


def _existing_ppt_names() -> list[str]:
    """환자 폴더들에 실제로 있는 PPT 파일 이름 — 형식 미리보기용.

    같은 이름이 여러 환자 폴더에 있을 이유는 없지만, 있어도 한 줄로 보이면 된다.
    PowerPoint 임시 파일(~$)은 사람이 만든 이름이 아니라 뺀다.
    """
    seen: dict[str, None] = {}
    for d in sorted(ROOT.iterdir(), key=lambda x: x.name.lower()):
        if not d.is_dir() or d.name.startswith((".", "_")):
            continue
        for f in _patient_files(d):
            if f.suffix.lower() == ".pptx" and not f.name.startswith("~$"):
                seen.setdefault(f.name)
    return list(seen)


@app.get("/api/pattern_check")
def pattern_check(pattern: str = "", kind: str = "folder"):
    """조립 중인 형식으로 실제 폴더·PPT 가 읽히는지 미리 본다 (설정 화면).

    ✓ 이 형식으로 인식 · ↩ 다른 등록 형식으로는 인식 · ✗ 어느 형식으로도 못 읽음.
    형식마다 **생성형 변형**(순번 등 인식 전용 블록과 빈 병원번호를 뺀 모습)도
    함께 시도한다 — 그렇게 만들어진 이름도 그 형식 소속이기 때문이다.
    """
    dig = dict(hospital_digits=cfg.identifiers.hospital_id.digits,
               ortho_digits=cfg.identifiers.ortho_id.digits,
               name_regex=cfg.identifiers.name.allow_regex)
    ppt = kind == "ppt"
    if ppt:
        pat = pattern or _ppt_pattern()
        if not pat.lower().endswith(".pptx"):
            pat += ".pptx"                  # 조립 중인 형식엔 확장자가 없다
        names = _existing_ppt_names()
    else:
        pat = pattern or _folder_pattern()
        names = [d.name for d in sorted(ROOT.iterdir(), key=lambda x: x.name.lower())
                 if d.is_dir() and not d.name.startswith((".", "_"))]

    def parse_one(name: str, q: str):
        if ppt:
            return N.parse_ppt_filename(name, q, **dig)
        return N.parse_pattern(name, q, label="폴더명", **dig)

    out = []
    for name in names:
        row = {"name": name, "match": False, "fallback": False}
        for q in (pat, N.strip_recognition(pat),
                  N.strip_recognition(pat, {"hospital_id"})):
            if not q:
                continue
            try:
                parse_one(name, q)
                row["match"] = True
                break
            except (N.NamingError, KeyError):
                pass
        if not row["match"]:
            try:
                _parse_ppt_name(name) if ppt else _parse_folder(name, label="폴더명", **dig)
                row["fallback"] = True      # 옛/기본 형식으로는 읽힌다
            except (N.NamingError, KeyError):
                pass
        out.append(row)
    return {"items": out}


@app.get("/api/uninstall/inventory")
def uninstall_inventory():
    """지워질 것과 남을 것. **환자 자료는 기본으로 남는다** — 의료 기록이다."""
    return Un.inventory(BACKEND_DIR.parents[1], ROOT).to_json()


@app.post("/api/uninstall/prepare")
def uninstall_prepare(body: dict = Body(default={})):
    """삭제 스크립트를 만들고 앱을 끝낸다.

    Windows 는 실행 중인 파일을 잠가서 앱이 자기 폴더를 못 지운다. 스크립트를
    바깥에 만들어 두고, 앱이 끝난 뒤 사용자가 한 번 더 실행한다 — 되돌릴 수 없는
    일이라 확인이 한 번 더 있는 편이 낫다.
    """
    if _busy():
        return {"ok": False, "detail": "확정하지 않은 작업이 있습니다"}
    drop = bool(body.get("drop_data"))
    tools = [str(t) for t in (body.get("drop_tools") or [])]
    if drop and body.get("confirm") != "삭제":
        return {"ok": False, "detail": "환자 자료를 지우려면 확인 문구가 필요합니다"}
    r = Un.prepare(BACKEND_DIR.parents[1], ROOT, drop_data=drop, drop_tools=tools)
    # 종료코드 0 — 재시작 루프를 끝낸다 (42 는 재시작)
    threading.Timer(1.0, lambda: os._exit(0)).start()
    return r


@app.get("/api/weights")
def weights_status():
    """가중치 준비 상태. 없으면 화면이 무엇을 어디서 받는지 안내한다."""
    try:
        import sys as _sys                                          # noqa: PLC0415
        _sys.path.insert(0, str(BACKEND_DIR.parents[1]))
        import weightstore                                          # noqa: PLC0415
        st = weightstore.scan(verify=False)
        return {"ready": st.ready,
                "items": [{"key": i.key, "state": i.state, "file": i.file,
                           "detail": i.detail, "url": i.drive_url} for i in st.items],
                "strays": [{"name": p.name, "why": w} for p, w in st.strays],
                "drop_dir": str(weightstore.DROP_DIR)}
    except Exception as e:                                          # noqa: BLE001
        return {"ready": False, "items": [], "strays": [],
                "error": f"{type(e).__name__}: {e}"}


app.mount("/static", _NoCacheStatic(directory=str(FRONTEND_DIR)), name="static")


def _port_free(port: int) -> bool:
    """그 포트에 바인딩할 수 있나 — 서버가 실제로 뜰 수 있는지와 같은 조건.

    POSIX 에서는 `SO_REUSEADDR` 를 켠다. 방금 죽은 프로세스가 남긴 TIME_WAIT 소켓이
    있으면 이 옵션 없이는 바인딩이 막히는데, 정작 uvicorn 은 이 옵션을 켜고 뜨므로
    "못 쓴다" 는 판정이 거짓이 된다 (좀비를 정리하고도 실패로 봤다).
    Windows 에서는 켜지 않는다 — 거기서는 이 옵션이 **남이 쓰는 포트도 빼앗는**
    뜻이라, 켜면 비어 있지도 않은 포트를 비었다고 보게 된다.
    """
    import socket                                                   # noqa: PLC0415
    with socket.socket() as s:
        if os.name != "nt":
            s.setsockopt(socket.SOL_SOCKET, socket.SO_REUSEADDR, 1)
        try:
            s.bind(("127.0.0.1", port))
            return True
        except OSError:
            return False


def _ours_at(port: int) -> bool:
    """그 포트에 떠 있는 것이 **살아 있는 이 앱**인가.

    두 번, 넉넉한 시간을 두고 물어본다. 한 번 놓치고 좀비로 판정하면 멀쩡히 돌던
    실행을 죽이게 되는데, 그 안에는 아직 확정하지 않은 작업(담아둔 사진·조정한
    구도·입력한 노트)이 통째로 들어 있다. 되살릴 수 없는 손실이라 판정을 느슨하게
    잡는 편이 맞다.
    """
    import urllib.request                                           # noqa: PLC0415
    for attempt in range(2):
        try:
            with urllib.request.urlopen(
                    f"http://127.0.0.1:{port}/api/health", timeout=4.0) as r:
                return bool(json.loads(r.read()).get("program_dir"))
        except Exception:                                           # noqa: BLE001
            if attempt == 0:
                time.sleep(0.6)
    return False


# 실행 중인 서버가 자기 pid 와 포트를 적어 두는 자리. 새로 뜨는 쪽이 "포트를 쥔
# 것이 우리 것인가" 를 확인하는 데 쓴다 — 이게 없으면 남의 프로그램을 죽일 수 있다.
LOCK_FILE = PROGRAM_DIR / ".server.json"


def _write_lock(port: int) -> None:
    try:
        LOCK_FILE.write_text(json.dumps({"pid": os.getpid(), "port": port}),
                             encoding="utf-8")
    except OSError:
        pass


def _kill_stale(port: int) -> bool:
    """응답하지 않으면서 포트만 쥐고 있는 **우리 프로세스**를 정리한다.

    창을 닫았는데 프로세스만 남았거나 멎어 버린 경우다. 죽일 대상은 우리가 적어
    둔 pid 와 포트가 모두 맞을 때뿐이다 — 포트만 보고 죽이면 하필 그 포트를 쓰던
    남의 프로그램을 끄게 된다.
    """
    import signal                                                   # noqa: PLC0415
    try:
        d = json.loads(LOCK_FILE.read_text(encoding="utf-8"))
    except Exception:                                               # noqa: BLE001
        return False
    pid = d.get("pid")
    if not isinstance(pid, int) or d.get("port") != port or pid == os.getpid():
        return False
    try:
        os.kill(pid, signal.SIGTERM)
    except OSError:
        return False                    # 이미 없는 pid — 포트는 남이 쥐고 있다
    for _ in range(20):                 # 최대 2초 기다린다
        time.sleep(0.1)
        if _port_free(port):
            return True
    try:                                # 안 죽으면 강제로 (Windows 는 위가 이미 강제)
        os.kill(pid, getattr(signal, "SIGKILL", signal.SIGTERM))
    except OSError:
        pass
    for _ in range(20):                 # 포트가 풀릴 때까지 다시 2초
        time.sleep(0.1)
        if _port_free(port):
            return True
    return False


def run():
    import threading
    import webbrowser
    import uvicorn
    port = int(os.environ.get("PORT", "8000"))
    # 이미 쓰이는 포트면 서버가 뜨지 못하고 창이 그냥 닫힌다 — 사용자 눈에는
    # "실행했는데 아무 일도 안 일어남" 이다. 그래서 먼저 살펴본다.
    if not _port_free(port):
        if _ours_at(port):
            # 앞서 띄운 것이 **살아서** 돌고 있다. 두 번째를 띄울 이유가 없고,
            # 죽여서도 안 된다 — 확정 전 작업이 그 안에 있다. 브라우저만 연다.
            print(f"[알림] 이미 실행 중입니다 (포트 {port}). 브라우저만 엽니다.")
            webbrowser.open(f"http://127.0.0.1:{port}/")
            return
        if _kill_stale(port):
            # 창은 닫혔는데 프로세스만 남아 포트를 쥐고 있던 경우 — 정리하고 뜬다.
            print(f"[알림] 응답하지 않는 이전 실행을 정리하고 시작합니다 (포트 {port}).")
        else:
            # 남이 쓰는 포트다. **옆 포트로 물러나지 않는다** — 주소가 늘 같아야
            # 즐겨찾기·바로가기가 계속 맞는다. 남의 프로그램을 끌 수도 없으니
            # 무엇을 하라는 것인지 적고 멈춘다.
            print(f"[오류] 포트 {port} 을 다른 프로그램이 쓰고 있습니다.")
            print("       그 프로그램을 닫고 다시 실행해 주세요.")
            print(f"       확인:  netstat -ano | findstr :{port}")
            print(f"       이번만 다른 포트로 띄우려면:  set PORT={port + 1} && run.bat")
            try:
                input("       Enter 를 누르면 닫습니다... ")
            except EOFError:
                pass
            return
    _write_lock(port)
    threading.Timer(1.2, lambda: webbrowser.open(f"http://127.0.0.1:{port}/")).start()
    uvicorn.run(app, host="127.0.0.1", port=port, log_level="info")


if __name__ == "__main__":
    run()
