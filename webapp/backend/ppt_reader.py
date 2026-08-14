"""
기존 PPT에서 기준영상 재구성 (Stage 4b, 사양 §5.1.1)

각 슬라이드의 슬롯별 사진 도형에서 변환값(a:xfrm의 rot/off/ext)과 원본 이미지를
읽어, "창(8.4×6.3cm)에서 실제로 보였던 화면"을 픽셀로 복원한다.
INFO_BOX 텍스트로 각 슬라이드의 차수를 식별한다.

복원된 창 이미지가 정합(registration)의 기준영상이 된다.
좌표 규약: PPT rot 양수 = 시계방향, 화면 픽셀은 y-아래. 회전행렬
R=[[c,-s],[s,c]]가 y-아래에서 시계방향과 일치(ppt_writer/coords와 동일 규약).
"""

from __future__ import annotations

import io
import math
import re
from dataclasses import dataclass, field

import cv2
import numpy as np
from PIL import Image

from lxml import etree
from pptx.oxml.ns import qn

import naming

from coords import WindowCm, emu_to_cm
from ppt_writer import PHOTO_NAME_PREFIX
from template import find_shape, shape_window_cm


@dataclass
class SlotRef:
    slot_name: str
    image: np.ndarray            # 복원된 창 픽셀 (BGR)
    off_cm: tuple[float, float]
    ext_cm: tuple[float, float]
    rot_deg: float
    flip_v: bool = False         # a:xfrm/@flipV — image 에는 이미 반영돼 있다


@dataclass
class VisitSlide:
    slide_index: int
    visit: str | None            # 차수 알파벳 (A, B, ...) 또는 None
    date: str | None             # "YY.MM.DD"
    kind: str                    # "first"|"revisit"|"unknown"
    label_text: str = ""         # 라벨 원문 — 표기 지문(label_fingerprint)용
    slots: dict[str, SlotRef] = field(default_factory=dict)
    # 우상단 상태 칸의 글. `Tx./App. Period` 의 기준일이 괄호로 여기 적혀 있어서,
    # 재진 때 그 날짜를 되읽는다 (`main._period_start`).
    status_text: str = ""


# ── INFO_BOX 파싱 ─────────────────────────────────────────────────────────────
# 사람이 적은 라벨은 형태가 흔들린다 — "24. 07. 18(재진)" 처럼 점 뒤 공백이
# 들어가거나 재진에 차수 글자가 없기도 하다. 날짜는 공백을 허용해 읽고
# YY.MM.DD 로 정규화하며, 재진의 차수 글자는 없어도 재진으로 인정한다.
_DATE = re.compile(r"(\d{2})\s*\.\s*(\d{1,2})\s*\.\s*(\d{1,2})")
_REVISIT = re.compile(r"재진\s*\(?([A-Z]+)?")
_FIRST = re.compile(r"초진")


def _norm_date(raw: str | None) -> str | None:
    """"26. 1. 5" · "26.01.05" → "26.01.05". 표기가 달라도 같은 날짜는 같은 글."""
    nums = re.findall(r"\d+", raw or "")
    if len(nums) < 3:
        return None
    y, mo, d = (int(x) for x in nums[:3])
    return f"{y:02d}.{mo:02d}.{d:02d}"


def parse_info_box(text: str) -> tuple[str | None, str | None, str]:
    """INFO_BOX 텍스트 → (visit, date, kind). visit 은 글자가 없으면 None.

    파서는 이것 하나다 — 표기는 점 뒤 공백·끝점·차수 글자 유무 정도만 흔들린다는
    전제로, 등록 양식 없이 관대하게 읽는다 (2026-08-13 결정).
    """
    date = None
    m = _DATE.search(text or "")
    if m:
        date = _norm_date(".".join(m.groups()))
    if _FIRST.search(text or ""):
        return "A", date, "first"
    m = _REVISIT.search(text or "")
    if m:
        return m.group(1), date, "revisit"
    return None, date, "unknown"


def label_fingerprint(text: str) -> dict | None:
    """라벨 원문의 표기 지문 — 새 라벨이 그 덱의 모양을 그대로 따라 쓰게 한다.

    한 덱 안에서도 손으로 적은 라벨은 띄어쓰기가 흔들린다("18(재진 C)" · "18 (재진F)").
    큰 틀(날짜 + 괄호 안 초진/재진 + 차수 글자)은 그대로이므로, 흔들리는 다섯 곳만
    본다: 점 뒤 공백 · 날짜 끝 마침표 · 괄호 유무 · 괄호 앞 공백 · 차수 글자 앞 공백
    (그리고 글자가 있는지). 마지막 차수 슬라이드를 기준으로 삼는다.
    """
    m = re.search(r"(\d{2})\s*\.(\s*)\d{1,2}\s*\.\s*(\d{1,2})(\s*\.)?", text or "")
    if not m:
        return None
    # 날짜 뒤 — 괄호(있다면)와 초진/재진, 그리고 차수 글자까지의 사이 공백
    k = re.search(r"(\s*)(\()?\s*(?:초진|재진)(\s*)([A-Z]+)?", text or "")
    return {
        "spaced": bool(m.group(2)),
        "trailing_dot": bool(m.group(4)),
        "paren": bool(k.group(2)) if k else True,
        "paren_space": bool(k.group(1)) if k else True,
        "letter_space": bool(k.group(3)) if k else True,
        "has_letter": bool(k.group(4)) if k else False,
    }


# ── 픽셀 복원 ─────────────────────────────────────────────────────────────────
def _pic_image_bgr(pic) -> np.ndarray:
    blob = pic.image.blob
    im = Image.open(io.BytesIO(blob)).convert("RGB")
    return cv2.cvtColor(np.asarray(im), cv2.COLOR_RGB2BGR)


def read_flip_v(pic) -> bool:
    """도형에 상하반전(`a:xfrm/@flipV`)이 걸려 있는가.

    이 앱이 쓴 PPT 는 교합면을 **원본 픽셀 + flipV** 로 저장한다. 반면 예전
    수작업으로 만든 PPT 는 이미 뒤집힌 픽셀을 그냥 넣었을 수 있는데, 그때는
    이 속성이 없으므로 복원 결과가 그대로 뒤집힌 상태가 된다. 어느 쪽이든
    "화면에서 보였던 그림"이 나오므로 정합 기준으로 바로 쓸 수 있다.
    """
    xfrm = getattr(pic._element.spPr, "xfrm", None)
    return xfrm is not None and str(xfrm.get("flipV", "")).lower() in ("1", "true")


def reconstruct_window(pic, window: WindowCm, px_per_cm: float,
                       canvas: np.ndarray | None = None) -> SlotRef:
    """
    사진 도형의 배치(off/ext/rot/flipV)를 이용해 창에서 보였던 픽셀을 복원.
    photo 픽셀 → window 캔버스 픽셀 아핀을 만들어 warpAffine.
    """
    img = _pic_image_bgr(pic)
    ph, pw = img.shape[:2]
    off_x, off_y = emu_to_cm(pic.left), emu_to_cm(pic.top)
    ext_w, ext_h = emu_to_cm(pic.width), emu_to_cm(pic.height)
    rot = float(pic.rotation or 0.0)  # 도, 시계방향(+)
    flip_v = read_flip_v(pic)

    # 1) photo px → box-local cm (중심 기준)
    #    u_cm = (px/pw*ext_w - ext_w/2, py/ph*ext_h - ext_h/2)
    # 2) 회전(시계, y-아래): R=[[c,-s],[s,c]]
    # 3) + center_cm(off+ext/2) → slide cm
    # 4) slide cm → window px: (cm - win_origin) * px_per_cm
    th = math.radians(rot)
    c, s = math.cos(th), math.sin(th)
    cx_cm, cy_cm = off_x + ext_w / 2.0, off_y + ext_h / 2.0

    # photo px → slide cm  (아핀 2x3)
    #   scale: sxx = ext_w/pw, syy = ext_h/ph
    sxx, syy = ext_w / pw, ext_h / ph
    # local_cm = (px*sxx - ext_w/2, py*syy - ext_h/2)
    # rot·local + center:
    #   X = c*lx - s*ly + cx ;  Y = s*lx + c*ly + cy
    # lx = px*sxx - ext_w/2 ; ly = py*syy - ext_h/2
    #
    # flipV 는 bbox 안에서 y 를 뒤집는다(회전보다 먼저):
    #   ly = (ph − py)*syy − ext_h/2 = −py*syy + ext_h/2
    # 즉 syy 의 부호를 뒤집고 상수항을 +ext_h/2 로 바꾸면 된다.
    fy = -1.0 if flip_v else 1.0
    syy_f = fy * syy
    ly0 = -fy * ext_h / 2
    a11 = c * sxx
    a12 = -s * syy_f
    a13 = c * (-ext_w / 2) - s * ly0 + cx_cm
    a21 = s * sxx
    a22 = c * syy_f
    a23 = s * (-ext_w / 2) + c * ly0 + cy_cm
    # slide cm → window px: (X - win.x)*ppc, (Y - win.y)*ppc
    ppc = px_per_cm
    A = np.array([
        [a11 * ppc, a12 * ppc, (a13 - window.x) * ppc],
        [a21 * ppc, a22 * ppc, (a23 - window.y) * ppc],
    ], dtype=np.float32)

    win_w_px = max(1, int(round(window.w * px_per_cm)))
    win_h_px = max(1, int(round(window.h * px_per_cm)))
    if canvas is not None:
        # 슬라이드 미리보기 합성용 — 기존 캔버스 위에 이 도형만 덧그린다.
        cv2.warpAffine(img, A, (canvas.shape[1], canvas.shape[0]), dst=canvas,
                       flags=cv2.INTER_LINEAR, borderMode=cv2.BORDER_TRANSPARENT)
        out = canvas
    else:
        out = cv2.warpAffine(img, A, (win_w_px, win_h_px),
                             flags=cv2.INTER_LINEAR, borderValue=(0, 0, 0))
    return SlotRef(
        slot_name="", image=out,
        off_cm=(off_x, off_y), ext_cm=(ext_w, ext_h), rot_deg=rot, flip_v=flip_v,
    )


# ── 슬라이드/프레젠테이션 순회 ────────────────────────────────────────────────
# 케이스 덱(초진)의 십자뷰에는 INFO_BOX 가 없다 — 날짜/차수는 양식의 날짜 칸에
# 적힌다. 이름을 여기 적어 두어 두 서식 어느 쪽이든 차수를 읽어낼 수 있게 한다.
# (case_deck.NOTE_DATE 와 같은 값. 임포트하면 순환 참조가 되므로 문자열로 둔다.)
_NOTE_DATE_BOX = "NOTE_DATE"
_NOTE_STATUS_BOX = "NOTE_STATUS"


def _visit_text(slide, cfg) -> str:
    """차수·날짜가 적힌 상자의 글.

    INFO_BOX → 날짜 칸 순으로 **이름**으로 찾고, 둘 다 없으면(수제 PPT) 슬라이드의
    모든 텍스트 박스를 훑어 날짜+초진/재진 라벨이 든 첫 박스를 쓴다.
    """
    for name in (cfg.ppt.info_box_name, _NOTE_DATE_BOX):
        sh = find_shape(slide, name)
        if sh is not None and sh.has_text_frame and sh.text_frame.text.strip():
            return sh.text_frame.text
    for sh in slide.shapes:
        if not getattr(sh, "has_text_frame", False):
            continue
        t = sh.text_frame.text or ""
        _v, dt, kind = parse_info_box(t)
        if dt and kind != "unknown":
            return t
    return ""


# 기간 칸을 글로 알아볼 때 쓰는 머리말. main.PERIOD_KEYS 와 같은 문자열이다
# (여기서 main 을 임포트하면 순환이 된다).
_PERIOD_MARKS = ("Tx. Period", "Rx. Period", "App. Period")


def _status_text(slide) -> str:
    """기간 칸(Tx/Rx/App)의 글.

    이름(NOTE_STATUS)으로 먼저 찾고, 없으면 **글 내용으로** 찾는다 — 라벨과 같은
    방식이다. 수제 PPT 의 상자에는 이름 규약이 없어서, 이 폴백이 없으면 기준일
    이력이 통째로 비고 기간이 차수마다 0 month 로 다시 시작한다.
    """
    sh = find_shape(slide, _NOTE_STATUS_BOX)
    if sh is not None and getattr(sh, "has_text_frame", False) \
            and (sh.text_frame.text or "").strip():
        return sh.text_frame.text
    for sh in slide.shapes:
        if not getattr(sh, "has_text_frame", False):
            continue
        t = sh.text_frame.text or ""
        if any(m in t for m in _PERIOD_MARKS):
            return t
    return ""


def read_visit_slide(slide, slide_index: int, cfg, px_per_cm: float,
                     windows: dict[str, WindowCm] | None = None,
                     slide_ctr: tuple[float, float] | None = None) -> VisitSlide:
    text = _visit_text(slide, cfg)
    visit, date, kind = parse_info_box(text)
    vs = VisitSlide(slide_index=slide_index, visit=visit, date=date, kind=kind,
                    label_text=text)
    vs.status_text = _status_text(slide)

    # 슬롯별 사진 도형 (PHOTO_SLOT_x). 창 좌표는 config 슬롯 좌표 사용
    # (앵커가 삭제되었으므로 템플릿 슬롯 기준을 쓴다 — ppt_reader.slot_windows 주입).
    for pic in slide.shapes:
        if not pic.name.startswith(PHOTO_NAME_PREFIX):
            continue
        slot_name = pic.name[len(PHOTO_NAME_PREFIX):]
        window = (windows or _SLOT_WINDOWS).get(slot_name)
        if window is None:
            continue
        ref = reconstruct_window(pic, window, px_per_cm)
        ref.slot_name = slot_name
        vs.slots[slot_name] = ref

    # 수제 PPT 폴백 — 이름 규약이 없어도, 차수 라벨이 있는 슬라이드라면 사진
    # 배치로 슬롯을 추정한다: 무게중심에 가장 가까운 사진 = 정면(FRONT),
    # 나머지는 정면 기준 상/하/좌/우. 절대 좌표를 안 쓰므로 사진이 몇 mm
    # 어긋나 있어도 십자 형태만 유지되면 맞는다. 복원 창은 각 사진의 실제
    # 중심에 슬롯 창 크기를 씌워 픽셀 스케일(px_per_cm)을 앱 규약과 맞춘다.
    if not vs.slots and (vs.date or vs.kind != "unknown"):
        # 유효 십자뷰 기준: **가로 8cm 이상 사진이 5장** (2026-08-12 결정).
        # 로고·썸네일 같은 작은 그림이 섞여 있어도 판정이 안 흔들린다.
        pics = [sh for sh in slide.shapes
                if getattr(sh, "shape_type", None) == 13
                and emu_to_cm(sh.width) >= 8.0]
        if len(pics) >= 5:
            ctr = [(emu_to_cm(p.left) + emu_to_cm(p.width) / 2,
                    emu_to_cm(p.top) + emu_to_cm(p.height) / 2) for p in pics]
            # 정면 = 슬라이드 중앙에 가장 가까운 사진 — 실제 덱들은 십자를
            # 중앙 근처에 두므로, 모서리에 낀 추가 사진에 안 끌려가는 기준이다.
            # 중앙을 모르면(호출부가 안 줬으면) 사진 중앙값으로 물러난다.
            if slide_ctr is not None:
                gx, gy = slide_ctr
            else:
                xs, ys = sorted(c[0] for c in ctr), sorted(c[1] for c in ctr)
                gx, gy = xs[len(xs) // 2], ys[len(ys) // 2]
            fi = min(range(len(pics)),
                     key=lambda k: (ctr[k][0] - gx) ** 2 + (ctr[k][1] - gy) ** 2)
            fx, fy = ctr[fi]
            picked = {"SLOT_FRONT": (fi, 0.0)}
            for k in range(len(pics)):
                if k == fi:
                    continue
                dx, dy = ctr[k][0] - fx, ctr[k][1] - fy
                slot = (("SLOT_UPPER" if dy < 0 else "SLOT_LOWER")
                        if abs(dy) >= abs(dx)
                        else ("SLOT_LEFT" if dx < 0 else "SLOT_RIGHT"))
                d2 = dx * dx + dy * dy
                if slot not in picked or d2 < picked[slot][1]:
                    picked[slot] = (k, d2)
            for slot, (k, _d) in picked.items():
                w0 = (windows or _SLOT_WINDOWS).get(slot)
                dw, dh = (w0.w, w0.h) if w0 else (8.4, 6.3)
                cx, cy = ctr[k]
                win = WindowCm(x=cx - dw / 2, y=cy - dh / 2, w=dw, h=dh)
                ref = reconstruct_window(pics[k], win, px_per_cm)
                ref.slot_name = slot
                vs.slots[slot] = ref
    return vs


# 슬롯 창 좌표는 템플릿에서 1회 주입해 캐시(앵커 삭제된 슬라이드도 복원 가능)
_SLOT_WINDOWS: dict[str, WindowCm] = {}


def set_slot_windows(windows: dict[str, WindowCm]) -> None:
    _SLOT_WINDOWS.clear()
    _SLOT_WINDOWS.update(windows)


def read_all_visits(prs, cfg, px_per_cm: float,
                    windows: dict[str, WindowCm] | None = None) -> list[VisitSlide]:
    """
    PPT의 모든 슬라이드에서 차수·슬롯 기준영상을 복원.

    windows를 주면 그 창으로 복원한다 — 그 PPT가 실제로 쓰던 레이아웃이다.
    주지 않으면 템플릿 창(전역 캐시)을 쓴다.
    """
    out = []
    ctr = (emu_to_cm(prs.slide_width) / 2, emu_to_cm(prs.slide_height) / 2)
    for i, slide in enumerate(prs.slides):
        vs = read_visit_slide(slide, i, cfg, px_per_cm, windows, slide_ctr=ctr)
        if vs.slots:  # 구내 사진이 있는 슬라이드만
            out.append(vs)
    assign_letterless_visits(out)   # "(재진)"처럼 글자 없는 라벨 — 날짜순으로 부여
    return out


def _date_key(d: str | None) -> tuple[int, ...] | None:
    """"24.09.26" → (24, 9, 26). 비교용 — 자릿수 표기가 달라도 같은 날짜다."""
    nums = re.findall(r"\d+", d or "")
    return tuple(int(x) for x in nums[:3]) if len(nums) >= 3 else None


def assign_letterless_visits(visits: list[VisitSlide]) -> None:
    """차수 글자가 없는 재진 슬라이드(예: "24.09.26 (재진)")에 날짜순으로 글자 부여.

    글자를 아는 **마지막 차수 이후** 날짜인 슬라이드에만 잇대어 부여한다.
    글자 차수들 사이에 끼워 넣는 것은 하지 않는다 — 어느 글자인지 확정할 수
    없고, 순서를 잘못 매기면 정합의 "직전 차수" 기준이 엉뚱한 회차가 된다.
    글자를 받은 슬라이드는 기준영상·겹쳐보기에 정상 차수처럼 참여한다.
    """
    lettered = [v for v in visits if v.visit]
    pending = [v for v in visits if v.visit is None and v.slots and _date_key(v.date)]
    if not pending:
        return
    last = (max(lettered, key=lambda v: naming.letter_to_num(v.visit))
            if lettered else None)
    floor = _date_key(last.date) if last else None
    cur = last.visit if last else None
    for v in sorted(pending, key=lambda v: (_date_key(v.date), v.slide_index)):
        if floor is not None and _date_key(v.date) < floor:
            continue                      # 글자 차수들 사이 — 애매하므로 그대로 둔다
        cur = naming.next_visit_letter([cur] if cur else [])
        v.visit = cur
        if v.kind == "unknown":
            v.kind = "revisit"


def _cross_like(slide) -> bool:
    """십자뷰 슬라이드인가 — 도형 메타데이터만 본다(픽셀 복원 없음).

    앱 슬라이드는 슬롯 이름(PHOTO_SLOT_*)으로, 수제 슬라이드는 read_visit_slide
    폴백과 같은 기준(폭 8cm 이상 사진 5장)으로 판정한다. 접두를 SLOT 까지 보는
    이유: 얼굴 슬라이드의 PHOTO_FACE_* 가 십자뷰로 오인되면 안 된다.
    """
    big = 0
    for sh in slide.shapes:
        if str(getattr(sh, "name", "")).startswith(PHOTO_NAME_PREFIX + "SLOT"):
            return True
        if getattr(sh, "shape_type", None) == 13 and emu_to_cm(sh.width) >= 8.0:
            big += 1
    return big >= 5


def scan_ppt_visits(prs, cfg) -> dict:
    """슬라이드 라벨을 읽어 차수 장부를 만든다 — 픽셀 복원 없이.

    차수는 **십자뷰 슬라이드의 라벨** 하나로 통일한다("라벨 없는 십자뷰는 없다"
    전제). 판정 순서는 라벨 확인 → 사진 기하 확인 — 라벨 없는 장은 기하 검사
    비용도 내지 않는다. 라벨은 있는데 십자뷰가 아닌 장(얼굴·엑스레이 등)은
    excluded 로 돌려 화면이 "왜 제외됐는지" 보여줄 수 있게 한다. 십자뷰를
    하나도 못 알아본 덱만 라벨 장 전체로 물러난다(fallback) — 8cm 문턱에 걸린
    수제 덱에서 이력이 통째로 사라지는 것을 막는다.

    반환: {"visits":  [{visit, date, slide_no}...],   # 차수 장부 (글자 부여 완료)
           "excluded": [{slide_no, visit, date}...],  # 라벨은 있으나 십자뷰 아님
           "fallback": bool,
           "slides":   int}                           # 덱의 전체 장수
    """
    labeled: list[VisitSlide] = []
    cross: list[VisitSlide] = []
    for i, slide in enumerate(prs.slides):
        if not any(getattr(sh, "shape_type", None) == 13 for sh in slide.shapes):
            continue                  # 글만 있는 장(환자정보 등)은 차수가 아니다
        text = _visit_text(slide, cfg)          # ① 라벨 먼저 — 없으면 여기서 끝
        if not text:
            continue
        v, date, kind = parse_info_box(text)
        if kind == "unknown" and not (v or date):
            continue
        vs = VisitSlide(slide_index=i, visit=v, date=date, kind=kind,
                        label_text=text)
        labeled.append(vs)
        if _cross_like(slide):                  # ② 라벨이 있는 장만 기하 확인
            cross.append(vs)
    picked = cross or labeled                   # ③ 십자뷰 0개 덱만 관대한 폴백
    for vs in picked:
        vs.slots = {"_": None}    # 글자 부여 규칙이 보는 '차수 슬라이드' 표식
    assign_letterless_visits(picked)
    chosen = {id(vs) for vs in picked}
    return {
        "visits": [{"visit": vs.visit, "date": vs.date,
                    "slide_no": vs.slide_index + 1}
                   for vs in picked if vs.visit],
        "excluded": [{"slide_no": vs.slide_index + 1, "visit": vs.visit,
                      "date": vs.date}
                     for vs in labeled if id(vs) not in chosen],
        "fallback": bool(labeled) and not cross,
        "slides": len(prs.slides._sldIdLst),
    }


def references_for_registration(visits: list[VisitSlide]) -> dict[str, dict[str, np.ndarray]]:
    """
    슬롯별로 정합에 쓸 기준영상 모음.
    반환: { slot_name: { 'A'(초진): img, 'D'(직전): img, ... } }
    상위(§5.1.2)에서 직전 차수·초진 두 기준을 골라 register_best에 넘긴다.
    """
    out: dict[str, dict[str, np.ndarray]] = {}
    for vs in visits:
        if vs.visit is None:
            continue
        for slot, ref in vs.slots.items():
            out.setdefault(slot, {})[vs.visit] = ref.image
    return out


def last_label_style(prs, cfg) -> dict | None:
    """수제 PPT 의 라벨 박스 스타일(위치·크기·폰트) — 새 차수 슬라이드가 상속한다.

    이름 있는 박스(INFO_BOX/날짜칸)가 발견되면 앱이 만든 PPT 다 — None 을 돌려
    (0.4, 0.4) 고정 규칙을 유지한다. 수제 라벨 박스만 스타일을 뽑는다.
    """
    best = None
    for slide in prs.slides:
        if any(find_shape(slide, n) is not None
               for n in (cfg.ppt.info_box_name, _NOTE_DATE_BOX)):
            continue                            # 앱 슬라이드만 건너뛴다
        for sh in slide.shapes:
            if not getattr(sh, "has_text_frame", False):
                continue
            _v, dt, kind = parse_info_box(sh.text_frame.text or "")
            if not (dt and kind != "unknown"):
                continue
            f = None
            for para in sh.text_frame.paragraphs:
                if para.runs:
                    f = para.runs[0].font
                    break
            try:
                # 테마 색(_SchemeColor)은 .rgb 가 없다 — RGB 로 안 적힌 색은 버린다
                color = str(f.color.rgb) if f is not None else None
            except (AttributeError, TypeError):
                color = None
            best = {"x": emu_to_cm(sh.left), "y": emu_to_cm(sh.top),
                    "w": emu_to_cm(sh.width), "h": emu_to_cm(sh.height),
                    "size_pt": float(f.size.pt) if f is not None and f.size else None,
                    "font": f.name if f is not None else None,
                    "bold": bool(f.bold) if f is not None and f.bold is not None else None,
                    "color": color}
            lst = sh.text_frame._txBody.find(qn("a:lstStyle"))
            if lst is not None and len(lst):
                best["lst_style_xml"] = etree.tostring(lst)
            break                               # 마지막 차수 것이 최종 채택
    return best


def last_status_style(prs, cfg) -> dict | None:
    """수제 PPT 의 상태 칸(Tx/Rx/App Period) 스타일 — 위치·크기·폰트·줄간격.

    이름 있는 박스(NOTE_STATUS)가 보이면 앱 PPT 다 — None (앱 규약 유지).
    """
    best = None
    for slide in prs.slides:
        if find_shape(slide, _NOTE_STATUS_BOX) is not None:
            continue                            # 앱 슬라이드만 건너뛴다
        for sh in slide.shapes:
            if not getattr(sh, "has_text_frame", False):
                continue
            t = sh.text_frame.text or ""
            if "Tx. Period" not in t and "Rx. Period" not in t:
                continue
            f, spacing = None, None
            for para in sh.text_frame.paragraphs:
                if para.runs:
                    f = para.runs[0].font
                    spacing = para.line_spacing
                    break
            try:
                color = str(f.color.rgb) if f is not None else None
            except (AttributeError, TypeError):
                color = None
            best = {"x": emu_to_cm(sh.left), "y": emu_to_cm(sh.top),
                    "w": emu_to_cm(sh.width), "h": emu_to_cm(sh.height),
                    "size_pt": float(f.size.pt) if f is not None and f.size else None,
                    "font": f.name if f is not None else None,
                    "bold": bool(f.bold) if f is not None and f.bold is not None else None,
                    "color": color, "line_spacing": spacing}
            # 변환기(.ppt→.pptx)가 상자 안에 구워 둔 문단 규칙(lstStyle) —
            # 단락 앞 간격(spcPct) 등이 여기 있다. 통째로 이식해야 재현된다.
            lst = sh.text_frame._txBody.find(qn("a:lstStyle"))
            if lst is not None and len(lst):
                best["lst_style_xml"] = etree.tostring(lst)
            break
    return best


def note_role(sh, slide_w_cm: float) -> str | None:
    """이 도형이 규약 상자 다섯 중 어느 역할인가 — 이름이 아니라 **내용과 자리**로.

    수제 덱의 상자에는 이름 규약이 없다. 라벨과 기간 칸은 글로, 나머지 세 칸은
    사분면 위치로 알아본다.

    판정이 한 곳에 있어야 '통째 복사 상속' 과 '나머지 도형 복사' 가 **같은 상자를
    두 번** 넣지 않는다. 실제로 그랬다 — '전부 복사' 를 켜면 이미 물려받은 날짜/차수·
    Tx/Rx/App 상자가 한 벌 더 얹혀 글자가 겹쳤다.

    반환: "label" | "status" | NOTE_SOAP | NOTE_LL | NOTE_NEXT | None
    """
    if not getattr(sh, "has_text_frame", False):
        return None
    t = sh.text_frame.text or ""
    if any(m in t for m in _PERIOD_MARKS):
        return "status"
    _v, dt, kind = parse_info_box(t)
    if dt and kind != "unknown":
        return "label"
    if not t.strip():
        return None                         # 빈 상자는 대응으로 안 본다
    x, y = emu_to_cm(sh.left), emu_to_cm(sh.top)
    half = slide_w_cm / 2
    if x < half and y < 1.2:
        return None                         # 라벨 자리인데 라벨이 아니다
    return ("NOTE_LL" if x < half and y >= 12.0
            else "NOTE_SOAP" if x < half
            else "NOTE_NEXT" if y >= 6.0 else None)


def last_label_status_xml(prs, cfg) -> dict:
    """수제 PPT 마지막 차수의 노트 상자들 **원본 XML** — 통째 복사 상속용.

    라벨/상태는 내용으로, 나머지(좌상단 s/p·좌하단·우하단)는 **위치(사분면)** 로
    대응 상자를 찾는다. 원본에 있는 상자는 복사돼 폰트가 유지되고, 없는 상자만
    설정 기본 크기로 새로 만들어진다. 이름 규약 박스가 보이면 앱 PPT — 빈 dict.
    """
    out: dict = {}
    sw = emu_to_cm(prs.slide_width)
    half = sw / 2
    for slide in prs.slides:
        # 이름 규약 박스가 있는 슬라이드 = 앱이 만든 것 — 그 장만 건너뛴다.
        # 예전엔 하나만 보여도 전체를 포기했는데, 그러면 한 번 이어붙인 수제
        # PPT 는 두 번째 커밋부터 상속이 통째로 꺼졌다.
        if any(find_shape(slide, n) is not None
               for n in (cfg.ppt.info_box_name, _NOTE_DATE_BOX, _NOTE_STATUS_BOX)):
            continue
        found: dict = {}
        for sh in slide.shapes:
            key = note_role(sh, sw)
            if key:
                found[key] = sh
        if "label" in found or "status" in found:   # 차수 슬라이드만 채택
            out = {k: etree.tostring(v._element) for k, v in found.items()}
    return out
