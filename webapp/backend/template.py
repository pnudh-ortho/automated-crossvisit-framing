"""
템플릿 파싱 (Stage 1)

PPT 좌표를 하드코딩하지 않고 도형 '이름'으로 슬롯 위치를 읽는다.
레이아웃이 수정되어도(창 좌표가 바뀌어도) 코드는 무변경.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

from coords import WindowCm, emu_to_cm


# ── 저장 속도: 이미 압축된 것을 다시 압축하지 않는다 ──────────────────────────
#
# 덱을 저장할 때마다 python-pptx 는 **모든 내용물에 deflate 를 건다.** 그런데 덱의
# 99% 는 JPEG 이라 줄지가 않는다 — 실측한 131MB 덱에서
#
#     ppt/media (사진)   130.8 MB → 130.8 MB   0.0% 줄어듦
#     xml 등               2.5 MB →   0.1 MB  94.5% 줄어듦
#
# 즉 줄지도 않는 131MB 를 매번 훑느라 저장 시간의 대부분을 쓰고 있었고, 그 값은
# 노트북에서 그대로 사람이 기다리는 시간이 된다. 사진만 **그대로 담고**(STORED)
# XML 은 종전대로 압축하면 결과 파일 크기는 사실상 같고 저장만 빨라진다.
# 미리보기와 확정 저장이 함께 덕을 본다.
#
# python-pptx 가 압축 방식을 열어 주지 않아 쓰는 자리(_ZipPkgWriter.write)를
# 갈아 끼운다. 사설 API 라 판이 바뀌면 조용히 종전 동작으로 물러난다 — 저장이
# 느려질 뿐 결과는 같다.
def _store_media_uncompressed() -> bool:
    try:
        import zipfile

        from pptx.opc.serialized import _ZipPkgWriter
    except Exception:                                        # noqa: BLE001
        return False
    if getattr(_ZipPkgWriter, "_crocs_fast_write", False):
        return True

    def write(self, pack_uri, blob):
        name = pack_uri.membername
        self._zipf.writestr(
            name, blob,
            compress_type=(zipfile.ZIP_STORED if name.startswith("ppt/media/")
                           else zipfile.ZIP_DEFLATED))

    _ZipPkgWriter.write = write
    _ZipPkgWriter._crocs_fast_write = True
    return True


FAST_SAVE = _store_media_uncompressed()


@dataclass(frozen=True)
class SlideGeom:
    width_cm: float
    height_cm: float


def load_presentation(path: str | Path) -> Presentation:
    return Presentation(str(path))


def slide_geom(prs: Presentation) -> SlideGeom:
    return SlideGeom(emu_to_cm(prs.slide_width), emu_to_cm(prs.slide_height))


def find_shape(slide, name: str):
    """이름으로 도형 하나 찾기. 없으면 None."""
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def require_shape(slide, name: str):
    sh = find_shape(slide, name)
    if sh is None:
        raise KeyError(f"템플릿에 도형 '{name}'가 없습니다")
    return sh


def shape_window_cm(shape) -> WindowCm:
    """도형의 위치·크기를 cm 창으로."""
    return WindowCm(
        x=emu_to_cm(shape.left),
        y=emu_to_cm(shape.top),
        w=emu_to_cm(shape.width),
        h=emu_to_cm(shape.height),
    )


def slot_windows(slide, slot_names: list[str]) -> dict[str, WindowCm]:
    """슬롯 앵커 이름 → 창(cm). 삽입 위치 기준."""
    out = {}
    for name in slot_names:
        out[name] = shape_window_cm(require_shape(slide, name))
    return out


def mask_shapes(slide, mask_prefix: str) -> list:
    return [sh for sh in slide.shapes if sh.name.startswith(mask_prefix)]


def validate_template(prs: Presentation, cfg) -> dict:
    """
    템플릿이 사양을 만족하는지 점검하고 요약을 반환.
    첫 슬라이드에 슬롯 5개 + INFO_BOX + 마스크가 있어야 한다.
    """
    slide = prs.slides[0]
    geom = slide_geom(prs)
    wins = slot_windows(slide, cfg.ppt.slot_names)
    masks = mask_shapes(slide, cfg.ppt.mask_prefix)
    info = require_shape(slide, cfg.ppt.info_box_name)
    return {
        "slide_cm": (round(geom.width_cm, 3), round(geom.height_cm, 3)),
        "slots": {k: (round(v.x, 3), round(v.y, 3), round(v.w, 3), round(v.h, 3)) for k, v in wins.items()},
        "mask_count": len(masks),
        "info_box": info.name,
    }
